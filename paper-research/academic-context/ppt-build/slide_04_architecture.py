"""
Slide 4 — System Architecture (Redesigned: 3-Tier Swim-Lane)
=============================================================
Layout philosophy:
  TIER 1  (top strip)   — Operator  →  Commander  ↔  VAPT Protocol
  TIER 2  (centre zone) — ASG  ←──────────────────→  APG  (dual graph)
  TIER 3  (bottom zone) — 6 Specialized Agents  |  Tool Adapter Layer + 11 tools

Each tier is a labelled swim-lane. Arrows cross tier boundaries only where
architecturally meaningful (Commander → Agents, Agents → ASG, ASG → APG).
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.enum.shapes

# ── Palette ───────────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
BG_T1       = RGBColor(0x06, 0x10, 0x20)   # tier 1 bg (orchestration)
BG_T2       = RGBColor(0x05, 0x14, 0x0A)   # tier 2 bg (dual graph — subtle green tint)
BG_T3       = RGBColor(0x0A, 0x08, 0x1C)   # tier 3 bg (agents / tools)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
GREY_DARK   = RGBColor(0x40, 0x48, 0x58)
CARD_DARK   = RGBColor(0x10, 0x16, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.CENTER, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

def arr(slide, x1, y1, x2, y2, color=GREY_MID, lw=1.2, bidirectional=False):
    c = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(lw)
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    if bidirectional:
        te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    else:
        te.set('type', 'none')
    return c

# ── Load + add slide ──────────────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome (left bar + top + bottom) ─────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_CYAN)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_CYAN)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_CYAN)

# ── Title ─────────────────────────────────────────────────────────────────────
txt(slide, "SYSTEM ARCHITECTURE", Inches(0.3), Inches(0.07), Inches(6), Inches(0.26),
    size=10, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
txt(slide, "CMatrix — Three-Tier Architecture Overview",
    Inches(0.3), Inches(0.33), Inches(11), Inches(0.5),
    size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
#  TIER 1  —  ORCHESTRATION  (top band, full width)
#  y: 0.88 → 2.28  (height 1.40")
# ═══════════════════════════════════════════════════════════════════════════════
T1_T = Inches(0.88)
T1_H = Inches(1.40)
T1_B = T1_T + T1_H

# Tier background + label
box(slide, Inches(0.18), T1_T, SLIDE_W-Inches(0.36), T1_H,
    fill=BG_T1, line_color=ACCENT_CYAN, lw=0.6)
txt(slide, "① ORCHESTRATION", Inches(0.22), T1_T+Inches(0.05),
    Inches(1.6), Inches(0.22), size=7.5, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

# — OPERATOR node —
op_l, op_t = Inches(0.35), T1_T+Inches(0.28)
op_w, op_h = Inches(1.7), Inches(0.96)
box(slide, op_l, op_t, op_w, op_h, fill=RGBColor(0x0C,0x18,0x2E), line_color=GREY_MID, lw=0.8)
txt(slide, "OPERATOR", op_l, op_t+Inches(0.06), op_w, Inches(0.24),
    size=9, bold=True, color=GREY_MID)
txt(slide, "Target · Scope · Mode", op_l, op_t+Inches(0.34), op_w, Inches(0.55),
    size=8, color=GREY_MID, italic=True)

# — Arrow: Operator → Commander —
arr(slide, op_l+op_w, op_t+op_h/2, Inches(2.28), op_t+op_h/2, color=ACCENT_CYAN, lw=1.5)
txt(slide, "mission config", Inches(2.08), op_t+op_h/2-Inches(0.3),
    Inches(0.9), Inches(0.26), size=7, italic=True, color=GREY_DARK)

# — COMMANDER block —
cmd_l, cmd_t = Inches(2.28), T1_T+Inches(0.15)
cmd_w, cmd_h = Inches(3.1), Inches(1.1)
box(slide, cmd_l, cmd_t, cmd_w, cmd_h, fill=RGBColor(0x04,0x16,0x2E), line_color=ACCENT_CYAN, lw=2.0)
# Cyan header bar
box(slide, cmd_l, cmd_t, cmd_w, Inches(0.3), fill=ACCENT_CYAN)
txt(slide, "COMMANDER AGENT", cmd_l, cmd_t+Inches(0.03), cmd_w, Inches(0.24),
    size=9.5, bold=True, color=BG_DARK)
bullets = ["Reads ASG + APG state", "Plans & delegates tasks", "Approves High-risk ops", "Writes to APG only"]
for i, b in enumerate(bullets):
    txt(slide, f"• {b}", cmd_l+Inches(0.1), cmd_t+Inches(0.35)+i*Inches(0.18),
        cmd_w-Inches(0.15), Inches(0.18), size=8, color=GREY_MID, align=PP_ALIGN.LEFT)

# — VAPT PROTOCOL block —
vp_l, vp_t = Inches(5.6), T1_T+Inches(0.28)
vp_w, vp_h = Inches(2.3), Inches(0.96)
box(slide, vp_l, vp_t, vp_w, vp_h, fill=RGBColor(0x10,0x0C,0x22), line_color=ACCENT_PURP, lw=1.2)
txt(slide, "VAPT PROTOCOL PROMPT", vp_l, vp_t+Inches(0.06), vp_w, Inches(0.22),
    size=8, bold=True, color=ACCENT_PURP)
txt(slide, "Phase rules · Re-plan triggers\nTermination conditions\nMethodology-as-config",
    vp_l+Inches(0.1), vp_t+Inches(0.32), vp_w-Inches(0.15), Inches(0.58),
    size=8, color=GREY_MID, align=PP_ALIGN.LEFT)

# Commander ↔ VAPT Protocol
arr(slide, cmd_l+cmd_w, cmd_t+cmd_h/2, vp_l, vp_t+vp_h/2, color=ACCENT_PURP, lw=1.0, bidirectional=True)

# ═══════════════════════════════════════════════════════════════════════════════
#  TIER 2  —  DUAL-GRAPH WORLD MODEL  (centre band)
#  y: 2.34 → 4.18  (height 1.84")
# ═══════════════════════════════════════════════════════════════════════════════
T2_T = T1_B + Inches(0.06)
T2_H = Inches(1.84)
T2_B = T2_T + T2_H

box(slide, Inches(0.18), T2_T, SLIDE_W-Inches(0.36), T2_H,
    fill=BG_T2, line_color=ACCENT_LIME, lw=0.6)
txt(slide, "② DUAL-GRAPH WORLD MODEL", Inches(0.22), T2_T+Inches(0.05),
    Inches(2.6), Inches(0.22), size=7.5, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)

# — ASG block (left) —
asg_l, asg_t = Inches(0.35), T2_T+Inches(0.3)
asg_w, asg_h = Inches(4.7), Inches(1.42)
box(slide, asg_l, asg_t, asg_w, asg_h, fill=RGBColor(0x04,0x18,0x0C), line_color=ACCENT_LIME, lw=1.8)
box(slide, asg_l, asg_t, asg_w, Inches(0.3), fill=ACCENT_LIME)
txt(slide, "ATTACK SURFACE GRAPH  (ASG)", asg_l, asg_t+Inches(0.04), asg_w, Inches(0.24),
    size=9.5, bold=True, color=BG_DARK)
# Node types row 1
node_cols = [
    ("Domain", ACCENT_LIME), ("Host", ACCENT_LIME), ("Port", ACCENT_LIME),
    ("Service", ACCENT_LIME), ("Technology", ACCENT_LIME),
]
nw = (asg_w - Inches(0.3)) / len(node_cols)
for i, (name, clr) in enumerate(node_cols):
    nl = asg_l + Inches(0.15) + i * nw
    box(slide, nl, asg_t+Inches(0.36), nw-Inches(0.06), Inches(0.25),
        fill=RGBColor(0x08,0x28,0x12), line_color=clr, lw=0.5)
    txt(slide, name, nl, asg_t+Inches(0.38), nw-Inches(0.06), Inches(0.22),
        size=7.5, bold=True, color=clr)
# Node types row 2
node_cols2 = [
    ("Endpoint", ACCENT_LIME), ("Parameter", ACCENT_LIME),
    ("Vulnerability", ACCENT_LIME), ("Evidence", ACCENT_LIME),
]
nw2 = (asg_w - Inches(0.3)) / 4
for i, (name, clr) in enumerate(node_cols2):
    nl = asg_l + Inches(0.15) + i * nw2
    box(slide, nl, asg_t+Inches(0.66), nw2-Inches(0.06), Inches(0.25),
        fill=RGBColor(0x08,0x28,0x12), line_color=clr, lw=0.5)
    txt(slide, name, nl, asg_t+Inches(0.68), nw2-Inches(0.06), Inches(0.22),
        size=7.5, bold=True, color=clr)
txt(slide, "Discovered Reality — confirmed facts only. Never contains hypotheses.",
    asg_l+Inches(0.12), asg_t+Inches(0.97), asg_w-Inches(0.2), Inches(0.38),
    size=8, italic=True, color=RGBColor(0x60,0xC0,0x70), align=PP_ALIGN.LEFT)

# — APG block (right) —
apg_l, apg_t = Inches(8.3), T2_T+Inches(0.3)
apg_w, apg_h = Inches(4.8), Inches(1.42)
box(slide, apg_l, apg_t, apg_w, apg_h, fill=RGBColor(0x1E,0x10,0x04), line_color=ACCENT_GOLD, lw=1.8)
box(slide, apg_l, apg_t, apg_w, Inches(0.3), fill=ACCENT_GOLD)
txt(slide, "ATTACK PATH GRAPH  (APG)", apg_l, apg_t+Inches(0.04), apg_w, Inches(0.24),
    size=9.5, bold=True, color=BG_DARK)
apg_nodes = [("AttackChain", ACCENT_GOLD), ("ChainStep", ACCENT_GOLD), ("Impact", ACCENT_GOLD)]
apg_nw = (apg_w - Inches(0.3)) / 3
for i, (name, clr) in enumerate(apg_nodes):
    nl = apg_l + Inches(0.15) + i * apg_nw
    box(slide, nl, apg_t+Inches(0.36), apg_nw-Inches(0.1), Inches(0.25),
        fill=RGBColor(0x2E,0x18,0x04), line_color=clr, lw=0.5)
    txt(slide, name, nl, apg_t+Inches(0.38), apg_nw-Inches(0.1), Inches(0.22),
        size=7.5, bold=True, color=clr)
txt(slide, "risk_score  ·  validation_status  ·  priority",
    apg_l+Inches(0.12), apg_t+Inches(0.68), apg_w-Inches(0.2), Inches(0.22),
    size=8, color=GREY_MID, align=PP_ALIGN.LEFT)
txt(slide, "Inferred Opportunity — attack reasoning only. Never contains raw scan data.",
    apg_l+Inches(0.12), apg_t+Inches(0.95), apg_w-Inches(0.2), Inches(0.38),
    size=8, italic=True, color=RGBColor(0xC0,0x90,0x20), align=PP_ALIGN.LEFT)

# — Separation zone between ASG and APG —
sep_l = asg_l + asg_w
sep_mid = (sep_l + apg_l) / 2
box(slide, sep_l, T2_T+Inches(0.3), apg_l-sep_l, asg_h,
    fill=RGBColor(0x06,0x10,0x08))
txt(slide, "STRICT\nSEPARATION", sep_l+Inches(0.04), T2_T+Inches(0.62),
    apg_l-sep_l-Inches(0.04), Inches(0.55),
    size=7, bold=True, color=GREY_MID)
# Dashed-style separator line
box(slide, sep_mid-Inches(0.01), T2_T+Inches(0.3), Inches(0.02), asg_h,
    fill=GREY_DARK)

# — ASG → APG arrow: Commander derives chains —
arr(slide, asg_l+asg_w/2, asg_t+asg_h, asg_l+asg_w/2, T2_B+Inches(0.1),
    color=ACCENT_GOLD, lw=0.8)

# — Commander → ASG (reads) —
arr(slide, cmd_l+cmd_w/2, T1_B, asg_l+asg_w*0.3, T2_T, color=ACCENT_CYAN, lw=1.2)
# — APG → Commander (feedback) —
arr(slide, apg_l+apg_w*0.5, T2_T, cmd_l+cmd_w*0.7, T1_B, color=ACCENT_GOLD, lw=1.0)

# ═══════════════════════════════════════════════════════════════════════════════
#  TIER 3  —  AGENTS + TOOL LAYER  (bottom band)
#  y: 4.26 → 7.24  (height 2.98")
# ═══════════════════════════════════════════════════════════════════════════════
T3_T = T2_B + Inches(0.06)
T3_H = SLIDE_H - T3_T - Inches(0.28)

box(slide, Inches(0.18), T3_T, SLIDE_W-Inches(0.36), T3_H,
    fill=BG_T3, line_color=ACCENT_PURP, lw=0.6)
txt(slide, "③ SPECIALIZED AGENTS  +  TOOL ADAPTER LAYER", Inches(0.22), T3_T+Inches(0.04),
    Inches(5), Inches(0.22), size=7.5, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

# — 6 Agent cards (left section) —
agents = [
    ("Recon",      "Amass · httpx · Nmap",                         ACCENT_LIME),
    ("Analysis",   "WhatWeb · Gobuster · ffuf · Nuclei · ZAP",     ACCENT_CYAN),
    ("Research",   "NVD · Exploit-DB · GitHub",                    ACCENT_GOLD),
    ("Validation", "SQLMap · Metasploit",                          ACCENT_RED),
    ("Evidence",   "EyeWitness",                                   ACCENT_PURP),
    ("Report",     "Reads ASG + APG",                              GREY_MID),
]
ag_col_w = Inches(8.9)
aw = (ag_col_w - Inches(0.45)) / 6
ah = T3_H - Inches(0.35)
ag_start = Inches(0.35)
ag_top = T3_T + Inches(0.28)

for i, (name, tools, clr) in enumerate(agents):
    al = ag_start + i * (aw + Inches(0.06))
    box(slide, al, ag_top, aw, ah, fill=RGBColor(0x0C,0x0E,0x22), line_color=clr, lw=1.2)
    # Color header
    box(slide, al, ag_top, aw, Inches(0.24), fill=clr)
    txt(slide, name, al, ag_top+Inches(0.02), aw, Inches(0.22),
        size=7.5, bold=True, color=BG_DARK)
    txt(slide, tools, al+Inches(0.06), ag_top+Inches(0.3), aw-Inches(0.1), ah-Inches(0.3),
        size=7, color=clr, align=PP_ALIGN.LEFT, wrap=True)
    # Arrow from Commander to agent (downward)
    arr(slide, cmd_l + cmd_w*0.4 + i*(cmd_w*0.1), T1_B+Inches(0.05),
        al+aw/2, ag_top, color=clr, lw=0.5)
    # Arrow from agent to ASG
    if name not in ("Report",):
        arr(slide, al+aw/2, ag_top,
            asg_l + asg_w*0.2 + i*(asg_w*0.08), T2_B,
            color=RGBColor(0x20,0x50,0x28), lw=0.4)

# — Tool Adapter Layer (right section, inside T3) —
tal_l = Inches(9.45)
tal_w = SLIDE_W - tal_l - Inches(0.22)
tal_t = T3_T + Inches(0.28)
tal_h = T3_H - Inches(0.35)
box(slide, tal_l, tal_t, tal_w, tal_h, fill=RGBColor(0x0A,0x08,0x1E), line_color=ACCENT_PURP, lw=1.2)
# Header
box(slide, tal_l, tal_t, tal_w, Inches(0.26), fill=ACCENT_PURP)
txt(slide, "TOOL ADAPTER LAYER", tal_l, tal_t+Inches(0.03), tal_w, Inches(0.22),
    size=8, bold=True, color=BG_DARK)

# Risk gate tiers
risk_defs = [
    ("LOW",  "Passive", "Execute immediately", ACCENT_LIME),
    ("MED",  "Active",  "LLM Classifier → exec/escalate", ACCENT_GOLD),
    ("HIGH", "Exploit", "Commander Mailbox approval", ACCENT_RED),
]
rh = (tal_h - Inches(0.32)) / 3
for i, (tier, kind, action, clr) in enumerate(risk_defs):
    rt = tal_t + Inches(0.3) + i * rh
    box(slide, tal_l+Inches(0.1), rt, tal_w-Inches(0.2), rh-Inches(0.06),
        fill=RGBColor(0x10,0x0C,0x24), line_color=clr, lw=0.7)
    box(slide, tal_l+Inches(0.1), rt, Inches(0.38), rh-Inches(0.06), fill=clr)
    txt(slide, tier, tal_l+Inches(0.1), rt+Inches(0.02), Inches(0.38), rh-Inches(0.1),
        size=7, bold=True, color=BG_DARK)
    txt(slide, kind, tal_l+Inches(0.52), rt+Inches(0.02), tal_w-Inches(0.65), Inches(0.2),
        size=7.5, bold=True, color=clr, align=PP_ALIGN.LEFT)
    txt(slide, action, tal_l+Inches(0.52), rt+Inches(0.22), tal_w-Inches(0.65), Inches(0.3),
        size=7, color=GREY_MID, align=PP_ALIGN.LEFT, italic=True)

# ── Slide number ──────────────────────────────────────────────────────────────
txt(slide, "04", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.52),
    Inches(0.35), Inches(0.42), size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 4 (Architecture) rewritten.")
