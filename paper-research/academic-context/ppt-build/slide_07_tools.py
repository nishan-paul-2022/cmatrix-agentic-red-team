"""
Slide 7 — VAPT Tool Catalogue (11 tools CMatrix autonomously operates)
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)
BG_DARK2    = RGBColor(0x08, 0x0C, 0x1A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome ─────────────────────────────────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_RED)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_RED)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_RED)

txt(slide, "OFFENSIVE TOOL CATALOGUE", Inches(0.3), Inches(0.08), Inches(8), Inches(0.3),
    size=11, bold=True, color=ACCENT_RED)
txt(slide, "11 Industry-Standard Security Tools — AI Agent Operated",
    Inches(0.3), Inches(0.38), Inches(12), Inches(0.62), size=30, bold=True, color=WHITE)
box(slide, Inches(0.3), Inches(1.0), Inches(6), Inches(0.03), fill=ACCENT_RED)

# ── Phase column headers ───────────────────────────────────────────────────────
phases = [
    ("PHASE 1  ·  RECONNAISSANCE", Inches(0.3),  Inches(4.0), ACCENT_LIME),
    ("PHASE 2  ·  ANALYSIS",        Inches(4.55), Inches(4.0), ACCENT_CYAN),
    ("PHASE 3  ·  VALIDATION",      Inches(9.2),  Inches(2.0), ACCENT_RED),
    ("PHASE 3  ·  EVIDENCE",        Inches(11.45),Inches(1.7), ACCENT_PURP),
]
ph_t = Inches(1.08)
ph_h = Inches(0.38)
for label, pl, pw, clr in phases:
    box(slide, pl, ph_t, pw, ph_h, fill=RGBColor(0x08,0x10,0x1C), line_color=clr, lw=1.2)
    txt(slide, label, pl+Inches(0.1), ph_t+Inches(0.06), pw-Inches(0.15), ph_h-Inches(0.1),
        size=9.5, bold=True, color=clr, align=PP_ALIGN.CENTER)

# ── Tool definitions ───────────────────────────────────────────────────────────
# (number, name, phase_color, agent, role_short, details)
tools = [
    (1,  "Amass",        ACCENT_LIME, "Recon Agent",      "Subdomain Enumeration",
     "DNS brute-force · certificate transparency logs · passive OSINT. "
     "Discovers all external subdomains and maps attack surface entry points."),
    (2,  "httpx",        ACCENT_LIME, "Recon Agent",      "Live Host Probing",
     "Validates which discovered hosts respond over HTTP/HTTPS. "
     "Identifies web servers, status codes, redirects, TLS details."),
    (3,  "Nmap",         ACCENT_LIME, "Recon Agent",      "Port + Service Scan",
     "Port scanning, service fingerprinting, OS detection. "
     "Optional NSE script execution for vuln detection on open services."),
    (4,  "WhatWeb",      ACCENT_CYAN, "Analysis Agent",   "Technology Fingerprinting",
     "Identifies CMS, frameworks, server software, JS libraries, and versions "
     "from HTTP responses. Seeds Technology nodes for the ASG."),
    (5,  "Gobuster",     ACCENT_CYAN, "Analysis Agent",   "Directory Brute-Force",
     "Discovers hidden paths, admin panels, backup files, and exposed resources "
     "via wordlist-based HTTP path enumeration."),
    (6,  "ffuf",         ACCENT_CYAN, "Analysis Agent",   "API + Parameter Fuzzing",
     "Fast web fuzzer for API route discovery, parameter fuzzing, "
     "virtual host enumeration, and undocumented endpoint identification."),
    (7,  "Nuclei",       ACCENT_CYAN, "Analysis Agent",   "Template-Based Vuln Scan",
     "Matches discovered services and technologies against a library of "
     "CVE and misconfiguration templates for automated vulnerability detection."),
    (8,  "OWASP ZAP",   ACCENT_CYAN, "Analysis Agent",   "Active Web App Scan",
     "Crawls and actively tests for OWASP Top 10: XSS, CSRF, injection flaws, "
     "authentication weaknesses, and insecure direct object references."),
    (9,  "SQLMap",       ACCENT_RED,  "Validation Agent", "SQL Injection Validation",
     "Automated SQLi detection and exploitation. Confirms injection points, "
     "extracts database data, tests for OS-level access escalation."),
    (10, "Metasploit",   ACCENT_RED,  "Validation Agent", "Exploit + RCE Demo",
     "Executes known exploits against identified vulnerabilities. "
     "Validates APG ChainSteps, demonstrates impact, achieves proof of exploitation."),
    (11, "EyeWitness",   ACCENT_PURP, "Evidence Agent",   "Screenshot Evidence",
     "Headless screenshot capture of web pages, exposed panels, and API responses. "
     "Produces visual proof artifacts linked to ASG Evidence nodes."),
]

# Layout: 3 columns — col0: tools 1-3, col1: tools 4-8, col2: tools 9-11
tool_t_start = Inches(1.55)
col_defs = [
    (Inches(0.3),  Inches(4.0),  [0,1,2]),
    (Inches(4.55), Inches(4.5),  [3,4,5,6,7]),
    (Inches(9.2),  Inches(3.9),  [8,9,10]),
]

for col_l, col_w, indices in col_defs:
    row_h = (SLIDE_H - tool_t_start - Inches(0.6)) / max(len(indices), 1)
    for row_i, tool_idx in enumerate(indices):
        num, name, clr, agent, role, detail = tools[tool_idx]
        t = tool_t_start + row_i * row_h
        h = row_h - Inches(0.08)
        # Card
        box(slide, col_l, t, col_w, h, fill=CARD_BG, line_color=clr, lw=1.0)
        # Number badge
        badge_w = Inches(0.4)
        box(slide, col_l, t, badge_w, h, fill=clr)
        txt(slide, str(num), col_l+Inches(0.02), t+h/2-Inches(0.18),
            badge_w-Inches(0.04), Inches(0.36), size=14, bold=True,
            color=BG_DARK, align=PP_ALIGN.CENTER)
        # Tool name
        txt(slide, name, col_l+badge_w+Inches(0.1), t+Inches(0.06),
            col_w-badge_w-Inches(0.15), Inches(0.3), size=13, bold=True, color=clr)
        # Role tag + agent
        txt(slide, f"{role}  ·  {agent}",
            col_l+badge_w+Inches(0.1), t+Inches(0.35),
            col_w-badge_w-Inches(0.15), Inches(0.22),
            size=8.5, bold=False, italic=True, color=GREY_MID)
        # Detail text (only if enough height)
        if h > Inches(0.75):
            txt(slide, detail,
                col_l+badge_w+Inches(0.1), t+Inches(0.58),
                col_w-badge_w-Inches(0.15), h-Inches(0.65),
                size=9, bold=False, color=GREY_MID, wrap=True)

# ── Safety note ────────────────────────────────────────────────────────────────
box(slide, Inches(0.3), SLIDE_H-Inches(0.62), SLIDE_W-Inches(0.6), Inches(0.48),
    fill=RGBColor(0x0E,0x08,0x18), line_color=ACCENT_PURP, lw=1.0)
txt(slide,
    "🔒  Every tool is wrapped in a Tool Adapter — agents reason about targets, not command syntax.  "
    "No irreversible offensive operation executes without Commander-level scope validation.",
    Inches(0.5), SLIDE_H-Inches(0.58), SLIDE_W-Inches(0.9), Inches(0.4),
    size=10.5, bold=False, color=ACCENT_PURP, wrap=False)

txt(slide, "07", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.55),
    Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_RED, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 7 added.")
