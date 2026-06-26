"""
Slide 8 — Real-World Scenario Walkthrough: shopvault.io
Visual phase-by-phase timeline showing CMatrix operating end-to-end.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.enum.shapes

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

def arrow_h(slide, x1, y, x2, color=GREY_MID, lw=2.0):
    c = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color; c.line.width = Pt(lw)
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    etree.SubElement(ln, qn('a:tailEnd')).set('type', 'none')
    return c

prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome ─────────────────────────────────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_GOLD)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)

txt(slide, "REAL-WORLD SCENARIO", Inches(0.3), Inches(0.08), Inches(8), Inches(0.3),
    size=11, bold=True, color=ACCENT_GOLD)
txt(slide, "End-to-End Mission: shopvault.io Black-Box Assessment",
    Inches(0.3), Inches(0.38), Inches(12), Inches(0.62), size=28, bold=True, color=WHITE)

# ── Scenario context banner ────────────────────────────────────────────────────
box(slide, Inches(0.3), Inches(1.0), SLIDE_W-Inches(0.6), Inches(0.55),
    fill=RGBColor(0x10,0x12,0x04), line_color=ACCENT_GOLD, lw=1.2)
txt(slide,
    "Scenario:  A security firm is hired to assess shopvault.io (mid-sized e-commerce).  "
    "Scope: all subdomains, web applications, REST APIs.  Mode: Black-Box (zero prior knowledge).  "
    "Operator configures CMatrix with root domain + scope — then presses start.",
    Inches(0.5), Inches(1.06), SLIDE_W-Inches(0.9), Inches(0.44),
    size=11, italic=True, color=ACCENT_GOLD, wrap=True)

# ══════════════════════════════════════════════
#  HORIZONTAL TIMELINE — 4 phases
# ══════════════════════════════════════════════
tl_top = Inches(1.7)
phase_defs = [
    # (label, color, x_start, width)
    ("Phase 1\nRecon",      ACCENT_LIME, Inches(0.3),   Inches(2.85)),
    ("Phase 2\nAnalysis",   ACCENT_CYAN, Inches(3.28),  Inches(4.0)),
    ("Phase 3\nValidation", ACCENT_RED,  Inches(7.41),  Inches(3.8)),
    ("Phase 4\nReport",     ACCENT_PURP, Inches(11.34), Inches(1.82)),
]

# Timeline spine
box(slide, Inches(0.3), tl_top+Inches(0.4), SLIDE_W-Inches(0.6), Inches(0.06), fill=GREY_MID)

for label, clr, px, pw in phase_defs:
    # Phase header box
    box(slide, px, tl_top, pw, Inches(0.46),
        fill=RGBColor(0x08,0x10,0x1C), line_color=clr, lw=1.5)
    txt(slide, label, px+Inches(0.08), tl_top+Inches(0.04), pw-Inches(0.12), Inches(0.38),
        size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
    # Dot on spine
    dot_x = px + pw/2 - Inches(0.12)
    box(slide, dot_x, tl_top+Inches(0.28), Inches(0.24), Inches(0.24), fill=clr)

# ══════════════════════════════════════════════
#  PHASE DETAIL PANELS (below timeline)
# ══════════════════════════════════════════════
panel_top = tl_top + Inches(0.65)
panel_bottom = SLIDE_H - Inches(0.6)
panel_h = panel_bottom - panel_top

# ── Phase 1 Panel: Reconnaissance ─────────────────────────────────────────────
p1_l = Inches(0.3); p1_w = Inches(2.85)
box(slide, p1_l, panel_top, p1_w, panel_h,
    fill=RGBColor(0x04,0x14,0x08), line_color=ACCENT_LIME, lw=1.0)
txt(slide, "🕵️  Recon Agent", p1_l+Inches(0.1), panel_top+Inches(0.08),
    p1_w-Inches(0.15), Inches(0.28), size=11, bold=True, color=ACCENT_LIME)
items1 = [
    ("Amass", "14 subdomains found\nincl. api · admin · staging · pay"),
    ("httpx",  "11 live hosts validated\nstaging returns unexpected 200"),
    ("Nmap",  "Ports 80, 443, 8080, 8443 open\npay.shopvault.io: expired TLS cert"),
    ("ASG →", "37 new nodes written"),
]
for i, (tool, detail) in enumerate(items1):
    t = panel_top + Inches(0.42) + i * Inches(0.86)
    box(slide, p1_l+Inches(0.1), t, p1_w-Inches(0.2), Inches(0.82),
        fill=RGBColor(0x08,0x20,0x0C), line_color=ACCENT_LIME, lw=0.5)
    txt(slide, tool, p1_l+Inches(0.18), t+Inches(0.05), p1_w-Inches(0.3), Inches(0.24),
        size=10, bold=True, color=ACCENT_LIME)
    txt(slide, detail, p1_l+Inches(0.18), t+Inches(0.3), p1_w-Inches(0.3), Inches(0.48),
        size=9, color=GREY_MID, wrap=True)

# ── Phase 2 Panel: Analysis + Intelligence ─────────────────────────────────────
p2_l = Inches(3.28); p2_w = Inches(4.0)
box(slide, p2_l, panel_top, p2_w, panel_h,
    fill=RGBColor(0x04,0x10,0x18), line_color=ACCENT_CYAN, lw=1.0)
txt(slide, "🔬  Analysis Agent  +  🔍  Research Agent",
    p2_l+Inches(0.1), panel_top+Inches(0.08), p2_w-Inches(0.15), Inches(0.28),
    size=10, bold=True, color=ACCENT_CYAN)
items2 = [
    ("WhatWeb",    "WordPress 5.9.3 · Django on api\nNginx 1.18.0 reverse proxy → CVE-2022-21661 seeded"),
    ("Gobuster",   "/admin/login · /admin/users\n/backup/db_export_2023.sql EXPOSED!"),
    ("ffuf",       "Undocumented: /api/v1/internal/users\nuser_id param accepts unsanitized input → IDOR"),
    ("Nuclei",     "pay.shopvault.io: expired TLS\nstaging: default creds · WooCommerce 6.1 CVE"),
    ("OWASP ZAP", "XSS on /search?q= · SQL errors on staging\nMissing CSRF on checkout"),
    ("APG →",     "3 AttackChains seeded: risk 8.8, 8.1, 7.5\n61 new ASG nodes written"),
]
for i, (tool, detail) in enumerate(items2):
    t = panel_top + Inches(0.42) + i * Inches(0.72)
    box(slide, p2_l+Inches(0.1), t, p2_w-Inches(0.2), Inches(0.68),
        fill=RGBColor(0x06,0x16,0x22), line_color=ACCENT_CYAN, lw=0.5)
    txt(slide, tool, p2_l+Inches(0.18), t+Inches(0.04), p2_w-Inches(0.3), Inches(0.22),
        size=10, bold=True, color=ACCENT_CYAN)
    txt(slide, detail, p2_l+Inches(0.18), t+Inches(0.28), p2_w-Inches(0.3), Inches(0.38),
        size=8.5, color=GREY_MID, wrap=True)

# ── Phase 3 Panel: Validation + Evidence ──────────────────────────────────────
p3_l = Inches(7.41); p3_w = Inches(3.8)
box(slide, p3_l, panel_top, p3_w, panel_h,
    fill=RGBColor(0x18,0x06,0x06), line_color=ACCENT_RED, lw=1.0)
txt(slide, "🎯  Validation Agent  +  📸  Evidence Agent",
    p3_l+Inches(0.1), panel_top+Inches(0.08), p3_w-Inches(0.15), Inches(0.28),
    size=10, bold=True, color=ACCENT_RED)
chains = [
    ("Chain-01  risk 8.8", ACCENT_RED,
     "SQLMap → CVE-2022-21661 confirmed\nAdmin hash extracted + cracked\nMetasploit → Web shell → RCE\nStatus: VALIDATED  (escalated to 9.1)"),
    ("Chain-02  risk 7.5", ACCENT_GOLD,
     "SQLMap on user_id param\nIDOR confirmed — any customer orders\nStatus: VALIDATED"),
    ("Chain-03  risk 8.1", ACCENT_GOLD,
     "SQLMap on staging login\nBlind SQLi → staging DB creds\nCredential reuse risk flagged\nStatus: VALIDATED"),
    ("Chain-04  (bonus)", ACCENT_PURP,
     "Direct HTTP GET on db_export_2023.sql\nFull customer PII accessible publicly\nStatus: VALIDATED immediately"),
]
for i, (chain, clr, detail) in enumerate(chains):
    t = panel_top + Inches(0.42) + i * Inches(1.08)
    box(slide, p3_l+Inches(0.1), t, p3_w-Inches(0.2), Inches(1.0),
        fill=RGBColor(0x20,0x08,0x08), line_color=clr, lw=0.8)
    txt(slide, chain, p3_l+Inches(0.18), t+Inches(0.05), p3_w-Inches(0.3), Inches(0.24),
        size=10, bold=True, color=clr)
    txt(slide, detail, p3_l+Inches(0.18), t+Inches(0.3), p3_w-Inches(0.3), Inches(0.66),
        size=8.5, color=GREY_MID, wrap=True)

# ── Phase 4 Panel: Report ─────────────────────────────────────────────────────
p4_l = Inches(11.34); p4_w = Inches(1.82)
box(slide, p4_l, panel_top, p4_w, panel_h,
    fill=RGBColor(0x10,0x08,0x1C), line_color=ACCENT_PURP, lw=1.0)
txt(slide, "📝 Report\nAgent",
    p4_l+Inches(0.08), panel_top+Inches(0.08), p4_w-Inches(0.12), Inches(0.4),
    size=10, bold=True, color=ACCENT_PURP, align=PP_ALIGN.CENTER)
report_items = [
    "Executive Summary",
    "4 validated chains",
    "RCE + PII exposure",
    "11 vuln entries",
    "Full ASG surface map",
    "Evidence at every step",
    "Remediation guidance",
    "ZERO manual commands",
]
for i, item in enumerate(report_items):
    txt(slide, f"• {item}",
        p4_l+Inches(0.1), panel_top+Inches(0.52)+i*Inches(0.55),
        p4_w-Inches(0.15), Inches(0.5),
        size=9, color=GREY_MID, wrap=True)

# ── ASG + APG counts strip ────────────────────────────────────────────────────
box(slide, Inches(0.3), SLIDE_H-Inches(0.58), SLIDE_W-Inches(0.6), Inches(0.44),
    fill=RGBColor(0x0A,0x10,0x20), line_color=ACCENT_CYAN, lw=1.0)
txt(slide,
    "Final ASG:  14 subdomains · 11 live hosts · 28 open ports · 19 endpoints · 7 parameters · 11 vulnerabilities · 4 validated chains  "
    "|  Operator issued ZERO manual commands",
    Inches(0.5), SLIDE_H-Inches(0.54), SLIDE_W-Inches(0.9), Inches(0.38),
    size=11, bold=True, color=ACCENT_CYAN, wrap=False)

txt(slide, "08", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.55),
    Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 8 added.")
