"""
Slide 9 — Attack Chain Lifecycle (Redesigned)
==============================================
Layout:
  TOP HALF:   Horizontal FSM —  HYPOTHESIZED → PARTIALLY VALIDATED
                                                    ↓ success      ↓ step fails (after cap)
                                               VALIDATED       RULED_OUT
              Each state box shows: name + what it means + visual status colour

  BOTTOM LEFT: Validation Agent Self-Debugging Loop
               Visual cycle: Attempt → Diagnose → Contextualize → Adapt → Retry
               Horizontal step-by-step with loop-back arrow

  BOTTOM RIGHT: Risk Scoring + chain priority example from shopvault.io
"""
from palette import *
import pptx.enum.shapes


BG_DARK2    = RGBColor(0x0A, 0x0D, 0x1A)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome ─────────────────────────────────────────────────────────────────────
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_GOLD)
    box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)
    box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)

    # ── Title ─────────────────────────────────────────────────────────────────────
    txt(slide, "ATTACK CHAIN LIFECYCLE", Inches(0.3), Inches(0.07), Inches(6), Inches(0.24),
        size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "APG Chain State Machine + Validation Agent Self-Debugging Loop",
        Inches(0.3), Inches(0.31), Inches(12), Inches(0.46),
        size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  TOP — HORIZONTAL FSM
    #  States: HYPOTHESIZED → PARTIALLY_VALIDATED → branch → VALIDATED / RULED_OUT
    # ═══════════════════════════════════════════════════════════════════════════════
    FSM_TOP = Inches(0.88)
    SW  = Inches(2.55)   # state width
    SH  = Inches(1.95)   # state height
    GAP = Inches(0.26)   # gap between states

    states = [
        ("HYPOTHESIZED",          ACCENT_GOLD,  RGBColor(0x50,0x40,0x00),
         "Commander inferred a possible\nchain from ASG Vulnerability nodes.\nNot yet tested."),
        ("PARTIALLY\nVALIDATED",  ACCENT_LIME,  RGBColor(0x28,0x40,0x08),
         "One or more ChainSteps confirmed.\nChain not complete end-to-end.\nValidation in progress."),
        ("VALIDATED",             ACCENT_LIME,  RGBColor(0x06,0x28,0x0E),
         "All ChainSteps confirmed with\nlinked Evidence. Impact demonstrated.\nMission success for this chain."),
        ("RULED_OUT",             ACCENT_RED,   RGBColor(0x28,0x08,0x08),
         "A required ChainStep failed after\nmax retries. Chain not exploitable\nas hypothesized. Commander re-plans."),
    ]

    # First 2 states are linear, then it branches to VALIDATED and RULED_OUT
    # Layout: state0 → state1 → (fork) → state2 (top branch) and state3 (bottom branch)
    FORK_X = Inches(0.25) + 2*(SW+GAP) + SW  # x where fork begins

    # Draw state 0 and 1
    for i in range(2):
        sl = Inches(0.25) + i * (SW + GAP)
        st = FSM_TOP
        label, clr, bg, desc = states[i]
        box(slide, sl, st, SW, SH, fill=bg, line_color=clr, lw=2.0)
        # Status indicator strip at top
        box(slide, sl, st, SW, Inches(0.3), fill=clr)
        txt(slide, label, sl, st+Inches(0.04), SW, Inches(0.25), size=9.5, bold=True, color=BG_DARK)
        txt(slide, desc, sl+Inches(0.1), st+Inches(0.38), SW-Inches(0.18), SH-Inches(0.46),
            size=9, color=GREY_MID, wrap=True)
        # Arrow →
        if i == 0:
            arr(slide, sl+SW, st+SH/2, sl+SW+GAP, st+SH/2, color=ACCENT_GOLD, lw=1.8)

    # Fork point visual (small diamond-like box)
    fork_cx = Inches(0.25) + 2*SW + 2*GAP - Inches(0.18)
    fork_cy = FSM_TOP + SH/2
    box(slide, fork_cx-Inches(0.18), fork_cy-Inches(0.18), Inches(0.36), Inches(0.36),
        fill=ACCENT_GOLD, line_color=GREY_MID, lw=0.5)

    # Arrow from state1 to fork
    arr(slide, Inches(0.25)+SW+GAP+SW, FSM_TOP+SH/2, fork_cx-Inches(0.18), fork_cy, color=ACCENT_GOLD, lw=1.8)

    # VALIDATED (top-right branch)
    v_l = fork_cx + GAP
    v_t = FSM_TOP
    label, clr, bg, desc = states[2]
    box(slide, v_l, v_t, SW, SH, fill=bg, line_color=clr, lw=2.0)
    box(slide, v_l, v_t, SW, Inches(0.3), fill=clr)
    txt(slide, label, v_l, v_t+Inches(0.04), SW, Inches(0.25), size=9.5, bold=True, color=BG_DARK)
    txt(slide, desc, v_l+Inches(0.1), v_t+Inches(0.38), SW-Inches(0.18), SH-Inches(0.46),
        size=9, color=GREY_MID, wrap=True)
    arr(slide, fork_cx+Inches(0.18), fork_cy-Inches(0.1), v_l, v_t+SH/2-Inches(0.1),
        color=ACCENT_LIME, lw=1.8)
    txt(slide, "chain succeeds", fork_cx+Inches(0.22), fork_cy-Inches(0.36),
        Inches(1.2), Inches(0.22), size=8, italic=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)

    # RULED_OUT (bottom-right branch)
    ro_l = fork_cx + GAP
    ro_t = FSM_TOP + SH + Inches(0.28)
    label, clr, bg, desc = states[3]
    box(slide, ro_l, ro_t, SW, SH, fill=bg, line_color=clr, lw=2.0)
    box(slide, ro_l, ro_t, SW, Inches(0.3), fill=clr)
    txt(slide, label, ro_l, ro_t+Inches(0.04), SW, Inches(0.25), size=9.5, bold=True, color=BG_DARK)
    txt(slide, desc, ro_l+Inches(0.1), ro_t+Inches(0.38), SW-Inches(0.18), SH-Inches(0.46),
        size=9, color=GREY_MID, wrap=True)
    arr(slide, fork_cx+Inches(0.18), fork_cy+Inches(0.1), ro_l, ro_t+SH/2-Inches(0.1),
        color=ACCENT_RED, lw=1.8)
    txt(slide, "step fails after cap", fork_cx+Inches(0.22), fork_cy+Inches(0.22),
        Inches(1.35), Inches(0.22), size=8, italic=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  BOTTOM LEFT — SELF-DEBUGGING LOOP
    #  Horizontal cycle: Attempt → Diagnose → Contextualize → Adapt → [Retry | Cap→RULED_OUT]
    # ═══════════════════════════════════════════════════════════════════════════════
    LOOP_T = FSM_TOP + SH + Inches(0.26) + SH + Inches(0.26)
    LOOP_L = Inches(0.18)
    LOOP_W = Inches(7.8)
    LOOP_H = SLIDE_H - LOOP_T - Inches(0.28)

    box(slide, LOOP_L, LOOP_T, LOOP_W, LOOP_H,
        fill=RGBColor(0x08,0x0C,0x1E), line_color=ACCENT_RED, lw=1.4)
    box(slide, LOOP_L, LOOP_T, LOOP_W, Inches(0.28), fill=ACCENT_RED)
    txt(slide, "🎯  Validation Agent — Self-Debugging Loop  (on ChainStep failure)",
        LOOP_L, LOOP_T+Inches(0.04), LOOP_W, Inches(0.22), size=9, bold=True, color=BG_DARK)

    loop_steps = [
        ("ATTEMPT",        "Execute tool\nagainst target",        ACCENT_GOLD),
        ("DIAGNOSE",       "Analyze failure:\nwrong param · auth\nversion mismatch",  ACCENT_GOLD),
        ("CONTEXTUALIZE",  "Query ASG for\nadditional node\nattributes",               ACCENT_CYAN),
        ("ADAPT",          "Modify tool call\nbased on diagnosis\n+ context",           ACCENT_LIME),
    ]
    step_count = len(loop_steps)
    step_w = (LOOP_W - Inches(0.3)) / (step_count + 0.8)
    step_h = LOOP_H - Inches(0.42)
    step_t = LOOP_T + Inches(0.35)

    prev_r = None
    for i, (step_lbl, desc, clr) in enumerate(loop_steps):
        sl = LOOP_L + Inches(0.15) + i * (step_w + Inches(0.08))
        box(slide, sl, step_t, step_w, step_h, fill=RGBColor(0x10,0x14,0x28), line_color=clr, lw=1.0)
        box(slide, sl, step_t, step_w, Inches(0.24), fill=clr)
        txt(slide, step_lbl, sl, step_t+Inches(0.03), step_w, Inches(0.2),
            size=8, bold=True, color=BG_DARK)
        txt(slide, desc, sl+Inches(0.06), step_t+Inches(0.28), step_w-Inches(0.1), step_h-Inches(0.32),
            size=8, color=GREY_MID, wrap=True)
        if prev_r is not None:
            arr(slide, prev_r, step_t+step_h/2, sl, step_t+step_h/2, color=clr, lw=1.2)
        prev_r = sl + step_w

    # Retry arrow (back from ADAPT to ATTEMPT) — curved feel via two-step
    last_r = LOOP_L + Inches(0.15) + 3 * (step_w + Inches(0.08)) + step_w
    retry_y = step_t + step_h + Inches(0.06)

    # CAP box on far right
    cap_l = last_r + Inches(0.08)
    cap_w = LOOP_W - (cap_l - LOOP_L) - Inches(0.15)
    box(slide, cap_l, step_t, cap_w, step_h, fill=RGBColor(0x22,0x08,0x08), line_color=ACCENT_RED, lw=1.4)
    box(slide, cap_l, step_t, cap_w, Inches(0.24), fill=ACCENT_RED)
    txt(slide, "CAP  (×3 max)", cap_l, step_t+Inches(0.03), cap_w, Inches(0.2), size=8, bold=True, color=BG_DARK)
    arr(slide, prev_r, step_t+step_h/2, cap_l, step_t+step_h/2, color=ACCENT_RED, lw=1.2)
    txt(slide, "→ Mark ChainStep\n   RULED_OUT\n→ Write failure\n   to ASG node",
        cap_l+Inches(0.06), step_t+Inches(0.28), cap_w-Inches(0.1), step_h-Inches(0.3),
        size=8, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Loop-back annotation
    arr(slide, LOOP_L+Inches(0.15)+3*(step_w+Inches(0.08))+step_w*0.5, step_t+step_h,
        LOOP_L+Inches(0.15)+step_w*0.5, step_t+step_h,
        color=ACCENT_LIME, lw=1.0)
    txt(slide, "retry (up to cap)", LOOP_L+Inches(0.15)+step_w, step_t+step_h+Inches(0.0),
        Inches(2.5), Inches(0.2), size=7.5, italic=True, color=ACCENT_LIME)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  BOTTOM RIGHT — Risk Scoring + Priority example
    # ═══════════════════════════════════════════════════════════════════════════════
    RS_L = LOOP_L + LOOP_W + Inches(0.1)
    RS_W = SLIDE_W - RS_L - Inches(0.2)
    RS_T = LOOP_T
    RS_H = LOOP_H

    box(slide, RS_L, RS_T, RS_W, RS_H, fill=RGBColor(0x0E,0x0C,0x04), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, RS_L, RS_T, RS_W, Inches(0.28), fill=ACCENT_GOLD)
    txt(slide, "Risk Score + Chain Priority", RS_L, RS_T+Inches(0.04), RS_W, Inches(0.22),
        size=9, bold=True, color=BG_DARK)

    txt(slide, "risk_score = CVSS × Exploitability × Impact",
        RS_L+Inches(0.1), RS_T+Inches(0.34), RS_W-Inches(0.18), Inches(0.28),
        size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "Priority = rank across all HYPOTHESIZED + PARTIALLY_VALIDATED chains.\n"
        "Commander pursues highest-priority chain first — re-ranks on every status change.",
        RS_L+Inches(0.1), RS_T+Inches(0.64), RS_W-Inches(0.18), Inches(0.55),
        size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # shopvault.io chain priority table
    box(slide, RS_L+Inches(0.1), RS_T+Inches(1.22), RS_W-Inches(0.18), Inches(0.26),
        fill=RGBColor(0x22,0x18,0x04))
    txt(slide, "shopvault.io  —  APG Chain Priority Order",
        RS_L+Inches(0.14), RS_T+Inches(1.25), RS_W-Inches(0.25), Inches(0.22),
        size=8.5, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

    chain_rows = [
        ("1", "Chain-01", "CVE-2022-21661 SQLi → RCE",          "9.1", ACCENT_RED),
        ("2", "Chain-03", "Blind SQLi on staging → DB creds",    "8.1", ACCENT_GOLD),
        ("3", "Chain-02", "IDOR on /api/v1/orders → PII",        "7.5", ACCENT_GOLD),
        ("4", "Chain-04", "Direct DB backup download → PII",     "triv", ACCENT_PURP),
    ]
    for i, (rank, cid, desc, score, clr) in enumerate(chain_rows):
        rt = RS_T + Inches(1.52) + i * Inches(0.4)
        bg = RGBColor(0x18,0x12,0x02) if i % 2 == 0 else RGBColor(0x12,0x0C,0x02)
        box(slide, RS_L+Inches(0.1), rt, RS_W-Inches(0.18), Inches(0.38), fill=bg, line_color=clr, lw=0.5)
        txt(slide, rank, RS_L+Inches(0.14), rt+Inches(0.06), Inches(0.22), Inches(0.26),
            size=9, bold=True, color=clr)
        txt(slide, f"{cid}  —  {desc}", RS_L+Inches(0.38), rt+Inches(0.06),
            RS_W-Inches(0.75), Inches(0.26), size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT)
        txt(slide, score, RS_L+RS_W-Inches(0.55), rt+Inches(0.06), Inches(0.42), Inches(0.26),
            size=9.5, bold=True, color=clr, align=PP_ALIGN.RIGHT)

    txt(slide, "09", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.52),
        Inches(0.35), Inches(0.42), size=13, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)


