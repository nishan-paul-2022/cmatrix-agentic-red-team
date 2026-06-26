"""
Slide 2 — The Problem CMatrix Solves
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)
CARD_RED    = RGBColor(0x2A, 0x10, 0x10)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

# ── Helpers ────────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=1):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_para(tf, text, size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

# ── Load existing PPTX and add slide ─────────────────────────────────────────
prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)

# ── Left accent bar ───────────────────────────────────────────────────────────
add_rect(slide, l=Inches(0), t=Inches(0), w=Inches(0.06), h=SLIDE_H, fill_color=ACCENT_RED)
add_rect(slide, l=Inches(0.06), t=Inches(0), w=SLIDE_W - Inches(0.06), h=Inches(0.04), fill_color=ACCENT_RED)
add_rect(slide, l=Inches(0.06), t=SLIDE_H - Inches(0.04), w=SLIDE_W - Inches(0.06), h=Inches(0.04), fill_color=ACCENT_RED)

# ── Slide label ───────────────────────────────────────────────────────────────
add_text(slide, "THE PROBLEM", Inches(0.3), Inches(0.1), Inches(6), Inches(0.35),
         size=11, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)

# ── Section title ─────────────────────────────────────────────────────────────
add_text(slide, "What Existing VAPT Systems Get Wrong",
         Inches(0.3), Inches(0.45), Inches(12), Inches(0.75),
         size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Divider ───────────────────────────────────────────────────────────────────
add_rect(slide, Inches(0.3), Inches(1.22), Inches(6), Inches(0.035), fill_color=ACCENT_RED)

# ── LEFT COLUMN — "Current State" problems (4 cards) ─────────────────────────
problems = [
    ("No Structured World Model",
     "Systems reason from flat conversation history or task queues — no graph of what the target actually is."),
    ("No Attack Path Reasoning",
     "Finished tasks are logged, not understood. Systems cannot represent which vulnerabilities chain together or why."),
    ("Fragile Re-Planning",
     "Dynamic re-planning is ad-hoc. There is no formal trigger grounded in discovered evidence."),
    ("Arbitrary Termination",
     "Missions end on time-limits or empty queues — not because the attack surface and all chains are truly exhausted."),
]

card_top  = Inches(1.45)
card_h    = Inches(1.25)
card_gap  = Inches(0.1)
card_l    = Inches(0.3)
card_w    = Inches(5.8)

for i, (title, body) in enumerate(problems):
    top = card_top + i * (card_h + card_gap)
    add_rect(slide, card_l, top, card_w, card_h, fill_color=CARD_RED,
             line_color=ACCENT_RED, line_width_pt=1)
    # red bullet
    add_text(slide, "✕", card_l + Inches(0.12), top + Inches(0.1),
             Inches(0.4), Inches(0.5), size=20, bold=True, color=ACCENT_RED)
    add_text(slide, title,
             card_l + Inches(0.48), top + Inches(0.08),
             card_w - Inches(0.6), Inches(0.38),
             size=14, bold=True, color=WHITE)
    add_text(slide, body,
             card_l + Inches(0.48), top + Inches(0.44),
             card_w - Inches(0.6), Inches(0.72),
             size=11.5, bold=False, color=GREY_MID, wrap=True)

# ── RIGHT COLUMN — quote block + fundamental failure statement ────────────────
# Big quote
quote_l = Inches(6.5)
quote_t = Inches(1.45)
quote_w = Inches(6.5)

add_rect(slide, quote_l, quote_t, quote_w, Inches(2.55),
         fill_color=CARD_BG, line_color=ACCENT_CYAN, line_width_pt=1.5)

add_text(slide, "\u201c",
         quote_l + Inches(0.15), quote_t + Inches(0.05),
         Inches(0.6), Inches(0.7),
         size=60, bold=True, color=ACCENT_CYAN)

add_text(slide,
         "They automate tool execution.\nCMatrix automates the reasoning\nof a professional penetration tester.",
         quote_l + Inches(0.2), quote_t + Inches(0.55),
         quote_w - Inches(0.4), Inches(1.6),
         size=20, bold=False, italic=True, color=WHITE, wrap=True)

add_text(slide, "— CMatrix Design Philosophy",
         quote_l + Inches(0.2), quote_t + Inches(2.2),
         quote_w - Inches(0.4), Inches(0.35),
         size=12, bold=False, color=GREY_MID)

# ── Root Cause Label ──────────────────────────────────────────────────────────
add_rect(slide, quote_l, Inches(4.15), quote_w, Inches(0.4),
         fill_color=ACCENT_RED)
add_text(slide, "⚠  ROOT CAUSE",
         quote_l + Inches(0.2), Inches(4.18), quote_w - Inches(0.3), Inches(0.35),
         size=13, bold=True, color=WHITE)

add_rect(slide, quote_l, Inches(4.55), quote_w, Inches(1.45),
         fill_color=CARD_BG, line_color=ACCENT_RED, line_width_pt=1)

root_cause_lines = [
    "Existing systems have no structured model of the target environment",
    "and no structured model of what attack paths are possible.",
    "",
    "They know what they did — not what the target is, or what can be done to it.",
]
tb = slide.shapes.add_textbox(quote_l + Inches(0.2), Inches(4.62),
                               quote_w - Inches(0.35), Inches(1.3))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(root_cause_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    run = p.add_run()
    run.text = line
    run.font.name = "Calibri"
    run.font.size = Pt(13)
    run.font.color.rgb = GREY_MID if line else WHITE

# ── CMatrix Solution teaser ───────────────────────────────────────────────────
add_rect(slide, Inches(0.3), Inches(6.22), SLIDE_W - Inches(0.6), Inches(0.85),
         fill_color=RGBColor(0x05, 0x14, 0x20), line_color=ACCENT_CYAN, line_width_pt=1)
add_text(slide,
         "CMatrix Solution →  A dual-graph world model: ASG captures discovered reality · APG captures inferred attack opportunity",
         Inches(0.5), Inches(6.3), SLIDE_W - Inches(1.0), Inches(0.65),
         size=14, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

# ── Slide number ──────────────────────────────────────────────────────────────
add_text(slide, "02", SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.55),
         Inches(0.7), Inches(0.45), size=13, bold=True, color=ACCENT_RED,
         align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 2 added.")
