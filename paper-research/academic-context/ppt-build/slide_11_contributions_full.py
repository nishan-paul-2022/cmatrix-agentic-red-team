"""
Slide 11 — Research Contributions Summary (C1 – C12)
Full table of all 12 novel contributions with short labels.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome ─────────────────────────────────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_CYAN)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_CYAN)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_CYAN)

txt(slide, "RESEARCH CONTRIBUTIONS", Inches(0.3), Inches(0.06), Inches(8), Inches(0.3),
    size=11, bold=True, color=ACCENT_CYAN)
txt(slide, "12 Novel Contributions — What CMatrix Adds to the Literature",
    Inches(0.3), Inches(0.36), Inches(12.5), Inches(0.6), size=28, bold=True, color=WHITE)
box(slide, Inches(0.3), Inches(0.98), Inches(7), Inches(0.03), fill=ACCENT_CYAN)

# ── 12 contribution rows in 2 columns ─────────────────────────────────────────
contributions = [
    ("C1",  ACCENT_LIME,  "Dual-Graph World Model",
     "ASG (discovered reality) + APG (inferred opportunity) — strict write ownership. No prior system separates these."),
    ("C2",  ACCENT_CYAN,  "Graph-State-Driven Re-Planning",
     "Re-planning triggered by explicit ASG/APG state changes, not fixed task-queue depletion."),
    ("C3",  ACCENT_GOLD,  "APG Attack Chain Lifecycle + Evidence Traceability",
     "Chains as first-class entities: risk-scored, lifecycle-tracked, every ChainStep linked to ASG Evidence."),
    ("C4",  ACCENT_TEAL,  "ASG-Aware Parallel Tool Dispatch",
     "Dependency-safe concurrent tool execution using the ASG as the dependency graph."),
    ("C5",  ACCENT_RED,   "Tool Risk Gate + Commander Mailbox",
     "Every tool call risk-classified before execution. High-risk ops require Commander approval."),
    ("C6",  ACCENT_PURP,  "ASG-Backed Lossless Context Compaction",
     "FullCompact compresses conversation history to near-zero without losing any findings — all live in the graph."),
    ("C7",  ACCENT_CYAN,  "Methodology-as-Configuration",
     "Commander's planning policy encoded as a versioned VAPT Protocol Prompt — swap methodologies without code changes."),
    ("C8",  ACCENT_LIME,  "Dual-Graph Termination Semantics",
     "Mission terminates when ASG is exhausted AND all APG chains reach terminal status — simultaneously."),
    ("C9",  ACCENT_GOLD,  "Live Vulnerability Intelligence Grounding",
     "Research Agent provides real-time CVE enrichment, PoC availability, exploit feasibility during active assessment."),
    ("C10", ACCENT_TEAL,  "Cross-Mission Experience Store",
     "Persistent RAG-backed store of validated exploitation outcomes queried at mission start to seed chain hypotheses."),
    ("C11", ACCENT_GOLD,  "Attack Strategy Library with Crystallization",
     "Generalizes validated chains into named, parameterized strategies indexed by technology fingerprint + confidence score."),
    ("C12", ACCENT_PURP,  "Structured Engagement Trajectory Export",
     "Complete machine-readable decision log per mission — enables reproducibility, ablation studies, and dataset generation."),
]

# 2 columns, 6 rows each
col_defs = [(Inches(0.3), Inches(6.45)), (Inches(6.95), Inches(6.2))]
row_h = Inches(0.49)
row_gap = Inches(0.055)
start_t = Inches(1.08)

for col_idx, (col_l, col_w) in enumerate(col_defs):
    for row_idx in range(6):
        item_idx = col_idx * 6 + row_idx
        cnum, clr, title, desc = contributions[item_idx]
        t = start_t + row_idx * (row_h + row_gap)

        # Row background
        bg = RGBColor(0x0C, 0x14, 0x26) if row_idx % 2 == 0 else CARD_BG
        box(slide, col_l, t, col_w, row_h, fill=bg, line_color=clr, lw=0.7)

        # Contribution ID badge
        badge_w = Inches(0.48)
        box(slide, col_l, t, badge_w, row_h, fill=clr)
        txt(slide, cnum, col_l + Inches(0.02), t + row_h/2 - Inches(0.15),
            badge_w - Inches(0.04), Inches(0.3),
            size=10, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        # Title
        txt(slide, title,
            col_l + badge_w + Inches(0.1), t + Inches(0.04),
            col_w - badge_w - Inches(0.15), Inches(0.22),
            size=10.5, bold=True, color=WHITE)

        # Description
        txt(slide, desc,
            col_l + badge_w + Inches(0.1), t + Inches(0.27),
            col_w - badge_w - Inches(0.15), Inches(0.2),
            size=8.8, bold=False, color=GREY_MID, wrap=True)

# ── Vertical divider ──────────────────────────────────────────────────────────
box(slide, Inches(6.82), start_t, Inches(0.03), 6*(row_h+row_gap) - row_gap, fill=GREY_MID)

# ── Bottom tagline ─────────────────────────────────────────────────────────────
bottom_t = start_t + 6 * (row_h + row_gap) + Inches(0.1)
box(slide, Inches(0.3), bottom_t, SLIDE_W - Inches(0.6), Inches(0.58),
    fill=RGBColor(0x06, 0x12, 0x24), line_color=ACCENT_CYAN, lw=1.2)
txt(slide,
    "Each contribution is independently evaluable and independently publishable.  "
    "Together they form a unified architecture that no existing VAPT system implements.",
    Inches(0.55), bottom_t + Inches(0.1), SLIDE_W - Inches(1.0), Inches(0.38),
    size=13, bold=False, italic=True, color=ACCENT_CYAN, wrap=True)

txt(slide, "11", SLIDE_W - Inches(0.4), SLIDE_H - Inches(0.55),
    Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 11 added.")
