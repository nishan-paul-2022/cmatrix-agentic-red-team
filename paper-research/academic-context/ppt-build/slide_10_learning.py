"""
Slide 10 — Cross-Mission Learning (Experience Store + Attack Strategy Library)
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.enum.shapes

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
ACCENT_TEAL = RGBColor(0x00, 0xBF, 0xD8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

def arrow_h(slide, x1, y, x2, color=GREY_MID, lw=1.5):
    c = slide.shapes.add_connector(pptx.enum.shapes.MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color; c.line.width = Pt(lw)
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    etree.SubElement(ln, qn('a:tailEnd')).set('type', 'none')

prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome ─────────────────────────────────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_TEAL)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_TEAL)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_TEAL)

txt(slide, "CROSS-MISSION LEARNING", Inches(0.3), Inches(0.08), Inches(8), Inches(0.3),
    size=11, bold=True, color=ACCENT_TEAL)
txt(slide, "Experience Store + Attack Strategy Library — CMatrix Gets Smarter Each Mission",
    Inches(0.3), Inches(0.38), Inches(12.5), Inches(0.62), size=26, bold=True, color=WHITE)
box(slide, Inches(0.3), Inches(1.0), Inches(7), Inches(0.03), fill=ACCENT_TEAL)

# ══════════════════════════════════════════════
#  FLOW DIAGRAM: Mission N → Store → Mission N+1
# ══════════════════════════════════════════════
flow_t = Inches(1.1)
flow_h = Inches(1.0)

# Mission N box
box(slide, Inches(0.3), flow_t, Inches(2.5), flow_h,
    fill=RGBColor(0x06,0x14,0x20), line_color=ACCENT_CYAN, lw=1.5)
txt(slide, "Mission N\n(completed)", Inches(0.45), flow_t+Inches(0.15),
    Inches(2.2), Inches(0.7), size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

# Arrow → Report Agent writes
arrow_h(slide, Inches(2.8), flow_t+flow_h/2, Inches(4.0), ACCENT_CYAN)
txt(slide, "Report Agent writes\nvalidated chains",
    Inches(2.82), flow_t+Inches(0.1), Inches(1.12), Inches(0.75),
    size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# Cross-Mission Experience Store box
box(slide, Inches(4.1), flow_t-Inches(0.3), Inches(2.6), flow_h+Inches(0.6),
    fill=RGBColor(0x04,0x20,0x1E), line_color=ACCENT_TEAL, lw=2.0)
txt(slide, "Cross-Mission\nExperience Store",
    Inches(4.22), flow_t-Inches(0.18), Inches(2.35), Inches(0.5),
    size=12, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
txt(slide, "RAG-backed · persistent\nacross all missions",
    Inches(4.22), flow_t+Inches(0.35), Inches(2.35), Inches(0.4),
    size=9, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# Arrow → Crystallize (if 2+ same fingerprint)
arrow_h(slide, Inches(6.7), flow_t+flow_h/2, Inches(7.9), ACCENT_GOLD)
txt(slide, "Crystallize\n(>=2 missions)",
    Inches(6.72), flow_t+Inches(0.05), Inches(1.12), Inches(0.75),
    size=8.5, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

# Attack Strategy Library box
box(slide, Inches(8.0), flow_t-Inches(0.3), Inches(2.6), flow_h+Inches(0.6),
    fill=RGBColor(0x20,0x18,0x04), line_color=ACCENT_GOLD, lw=2.0)
txt(slide, "Attack Strategy\nLibrary",
    Inches(8.12), flow_t-Inches(0.18), Inches(2.35), Inches(0.5),
    size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
txt(slide, "Named strategies · confidence scores\ntechnology-fingerprint indexed",
    Inches(8.12), flow_t+Inches(0.35), Inches(2.35), Inches(0.4),
    size=9, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# Arrow → Mission N+1
arrow_h(slide, Inches(10.6), flow_t+flow_h/2, Inches(11.8), ACCENT_LIME)
txt(slide, "Pre-ranked seeds\nat mission start",
    Inches(10.62), flow_t+Inches(0.1), Inches(1.12), Inches(0.75),
    size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# Mission N+1 box
box(slide, Inches(11.9), flow_t, Inches(1.25), flow_h,
    fill=RGBColor(0x06,0x20,0x10), line_color=ACCENT_LIME, lw=1.5)
txt(slide, "Mission\nN+1", Inches(11.95), flow_t+Inches(0.12),
    Inches(1.15), Inches(0.75), size=12, bold=True, color=ACCENT_LIME, align=PP_ALIGN.CENTER)

# Also: Experience Store feeds Mission N+1
arrow_h(slide, Inches(5.4), flow_t+flow_h*0.9, Inches(11.9), ACCENT_TEAL)
txt(slide, "Candidate chain hypotheses injected into Commander context",
    Inches(6.6), flow_t+flow_h*0.9+Inches(0.04), Inches(4.0), Inches(0.28),
    size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
#  TWO CARDS BELOW: Experience Store vs Library
# ══════════════════════════════════════════════
card_t = Inches(2.42)
card_h = Inches(4.55)

# Experience Store card
es_l = Inches(0.3); es_w = Inches(6.2)
box(slide, es_l, card_t, es_w, card_h,
    fill=RGBColor(0x04,0x18,0x18), line_color=ACCENT_TEAL, lw=1.5)
box(slide, es_l, card_t, es_w, Inches(0.4), fill=ACCENT_TEAL)
txt(slide, "C10  ·  Cross-Mission Experience Store",
    es_l+Inches(0.12), card_t+Inches(0.06), es_w-Inches(0.2), Inches(0.28),
    size=12, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

es_items = [
    ("What it stores",
     "Per-mission · per-chain raw records of validated exploitation outcomes.\n"
     "Target technology fingerprint · CVE · successful tool invocation · ChainStep sequence · mission outcome."),
    ("When it's written",
     "Report Agent writes one entry for every chain with terminal status VALIDATED at mission close."),
    ("When it's read",
     "Commander queries immediately after Recon Agent writes first Technology nodes to ASG — "
     "before Analysis begins. Retrieved records become candidate chain hypotheses."),
    ("Research origin",
     "Generalizes AutoAttacker's within-mission experience-reuse mechanism to cross-mission scope. "
     "AutoAttacker reuses subtasks within one session; CMatrix accumulates across every mission ever run."),
]
for i, (heading, body) in enumerate(es_items):
    t = card_t + Inches(0.5) + i * Inches(0.98)
    box(slide, es_l+Inches(0.12), t, es_w-Inches(0.24), Inches(0.94),
        fill=RGBColor(0x06,0x22,0x22), line_color=ACCENT_TEAL, lw=0.5)
    txt(slide, heading, es_l+Inches(0.22), t+Inches(0.06), es_w-Inches(0.38), Inches(0.26),
        size=11, bold=True, color=ACCENT_TEAL)
    txt(slide, body, es_l+Inches(0.22), t+Inches(0.34), es_w-Inches(0.38), Inches(0.56),
        size=9.5, color=GREY_MID, wrap=True)

# Attack Strategy Library card
sl_l = Inches(6.75); sl_w = Inches(6.38)
box(slide, sl_l, card_t, sl_w, card_h,
    fill=RGBColor(0x1C,0x14,0x04), line_color=ACCENT_GOLD, lw=1.5)
box(slide, sl_l, card_t, sl_w, Inches(0.4), fill=ACCENT_GOLD)
txt(slide, "C11  ·  Attack Strategy Library",
    sl_l+Inches(0.12), card_t+Inches(0.06), sl_w-Inches(0.2), Inches(0.28),
    size=12, bold=True, color=BG_DARK)

sl_items = [
    ("What it stores",
     "Generalized attack procedures distilled from multiple validated outcomes against the same "
     "target technology class. Named strategies (e.g. STRAT-WP-SQLI-001) with confidence scores."),
    ("Crystallization threshold",
     "Same technology fingerprint (e.g. WordPress 5.x + WooCommerce + Nginx) produces validated "
     "AttackChain in 2 or more independent missions → Commander triggers crystallization."),
    ("When it's used",
     "Retrieved at mission start. Strategies are injected as pre-ranked APG AttackChain seeds — "
     "prioritized above zero-prior chains because they carry validated track record."),
    ("Why it matters",
     "No existing autonomous VAPT system accumulates and generalizes validated exploitation "
     "procedures across sessions. Every prior system resets to zero knowledge on each mission. "
     "CMatrix's library makes it measurably more efficient on repeat target-type engagements."),
]
for i, (heading, body) in enumerate(sl_items):
    t = card_t + Inches(0.5) + i * Inches(0.98)
    box(slide, sl_l+Inches(0.12), t, sl_w-Inches(0.24), Inches(0.94),
        fill=RGBColor(0x22,0x18,0x06), line_color=ACCENT_GOLD, lw=0.5)
    txt(slide, heading, sl_l+Inches(0.22), t+Inches(0.06), sl_w-Inches(0.38), Inches(0.26),
        size=11, bold=True, color=ACCENT_GOLD)
    txt(slide, body, sl_l+Inches(0.22), t+Inches(0.34), sl_w-Inches(0.38), Inches(0.56),
        size=9.5, color=GREY_MID, wrap=True)

txt(slide, "10", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.55),
    Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 10 added.")
