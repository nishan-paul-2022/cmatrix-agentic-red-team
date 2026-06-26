"""
Slide 5 — Dual-Graph World Model (Deep Dive — Visual Graph Diagram)
====================================================================
Instead of two table columns, this slide draws an ACTUAL graph:

LEFT  (ASG):  A miniature knowledge graph with real node-types shown as
              labelled circles/rounded-rects connected with typed edges.
              Sample: Domain → Host → Port → Service → Vulnerability → Evidence

RIGHT (APG):  An attack-chain visual: AttackChain containing sequenced
              ChainStep nodes (numbered) with risk_score + validation_status.
              Final ChainStep → Impact node.

CENTRE:       A vertical "STRICT SEPARATION" barrier with:
              Discovery Agents → write ASG only (left arrow)
              Commander Agent  → writes APG only (right arrow)

Bottom:       The architectural principle banner.
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.enum.shapes

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
ACCENT_RED  = RGBColor(0xFF, 0x45, 0x45)
ACCENT_PURP = RGBColor(0xBD, 0x93, 0xF9)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
GREY_DARK   = RGBColor(0x30, 0x38, 0x48)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else: shp.fill.background()
    if line_color: shp.line.color.rgb = line_color; shp.line.width = Pt(lw)
    else: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def txt(slide, text, l, t, w, h, size=10, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = "Calibri"; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

def connector(slide, x1, y1, x2, y2, color=GREY_MID, lw=1.2, label="", bidirectional=False):
    c = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color; c.line.width = Pt(lw)
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'sm'); he.set('len', 'med')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    if bidirectional:
        te.set('type', 'arrow'); te.set('w', 'sm'); te.set('len', 'med')
    else:
        te.set('type', 'none')
    # Label near midpoint
    if label:
        mx = (x1 + x2) / 2; my = (y1 + y2) / 2
        txt(slide, label, mx - Inches(0.55), my - Inches(0.22),
            Inches(1.1), Inches(0.22), size=7, italic=True, color=GREY_MID)
    return c

def node(slide, cx, cy, w, h, label, sublabel="", fill_col=None, border_col=None, label_size=9):
    """Draw a graph node (rectangle) centred at cx, cy."""
    fc = fill_col or RGBColor(0x10, 0x18, 0x2C)
    bc = border_col or ACCENT_CYAN
    l = cx - w/2; t = cy - h/2
    box(slide, l, t, w, h, fill=fc, line_color=bc, lw=1.4)
    # Main label
    ty = t + Inches(0.04) if sublabel else t
    th = h - Inches(0.06)
    txt(slide, label, l, ty, w, th if not sublabel else h/2,
        size=label_size, bold=True, color=bc)
    if sublabel:
        txt(slide, sublabel, l, cy, w, h/2 - Inches(0.04),
            size=7, color=GREY_MID, italic=True)

# ── Load + add slide ──────────────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

# ── Chrome (lime accent this slide) ──────────────────────────────────────────
box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_LIME)
box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_LIME)
box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_LIME)

# ── Title ─────────────────────────────────────────────────────────────────────
txt(slide, "DUAL-GRAPH WORLD MODEL", Inches(0.3), Inches(0.06), Inches(6), Inches(0.26),
    size=10, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
txt(slide, "Two Strictly Separated Knowledge Layers — Visualised",
    Inches(0.3), Inches(0.32), Inches(11), Inches(0.48),
    size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# ── Two zone backgrounds ──────────────────────────────────────────────────────
# ASG zone (left)
box(slide, Inches(0.18), Inches(0.88), Inches(5.85), Inches(5.62),
    fill=RGBColor(0x04,0x18,0x0A), line_color=ACCENT_LIME, lw=1.5)
txt(slide, "ASG — Attack Surface Graph", Inches(0.22), Inches(0.9),
    Inches(3.5), Inches(0.26), size=11, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
txt(slide, "What does the target look like?  (Discovered Reality)",
    Inches(0.22), Inches(1.14), Inches(5.5), Inches(0.24),
    size=9, italic=True, color=RGBColor(0x60,0xC0,0x70), align=PP_ALIGN.LEFT)

# APG zone (right)
box(slide, Inches(7.3), Inches(0.88), Inches(5.85), Inches(5.62),
    fill=RGBColor(0x1C,0x10,0x04), line_color=ACCENT_GOLD, lw=1.5)
txt(slide, "APG — Attack Path Graph", Inches(7.34), Inches(0.9),
    Inches(3.5), Inches(0.26), size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
txt(slide, "What can be done to it?  (Inferred Opportunity)",
    Inches(7.34), Inches(1.14), Inches(5.5), Inches(0.24),
    size=9, italic=True, color=RGBColor(0xC0,0x90,0x20), align=PP_ALIGN.LEFT)

# ── SEPARATION BARRIER (centre) ──────────────────────────────────────────────
sep_cx = Inches(6.667)
box(slide, sep_cx - Inches(0.55), Inches(0.88), Inches(1.1), Inches(5.62),
    fill=RGBColor(0x08,0x10,0x18), line_color=GREY_DARK, lw=0.8)
# Dashed line
box(slide, sep_cx - Inches(0.01), Inches(0.88), Inches(0.02), Inches(5.62), fill=GREY_DARK)
# Labels
txt(slide, "STRICT\nSEPARATION", sep_cx - Inches(0.52), Inches(2.5),
    Inches(1.04), Inches(0.7), size=8, bold=True, color=GREY_MID)
txt(slide, "No agent\ncrosses this\nboundary", sep_cx - Inches(0.5), Inches(3.3),
    Inches(1.0), Inches(0.75), size=7.5, italic=True, color=GREY_MID)
# Write arrows (pointing at each zone)
connector(slide, sep_cx - Inches(0.04), Inches(4.25),
          Inches(4.5), Inches(4.25), color=ACCENT_LIME, lw=1.2, label="Discovery\nAgents write")
connector(slide, sep_cx + Inches(0.04), Inches(4.6),
          Inches(8.5), Inches(4.6), color=ACCENT_GOLD, lw=1.2, label="Commander\nwrites")

# ═══════════════════════════════════════════════════════════════════════════════
#  ASG GRAPH — visual node-edge diagram
#  Chain: Domain → Host → Port → Service → Vulnerability → Evidence
#  Branch: Service → Technology; Host → Endpoint → Parameter
# ═══════════════════════════════════════════════════════════════════════════════
NW = Inches(1.22)  # node width
NH = Inches(0.44)  # node height

# Node positions (cx, cy)
ASG_NODES = {
    "Domain":        (Inches(1.5),  Inches(1.72)),
    "Host":          (Inches(1.5),  Inches(2.42)),
    "Port":          (Inches(1.5),  Inches(3.12)),
    "Service":       (Inches(1.5),  Inches(3.82)),
    "Technology":    (Inches(3.1),  Inches(3.12)),
    "Endpoint":      (Inches(3.1),  Inches(3.82)),
    "Parameter":     (Inches(3.1),  Inches(4.52)),
    "Vulnerability": (Inches(1.5),  Inches(4.52)),
    "Evidence":      (Inches(1.5),  Inches(5.22)),
}

ASG_COLORS = {
    "Domain":        (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
    "Host":          (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
    "Port":          (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
    "Service":       (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
    "Technology":    (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
    "Endpoint":      (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
    "Parameter":     (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
    "Vulnerability": (RGBColor(0x22,0x06,0x06), ACCENT_RED),
    "Evidence":      (RGBColor(0x12,0x08,0x20), ACCENT_PURP),
}

ASG_SUBLABELS = {
    "Domain":        "shopvault.io",
    "Host":          "192.168.1.10",
    "Port":          ":443 / :8080",
    "Service":       "Nginx 1.18.0",
    "Technology":    "WordPress 5.9.3",
    "Endpoint":      "/api/v1/orders",
    "Parameter":     "user_id=?",
    "Vulnerability": "CVE-2022-21661",
    "Evidence":      "screenshot.png",
}

for name, (cx, cy) in ASG_NODES.items():
    fc, bc = ASG_COLORS[name]
    node(slide, cx, cy, NW, NH, name, ASG_SUBLABELS[name], fc, bc, label_size=8.5)

# ASG Edges (vertical spine)
vertical_chain = ["Domain", "Host", "Port", "Service", "Vulnerability", "Evidence"]
for i in range(len(vertical_chain) - 1):
    src = ASG_NODES[vertical_chain[i]]
    dst = ASG_NODES[vertical_chain[i+1]]
    edge_labels = {
        ("Domain","Host"): "has_host",
        ("Host","Port"): "has_port",
        ("Port","Service"): "runs",
        ("Service","Vulnerability"): "affected_by",
        ("Vulnerability","Evidence"): "validated_by",
    }
    lbl = edge_labels.get((vertical_chain[i], vertical_chain[i+1]), "")
    connector(slide, src[0], src[1]+NH/2, dst[0], dst[1]-NH/2,
              color=ACCENT_LIME, lw=0.9, label=lbl)

# Branch edges: Port → Technology, Port → Endpoint, Endpoint → Parameter
connector(slide, ASG_NODES["Port"][0]+NW/2, ASG_NODES["Port"][1],
          ASG_NODES["Technology"][0]-NW/2, ASG_NODES["Technology"][1],
          color=ACCENT_CYAN, lw=0.8, label="uses")
connector(slide, ASG_NODES["Service"][0]+NW/2, ASG_NODES["Service"][1],
          ASG_NODES["Endpoint"][0]-NW/2, ASG_NODES["Endpoint"][1],
          color=ACCENT_CYAN, lw=0.8, label="has_endpoint")
connector(slide, ASG_NODES["Endpoint"][0], ASG_NODES["Endpoint"][1]+NH/2,
          ASG_NODES["Parameter"][0], ASG_NODES["Parameter"][1]-NH/2,
          color=ACCENT_CYAN, lw=0.8, label="has_parameter")

# ── ASG Edge legend strip ─────────────────────────────────────────────────────
box(slide, Inches(0.22), Inches(5.72), Inches(5.78), Inches(0.62),
    fill=RGBColor(0x06,0x1C,0x0C), line_color=ACCENT_LIME, lw=0.6)
txt(slide, "EDGES:  has_host  ·  has_port  ·  runs  ·  uses  ·  has_endpoint  ·  "
    "has_parameter  ·  affected_by  ·  validated_by",
    Inches(0.3), Inches(5.78), Inches(5.6), Inches(0.46),
    size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
#  APG GRAPH — attack chain visual
# ═══════════════════════════════════════════════════════════════════════════════
# AttackChain container
ac_l, ac_t = Inches(7.42), Inches(1.44)
ac_w, ac_h = Inches(5.5), Inches(4.7)
box(slide, ac_l, ac_t, ac_w, ac_h, fill=RGBColor(0x14,0x0C,0x02), line_color=ACCENT_GOLD, lw=1.4)
box(slide, ac_l, ac_t, ac_w, Inches(0.28), fill=ACCENT_GOLD)
txt(slide, "AttackChain  ·  risk_score: 9.1  ·  status: VALIDATED",
    ac_l, ac_t+Inches(0.04), ac_w, Inches(0.22), size=8, bold=True, color=BG_DARK)
txt(slide, "Chain-01: CVE-2022-21661 → SQL Injection → RCE → Customer PII",
    ac_l+Inches(0.12), ac_t+Inches(0.32), ac_w-Inches(0.2), Inches(0.22),
    size=8.5, bold=False, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

# ChainStep nodes (vertical)
steps = [
    ("STEP 1", "SQLMap → WP_Query\nSQLi confirmed",         ACCENT_GOLD,    "VALIDATED"),
    ("STEP 2", "Admin hash extracted\n+ cracked offline",   ACCENT_GOLD,    "VALIDATED"),
    ("STEP 3", "Metasploit → Web shell\nRCE achieved",      ACCENT_RED,     "VALIDATED"),
    ("IMPACT", "Full server access\nCustomer PII exposed",  ACCENT_PURP,    "DEMONSTRATED"),
]
snw = Inches(2.1); snh = Inches(0.65)
s_start_t = ac_t + Inches(0.65)
s_cx = ac_l + ac_w/2

for i, (step_lbl, detail, clr, status) in enumerate(steps):
    sy = s_start_t + i * (snh + Inches(0.22))
    sl = s_cx - snw/2
    is_impact = (step_lbl == "IMPACT")
    bg = RGBColor(0x16,0x08,0x20) if is_impact else RGBColor(0x1E,0x14,0x02)
    box(slide, sl, sy, snw, snh, fill=bg, line_color=clr, lw=1.6)
    # Step label + status pill on right
    txt(slide, step_lbl, sl+Inches(0.1), sy+Inches(0.04), Inches(0.8), Inches(0.28),
        size=8.5, bold=True, color=clr, align=PP_ALIGN.LEFT)
    # Status badge
    badge_w = Inches(1.05)
    badge_color = ACCENT_LIME if status == "VALIDATED" else ACCENT_PURP
    box(slide, sl+snw-badge_w-Inches(0.06), sy+Inches(0.06), badge_w, Inches(0.2), fill=badge_color)
    txt(slide, status, sl+snw-badge_w-Inches(0.06), sy+Inches(0.07),
        badge_w, Inches(0.18), size=6.5, bold=True, color=BG_DARK)
    # Detail
    txt(slide, detail, sl+Inches(0.1), sy+Inches(0.3), snw-Inches(0.16), Inches(0.34),
        size=8, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)
    # Connector to next step
    if i < len(steps) - 1:
        connector(slide, s_cx, sy+snh, s_cx, sy+snh+Inches(0.22), color=clr, lw=1.2)
        txt(slide, "next_step", s_cx+Inches(0.06), sy+snh+Inches(0.02),
            Inches(0.8), Inches(0.2), size=6.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)
    # supported_by label (right of each step)
    txt(slide, "↗ supported_by\nASG Evidence",
        sl+snw+Inches(0.08), sy+Inches(0.15), Inches(0.85), Inches(0.36),
        size=6.5, italic=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

# starts_at label (top)
txt(slide, "starts_at: ASG Vulnerability node",
    ac_l+Inches(0.12), ac_t+Inches(0.54), ac_w-Inches(0.2), Inches(0.2),
    size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

# ── APG Edge legend strip ─────────────────────────────────────────────────────
box(slide, Inches(7.34), Inches(5.72), Inches(5.78), Inches(0.62),
    fill=RGBColor(0x1C,0x10,0x02), line_color=ACCENT_GOLD, lw=0.6)
txt(slide, "EDGES:  starts_at  ·  next_step  ·  achieves  ·  supported_by\n"
    "VALIDATION STATUS:  HYPOTHESIZED → PARTIALLY_VALIDATED → VALIDATED / RULED_OUT",
    Inches(7.42), Inches(5.78), Inches(5.6), Inches(0.5),
    size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

# ── Bottom principle bar ──────────────────────────────────────────────────────
box(slide, Inches(0.18), Inches(6.42), SLIDE_W-Inches(0.36), Inches(0.72),
    fill=RGBColor(0x06,0x14,0x22), line_color=ACCENT_CYAN, lw=1.5)
txt(slide,
    "Separation Principle:  Discovery agents write only to the ASG — they never reason about chains.  "
    "The Commander writes only to the APG — it never runs tools.  "
    "No agent conflates discovered facts with hypothesised attack reasoning.",
    Inches(0.5), Inches(6.5), SLIDE_W-Inches(0.9), Inches(0.58),
    size=11, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT, wrap=True)

txt(slide, "05", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.52),
    Inches(0.35), Inches(0.42), size=13, bold=True, color=ACCENT_LIME, align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 5 (Dual Graph) rewritten.")
