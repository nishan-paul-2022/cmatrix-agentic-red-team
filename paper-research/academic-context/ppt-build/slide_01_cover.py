"""
Slide 1 — Cover Slide
CMatrix Research Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import os, copy

# ── Palette ────────────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0A, 0x0D, 0x1A)   # near-black navy
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan
ACCENT_LIME  = RGBColor(0x39, 0xFF, 0x14)   # neon green
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID     = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG      = RGBColor(0x10, 0x16, 0x2B)   # slightly lighter dark

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
PPTX_PATH = os.path.join(OUT_DIR, "CMatrix_Presentation.pptx")

# ── Helpers ────────────────────────────────────────────────────────────────────

def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

def set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE=1
        l, t, w, h
    )
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(line_width_pt)
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_name="Calibri", font_size=24,
                 bold=False, italic=False,
                 color=WHITE, align=PP_ALIGN.LEFT,
                 word_wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_multiline_text(slide, lines, l, t, w, h,
                       font_name="Calibri", font_size=18,
                       bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                       line_spacing_pt=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bld, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing_pt:
            p.line_spacing = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(sz)
        run.font.bold = bld
        run.font.color.rgb = clr
    return txBox

# ── Build Presentation ─────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # index 6 = blank

slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)

# ── Left accent bar (vertical, full height, thin cyan) ──────────────────────
add_rect(slide,
         l=Inches(0), t=Inches(0),
         w=Inches(0.06), h=SLIDE_H,
         fill_color=ACCENT_CYAN)

# ── Top horizontal accent line (thin, lime) ──────────────────────────────────
add_rect(slide,
         l=Inches(0.06), t=Inches(0),
         w=SLIDE_W - Inches(0.06), h=Inches(0.04),
         fill_color=ACCENT_LIME)

# ── Bottom horizontal accent line ────────────────────────────────────────────
add_rect(slide,
         l=Inches(0.06), t=SLIDE_H - Inches(0.04),
         w=SLIDE_W - Inches(0.06), h=Inches(0.04),
         fill_color=ACCENT_CYAN)

# ── Big "C M A T R I X" glyph — ghost watermark on right ────────────────────
add_text_box(slide,
             "CMATRIX",
             l=Inches(6.5), t=Inches(1.3),
             w=Inches(6.5), h=Inches(5.5),
             font_name="Calibri", font_size=130,
             bold=True, color=RGBColor(0x00, 0xE5, 0xFF),
             align=PP_ALIGN.RIGHT)

# Make it very transparent-ish by using a very low-alpha trick (we use a light
# version of the colour since pptx alpha for run text isn't straightforward)
# Actually let's use a darker shade for the ghost
# We'll add it as a separate shape with low opacity using XML tweak
# -- simpler: just use very dim colour
ghost = slide.shapes[-1]
# Set alpha via XML
sp = ghost.element
# find <a:solidFill> inside the run
from pptx.oxml.ns import qn
for r in sp.iter(qn('a:r')):
    for srgb in r.iter(qn('a:srgbClr')):
        # wrap in a:solidFill with lumMod for dim effect
        pass  # leave as-is; we set color below
# Re-set colour to a dim version
for run_el in sp.iter(qn('a:r')):
    for solidFill in run_el.iter(qn('a:solidFill')):
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            srgb.set('val', '0D2535')  # very dark teal ghost

# ── Tag line above main title ─────────────────────────────────────────────────
add_text_box(slide,
             "RESEARCH PRESENTATION",
             l=Inches(0.5), t=Inches(0.8),
             w=Inches(8), h=Inches(0.55),
             font_name="Calibri", font_size=13,
             bold=False, color=ACCENT_CYAN,
             align=PP_ALIGN.LEFT)

# ── Main title ────────────────────────────────────────────────────────────────
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(8.8), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "CMatrix"
run.font.name = "Calibri"
run.font.size = Pt(72)
run.font.bold = True
run.font.color.rgb = WHITE

# ── Sub-title (research descriptor) ─────────────────────────────────────────
add_text_box(slide,
             "Dual-Graph-Guided LLM-Orchestrated\nMulti-Agent Framework for Autonomous VAPT",
             l=Inches(0.5), t=Inches(3.0),
             w=Inches(9.5), h=Inches(1.2),
             font_name="Calibri", font_size=22,
             bold=False, italic=True,
             color=ACCENT_CYAN,
             align=PP_ALIGN.LEFT)

# ── Divider line ──────────────────────────────────────────────────────────────
add_rect(slide,
         l=Inches(0.5), t=Inches(4.25),
         w=Inches(6.5), h=Inches(0.03),
         fill_color=GREY_MID)

# ── Three keyword badges ──────────────────────────────────────────────────────
badges = [
    ("Autonomous Pentesting",  Inches(0.5)),
    ("AI Agent Orchestration", Inches(3.5)),
    ("Attack Graph Reasoning", Inches(6.5)),
]
for label, lx in badges:
    # badge background
    add_rect(slide,
             l=lx, t=Inches(4.4),
             w=Inches(2.7), h=Inches(0.42),
             fill_color=CARD_BG,
             line_color=ACCENT_CYAN, line_width_pt=1)
    add_text_box(slide, label,
                 l=lx + Inches(0.08), t=Inches(4.43),
                 w=Inches(2.55), h=Inches(0.38),
                 font_name="Calibri", font_size=12,
                 bold=True, color=ACCENT_LIME,
                 align=PP_ALIGN.CENTER)

# ── Presenter info block ──────────────────────────────────────────────────────
add_text_box(slide,
             "Supervisor Meeting  ·  June 2026",
             l=Inches(0.5), t=Inches(5.15),
             w=Inches(8), h=Inches(0.4),
             font_name="Calibri", font_size=14,
             bold=False, color=GREY_MID,
             align=PP_ALIGN.LEFT)

add_text_box(slide,
             "Black-Box & Grey-Box Assessment  ·  Network · Web · REST API",
             l=Inches(0.5), t=Inches(5.6),
             w=Inches(10), h=Inches(0.4),
             font_name="Calibri", font_size=13,
             bold=False, color=GREY_MID,
             align=PP_ALIGN.LEFT)

# ── Slide number dot ─────────────────────────────────────────────────────────
add_text_box(slide, "01",
             l=SLIDE_W - Inches(0.9), t=SLIDE_H - Inches(0.55),
             w=Inches(0.7), h=Inches(0.45),
             font_name="Calibri", font_size=13,
             bold=True, color=ACCENT_CYAN,
             align=PP_ALIGN.RIGHT)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(PPTX_PATH)
print(f"✅  Saved: {PPTX_PATH}")
