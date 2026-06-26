"""
Slide 3 — Research Contributions (Novel Work)
"""
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG_DARK     = RGBColor(0x0A, 0x0D, 0x1A)
ACCENT_CYAN = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_LIME = RGBColor(0x39, 0xFF, 0x14)
ACCENT_GOLD = RGBColor(0xFF, 0xD7, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID    = RGBColor(0xA0, 0xAA, 0xB8)
CARD_BG     = RGBColor(0x10, 0x16, 0x2B)
CARD_GOLD   = RGBColor(0x1E, 0x19, 0x05)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PPTX_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "CMatrix_Presentation.pptx")

def set_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=1):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    if fill_color: shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else: shape.fill.background()
    if line_color: shape.line.color.rgb = line_color; shape.line.width = Pt(line_width_pt)
    else: shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic; run.font.color.rgb = color
    return tb

prs = Presentation(PPTX_PATH)
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, BG_DARK)

# ── Accent bars ───────────────────────────────────────────────────────────────
add_rect(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill_color=ACCENT_GOLD)
add_rect(slide, Inches(0.06), Inches(0), SLIDE_W - Inches(0.06), Inches(0.04), fill_color=ACCENT_GOLD)
add_rect(slide, Inches(0.06), SLIDE_H - Inches(0.04), SLIDE_W - Inches(0.06), Inches(0.04), fill_color=ACCENT_GOLD)

# ── Slide label + title ───────────────────────────────────────────────────────
add_text(slide, "NOVEL CONTRIBUTION", Inches(0.3), Inches(0.08), Inches(8), Inches(0.35),
         size=11, bold=True, color=ACCENT_GOLD)
add_text(slide, "What Makes CMatrix a Research Contribution",
         Inches(0.3), Inches(0.43), Inches(12.5), Inches(0.75),
         size=34, bold=True, color=WHITE)
add_rect(slide, Inches(0.3), Inches(1.2), Inches(7), Inches(0.035), fill_color=ACCENT_GOLD)

# ── Primary insight banner ────────────────────────────────────────────────────
add_rect(slide, Inches(0.3), Inches(1.35), SLIDE_W - Inches(0.6), Inches(0.58),
         fill_color=RGBColor(0x12, 0x10, 0x02), line_color=ACCENT_GOLD, line_width_pt=1.5)
add_text(slide,
         "CMatrix is the first autonomous VAPT system to maintain two strictly separated, "
         "continuously evolving graph structures — one for discovered reality, one for inferred attack opportunity.",
         Inches(0.55), Inches(1.42), SLIDE_W - Inches(1.0), Inches(0.45),
         size=13, bold=False, italic=True, color=ACCENT_GOLD, wrap=True)

# ── 6 contribution cards in a 3×2 grid ───────────────────────────────────────
contributions = [
    ("C1", "Dual-Graph World Model",
     "ASG (discovered reality) + APG (inferred opportunity) — strictly separated with enforced write ownership. "
     "No prior system maintains these as distinct structures."),
    ("C2", "Graph-State-Driven Re-Planning",
     "Re-planning fires on explicit graph triggers (new vuln node seeds a chain, chain status changes) — "
     "not fixed task completion or arbitrary schedules."),
    ("C3", "APG Attack Chain Lifecycle",
     "Attack chains are first-class entities: risk-scored, prioritized, lifecycle-tracked "
     "(HYPOTHESIZED → VALIDATED / RULED_OUT) with evidence-linked ChainSteps."),
    ("C5", "Tool Risk Gate + Commander Mailbox",
     "Every tool call is risk-classified before execution. High-risk ops route to Commander for approval — "
     "a zero-code insertion point for human-in-the-loop."),
    ("C10", "Cross-Mission Experience Store",
     "Persistent RAG-backed store of validated exploitation outcomes across missions. "
     "Commander seeds APG hypotheses from prior analogous engagements at mission start."),
    ("C11", "Attack Strategy Library",
     "Crystallizes validated chains into named, parameterized strategies when the same target fingerprint "
     "succeeds across 2+ missions. No prior system accumulates strategies across sessions."),
]

cols = 3
rows = 2
card_w = Inches(4.15)
card_h = Inches(1.75)
gap_x  = Inches(0.13)
gap_y  = Inches(0.12)
start_l = Inches(0.3)
start_t = Inches(2.1)

for idx, (cnum, title, body) in enumerate(contributions):
    col = idx % cols
    row = idx // cols
    l = start_l + col * (card_w + gap_x)
    t = start_t + row * (card_h + gap_y)

    add_rect(slide, l, t, card_w, card_h,
             fill_color=CARD_BG, line_color=ACCENT_GOLD, line_width_pt=0.8)
    # Contribution badge
    add_rect(slide, l, t, Inches(0.55), Inches(0.42),
             fill_color=ACCENT_GOLD)
    add_text(slide, cnum, l + Inches(0.02), t + Inches(0.03),
             Inches(0.52), Inches(0.36), size=11, bold=True,
             color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, title, l + Inches(0.62), t + Inches(0.07),
             card_w - Inches(0.7), Inches(0.38),
             size=13, bold=True, color=WHITE)
    add_text(slide, body, l + Inches(0.12), t + Inches(0.52),
             card_w - Inches(0.2), Inches(1.15),
             size=11, bold=False, color=GREY_MID, wrap=True)

# ── Bottom strip — count of contributions ─────────────────────────────────────
add_rect(slide, Inches(0.3), Inches(5.78), SLIDE_W - Inches(0.6), Inches(0.5),
         fill_color=RGBColor(0x0D, 0x12, 0x20), line_color=ACCENT_CYAN, line_width_pt=0.8)
add_text(slide,
         "12 novel research contributions in total  ·  C4: ASG-aware parallel dispatch  ·  "
         "C6: Lossless context compaction  ·  C7: Methodology-as-config  ·  C8: Dual-graph termination  ·  "
         "C9: Live CVE grounding  ·  C12: Trajectory export dataset",
         Inches(0.55), Inches(5.84), SLIDE_W - Inches(0.9), Inches(0.4),
         size=10.5, bold=False, color=GREY_MID, wrap=False)

# ── Slide number ──────────────────────────────────────────────────────────────
add_text(slide, "03", SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.55),
         Inches(0.7), Inches(0.45), size=13, bold=True, color=ACCENT_GOLD,
         align=PP_ALIGN.RIGHT)

prs.save(PPTX_PATH)
print("✅  Slide 3 added.")
