"""
Slide 6 — Agent Architecture (Redesigned)
==========================================
Layout:
  LEFT panel (40%): Context-Isolated Spawning model — visual showing
    Commander → spawn → isolated box (ASG slice + APG slice + toolset + task) → return delta
    This is the KEY architectural insight that makes CMatrix different.

  RIGHT panel (58%): 6 agent cards in 2-col × 3-row grid
    Each card: phase tag | name | role sentence | tools | key property highlight
"""
from palette import *
import pptx.enum.shapes


BG_DARK2    = RGBColor(0x0A, 0x0D, 0x1A)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome ────────────────────────────────────────────────────────────────────
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_PURP)
    box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_PURP)
    box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_PURP)

    # ── Title ─────────────────────────────────────────────────────────────────────
    txt(slide, "AGENT ARCHITECTURE", Inches(0.3), Inches(0.07), Inches(6), Inches(0.26),
        size=10, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)
    txt(slide, "Context-Isolated Agents — Spawn · Execute · Return · Discard",
        Inches(0.3), Inches(0.32), Inches(11), Inches(0.48),
        size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  LEFT — Context Isolation Model (visual diagram)
    #  Shows: Commander → spawn with scoped context → isolated agent → return structured delta
    # ═══════════════════════════════════════════════════════════════════════════════
    LP_L = Inches(0.18)
    LP_W = Inches(4.85)
    LP_T = Inches(0.88)
    LP_H = SLIDE_H - LP_T - Inches(0.28)

    box(slide, LP_L, LP_T, LP_W, LP_H, fill=RGBColor(0x08,0x0C,0x20), line_color=ACCENT_PURP, lw=1.2)
    txt(slide, "CONTEXT ISOLATION MODEL", LP_L, LP_T+Inches(0.06),
        LP_W, Inches(0.22), size=9, bold=True, color=ACCENT_PURP)

    # — Commander node —
    cmd_cx = LP_L + LP_W/2
    box(slide, LP_L+Inches(0.9), LP_T+Inches(0.38), Inches(3.0), Inches(0.55),
        fill=RGBColor(0x04,0x18,0x32), line_color=ACCENT_CYAN, lw=1.8)
    box(slide, LP_L+Inches(0.9), LP_T+Inches(0.38), Inches(3.0), Inches(0.2), fill=ACCENT_CYAN)
    txt(slide, "COMMANDER AGENT", LP_L+Inches(0.9), LP_T+Inches(0.4), Inches(3.0), Inches(0.18),
        size=8, bold=True, color=BG_DARK)
    txt(slide, "Reads ASG + APG → decides next action", LP_L+Inches(0.9), LP_T+Inches(0.6),
        Inches(3.0), Inches(0.26), size=8, color=GREY_MID, italic=True)

    # Arrow: spawn ↓
    arr(slide, cmd_cx, LP_T+Inches(0.93), cmd_cx, LP_T+Inches(1.18), color=ACCENT_PURP, lw=1.5)
    txt(slide, "spawn with scoped context", cmd_cx+Inches(0.08), LP_T+Inches(0.95),
        Inches(1.8), Inches(0.22), size=7.5, italic=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

    # — Isolated Context Box —
    ic_l = LP_L + Inches(0.22)
    ic_t = LP_T + Inches(1.25)
    ic_w = LP_W - Inches(0.44)
    ic_h = Inches(3.5)
    box(slide, ic_l, ic_t, ic_w, ic_h, fill=RGBColor(0x0A,0x10,0x28), line_color=ACCENT_PURP, lw=1.0)
    box(slide, ic_l, ic_t, ic_w, Inches(0.24), fill=RGBColor(0x20,0x10,0x30))
    txt(slide, "ISOLATED AGENT CONTEXT  (fresh per task)", ic_l, ic_t+Inches(0.04),
        ic_w, Inches(0.18), size=7.5, bold=True, color=ACCENT_PURP)

    # 4 context components inside the isolated box
    ctx_items = [
        ("ASG SLICE", "Only nodes relevant to this task",  ACCENT_LIME, RGBColor(0x04,0x16,0x0A)),
        ("APG SLICE", "Relevant AttackChains (if any)",    ACCENT_GOLD, RGBColor(0x1A,0x10,0x02)),
        ("TOOL SET",  "Authorized tools only — no others", ACCENT_RED,  RGBColor(0x18,0x06,0x06)),
        ("TASK SPEC", "Commander's current plan item",     ACCENT_CYAN, RGBColor(0x04,0x14,0x20)),
    ]
    ctx_h = (ic_h - Inches(0.35)) / 4
    for i, (label, detail, clr, bg) in enumerate(ctx_items):
        ct = ic_t + Inches(0.3) + i * ctx_h
        box(slide, ic_l+Inches(0.1), ct, ic_w-Inches(0.2), ctx_h-Inches(0.06),
            fill=bg, line_color=clr, lw=0.8)
        txt(slide, label, ic_l+Inches(0.18), ct+Inches(0.04), Inches(1.1), ctx_h-Inches(0.14),
            size=8, bold=True, color=clr, align=PP_ALIGN.LEFT)
        txt(slide, detail, ic_l+Inches(1.32), ct+Inches(0.04), ic_w-Inches(1.5), ctx_h-Inches(0.14),
            size=8, color=GREY_MID, italic=True, align=PP_ALIGN.LEFT)

    # Arrow: return ↓
    ret_y_top = ic_t + ic_h
    arr(slide, cmd_cx, ret_y_top, cmd_cx, ret_y_top+Inches(0.32), color=ACCENT_LIME, lw=1.5)

    # — Returned delta box —
    rd_t = ret_y_top + Inches(0.38)
    rd_h = Inches(0.56)
    box(slide, ic_l, rd_t, ic_w, rd_h, fill=RGBColor(0x04,0x1E,0x0C), line_color=ACCENT_LIME, lw=1.4)
    txt(slide, "RETURNS: Structured ASG Delta (new nodes + edges only)", ic_l, rd_t+Inches(0.06),
        ic_w, Inches(0.22), size=8, bold=True, color=ACCENT_LIME)
    txt(slide, "Working context discarded — no raw history passes to Commander",
        ic_l, rd_t+Inches(0.3), ic_w, Inches(0.22), size=7.5, italic=True, color=GREY_MID)

    # — 3 properties of context isolation —
    props_t = rd_t + rd_h + Inches(0.1)
    props = [
        ("Commander context stays clean",     ACCENT_CYAN),
        ("Agents can't contaminate each other", ACCENT_PURP),
        ("High-risk refusals don't bias planning", ACCENT_RED),
    ]
    for i, (prop, clr) in enumerate(props):
        pt = props_t + i * Inches(0.27)
        box(slide, ic_l, pt, ic_w, Inches(0.24), fill=RGBColor(0x0C,0x10,0x22),
            line_color=clr, lw=0.6)
        txt(slide, f"✓  {prop}", ic_l+Inches(0.1), pt+Inches(0.03),
            ic_w-Inches(0.15), Inches(0.2), size=8, color=clr, align=PP_ALIGN.LEFT)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  RIGHT — 6 Agent cards (2 columns × 3 rows)
    # ═══════════════════════════════════════════════════════════════════════════════
    RP_L = Inches(5.22)
    RP_W = SLIDE_W - RP_L - Inches(0.18)

    agents_def = [
        ("👑 Commander",  ACCENT_CYAN,  "ORCHESTRATION",
         "Reads the dual graph. Plans. Delegates. Never runs tools.",
         "No tools — reasons over ASG + APG",
         "Only agent that writes to APG. Approves High-risk ops via mailbox."),
        ("🕵️ Recon",     ACCENT_LIME,  "PHASE 1",
         "External reconnaissance and host discovery.",
         "Amass  ·  httpx  ·  Nmap",
         "Writes Domain, Host, Port, Service nodes to ASG."),
        ("🔬 Analysis",  ACCENT_TEAL,  "PHASE 2",
         "Deep enumeration and vulnerability discovery.",
         "WhatWeb · Gobuster · ffuf · Nuclei · ZAP",
         "Writes Technology, Endpoint, Parameter, Vulnerability nodes."),
        ("🔍 Research",  ACCENT_GOLD,  "INTELLIGENCE",
         "Live CVE grounding — only agent with outbound internet access.",
         "NVD · Exploit-DB · GitHub · Vendor Advisories",
         "Enriches ASG Vulnerability nodes with real-time CVSS + PoC data."),
        ("🎯 Validation", ACCENT_RED,  "PHASE 3",
         "Proves vulnerabilities are real and exploitable.",
         "SQLMap  ·  Metasploit",
         "Self-debugging loop: Diagnose → Contextualize → Adapt → Cap (×3)."),
        ("📸 Evidence",  ACCENT_PURP, "PHASE 3",
         "Captures proof artifacts for every validated finding.",
         "EyeWitness",
         "Links Evidence nodes to ASG findings via validated_by edges."),
    ]

    card_cols = 2
    card_w = (RP_W - Inches(0.12)) / card_cols
    card_h = (SLIDE_H - Inches(1.0) - Inches(0.28)) / 3
    gap_x = Inches(0.12); gap_y = Inches(0.1)
    c_start_t = Inches(0.9)

    for idx, (name, clr, phase, role, tools, behavior) in enumerate(agents_def):
        col = idx % card_cols
        row = idx // card_cols
        cl = RP_L + col * (card_w + gap_x)
        ct = c_start_t + row * (card_h + gap_y)

        # Card body
        box(slide, cl, ct, card_w, card_h, fill=CARD_BG, line_color=clr, lw=1.4)
        # Phase tag
        tag_w = Inches(1.1)
        box(slide, cl+card_w-tag_w, ct, tag_w, Inches(0.24), fill=clr)
        txt(slide, phase, cl+card_w-tag_w, ct+Inches(0.03), tag_w, Inches(0.2),
            size=7.5, bold=True, color=BG_DARK)
        # Name
        txt(slide, name, cl+Inches(0.1), ct+Inches(0.04), card_w-tag_w-Inches(0.15), Inches(0.26),
            size=12, bold=True, color=clr, align=PP_ALIGN.LEFT)
        # Divider
        box(slide, cl+Inches(0.1), ct+Inches(0.34), card_w-Inches(0.2), Inches(0.02), fill=clr)
        # Role (italic)
        txt(slide, role, cl+Inches(0.1), ct+Inches(0.38), card_w-Inches(0.18), Inches(0.24),
            size=8.5, italic=True, color=WHITE, align=PP_ALIGN.LEFT)
        # Tools strip
        box(slide, cl+Inches(0.1), ct+Inches(0.65), card_w-Inches(0.2), Inches(0.24),
            fill=RGBColor(0x08,0x10,0x22), line_color=clr, lw=0.5)
        txt(slide, tools, cl+Inches(0.18), ct+Inches(0.67), card_w-Inches(0.3), Inches(0.2),
            size=7.5, bold=True, color=clr, align=PP_ALIGN.LEFT)
        # Behavior
        txt(slide, behavior, cl+Inches(0.1), ct+Inches(0.93), card_w-Inches(0.18), card_h-Inches(1.0),
            size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # ── Slide number ──────────────────────────────────────────────────────────────
    txt(slide, "06", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.52),
        Inches(0.35), Inches(0.42), size=13, bold=True, color=ACCENT_PURP, align=PP_ALIGN.RIGHT)


