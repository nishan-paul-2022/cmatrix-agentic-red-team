"""
Slide 08 — Agent Spawn Lifecycle
==================================
Shows the detailed spawn sequence diagram with the correct spawn call notation
including APG slice [optional], and shows 5 context components.
Also annotates C4 parallel dispatch.
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_PURP)
    slide_header(slide, "AGENT SPAWN LIFECYCLE",
                 "Context-Isolated Spawn · Execute · Return · Discard",
                 ACCENT_PURP, title_size=28, divider_w=10)

    # ── LEFT: Spawn sequence diagram ─────────────────────────────────────────
    LP_L, LP_W = Inches(0.22), Inches(7.2)
    LP_T, LP_H = Inches(0.98), SLIDE_H - Inches(1.26)
    box(slide, LP_L, LP_T, LP_W, LP_H,
        fill=RGBColor(0x06, 0x08, 0x1C), line_color=ACCENT_PURP, lw=1.2)
    txt(slide, "SPAWN SEQUENCE (single-agent case)", LP_L + Inches(0.1), LP_T + Inches(0.06),
        LP_W - Inches(0.2), Inches(0.2), size=9, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

    # Commander → spawn call
    cmd_cx = LP_L + LP_W / 2
    cmd_t = LP_T + Inches(0.38)
    box(slide, LP_L + Inches(0.6), cmd_t, Inches(5.95), Inches(0.5),
        fill=RGBColor(0x04, 0x16, 0x30), line_color=ACCENT_CYAN, lw=1.8)
    box(slide, LP_L + Inches(0.6), cmd_t, Inches(5.95), Inches(0.2), fill=ACCENT_CYAN)
    txt(slide, "COMMANDER AGENT", LP_L + Inches(0.6), cmd_t + Inches(0.02),
        Inches(5.95), Inches(0.18), size=8, bold=True, color=BG_DARK)
    txt(slide, "Reads ASG + APG → decides next action",
        LP_L + Inches(0.65), cmd_t + Inches(0.22), Inches(5.85), Inches(0.22),
        size=8, color=GREY_MID, italic=True)

    # Spawn arrow
    arr(slide, cmd_cx, cmd_t + Inches(0.5), cmd_cx, cmd_t + Inches(0.82),
        color=ACCENT_PURP, lw=1.8)
    txt(slide, "spawn(ASG slice + APG slice[optional] + task + toolset)",
        cmd_cx + Inches(0.08), cmd_t + Inches(0.52),
        Inches(3.4), Inches(0.26), size=8, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

    # Isolated agent box
    ic_t = cmd_t + Inches(0.88)
    ic_l = LP_L + Inches(0.22)
    ic_w = LP_W - Inches(0.44)
    ic_h = Inches(2.9)
    box(slide, ic_l, ic_t, ic_w, ic_h,
        fill=RGBColor(0x08, 0x0C, 0x22), line_color=ACCENT_PURP, lw=1.0)
    box(slide, ic_l, ic_t, ic_w, Inches(0.22), fill=RGBColor(0x18, 0x0C, 0x28))
    txt(slide, "ISOLATED AGENT CONTEXT  (fresh per task — never shared)",
        ic_l + Inches(0.08), ic_t + Inches(0.03), ic_w - Inches(0.12), Inches(0.18),
        size=8, bold=True, color=ACCENT_PURP)

    ctx_items = [
        ("ASG SLICE", "Subgraph relevant to this task only — no global state", ACCENT_LIME, RGBColor(0x04, 0x18, 0x0C)),
        ("APG SLICE", "Relevant AttackChains only — Validation agents only",   ACCENT_GOLD, RGBColor(0x18, 0x10, 0x02)),
        ("TOOL SET",  "Authorized tools only — no cross-phase access",          ACCENT_RED,  RGBColor(0x1A, 0x06, 0x06)),
        ("TASK SPEC", "Commander's current plan item",                          ACCENT_CYAN, RGBColor(0x04, 0x12, 0x20)),
        ("KNOWLEDGE", "Vulnerability-class docs injected at spawn (§7)",        ACCENT_PURP, RGBColor(0x12, 0x06, 0x1E)),
    ]
    ctx_row_h = (ic_h - Inches(0.28)) / len(ctx_items)
    for i, (label, detail, clr, bg) in enumerate(ctx_items):
        ct = ic_t + Inches(0.26) + i * ctx_row_h
        box(slide, ic_l + Inches(0.1), ct, ic_w - Inches(0.2), ctx_row_h - Inches(0.04),
            fill=bg, line_color=clr, lw=0.7)
        txt(slide, label, ic_l + Inches(0.18), ct + Inches(0.04), Inches(1.1), ctx_row_h - Inches(0.1),
            size=8, bold=True, color=clr, align=PP_ALIGN.LEFT)
        txt(slide, detail, ic_l + Inches(1.32), ct + Inches(0.04), ic_w - Inches(1.5), ctx_row_h - Inches(0.1),
            size=8, color=GREY_MID, italic=True, align=PP_ALIGN.LEFT)

    # Return arrow
    ret_y = ic_t + ic_h
    arr(slide, cmd_cx, ret_y, cmd_cx, ret_y + Inches(0.3), color=ACCENT_LIME, lw=1.8)
    rd_t = ret_y + Inches(0.35)
    rd_h = Inches(0.52)
    box(slide, ic_l, rd_t, ic_w, rd_h,
        fill=RGBColor(0x04, 0x1C, 0x0C), line_color=ACCENT_LIME, lw=1.4)
    txt(slide, "RETURNS: Structured ASG Delta only  — working context discarded",
        ic_l + Inches(0.1), rd_t + Inches(0.04), ic_w - Inches(0.16), Inches(0.22),
        size=8.5, bold=True, color=ACCENT_LIME)
    txt(slide, "No raw history, no conversation, no cross-contamination between agents.",
        ic_l + Inches(0.1), rd_t + Inches(0.27), ic_w - Inches(0.16), Inches(0.22),
        size=8, italic=True, color=GREY_MID)

    # Properties
    props = [
        ("Commander context stays clean between all delegated tasks", ACCENT_CYAN),
        ("Agent failures cannot contaminate other agents' reasoning",  ACCENT_PURP),
        ("High-risk refusals do not bias future planning decisions",   ACCENT_RED),
    ]
    for i, (prop, clr) in enumerate(props):
        pt = rd_t + rd_h + Inches(0.1) + i * Inches(0.26)
        box(slide, ic_l, pt, ic_w, Inches(0.22), fill=RGBColor(0x0C, 0x0E, 0x22), line_color=clr, lw=0.5)
        txt(slide, f"✓  {prop}", ic_l + Inches(0.08), pt + Inches(0.03),
            ic_w - Inches(0.12), Inches(0.18), size=8.5, color=clr, align=PP_ALIGN.LEFT)

    # ── RIGHT: C4 Parallel Dispatch ───────────────────────────────────────────
    RP_L = LP_L + LP_W + Inches(0.14)
    RP_W = SLIDE_W - RP_L - Inches(0.22)
    RP_T, RP_H = Inches(0.98), SLIDE_H - Inches(1.26)

    box(slide, RP_L, RP_T, RP_W, RP_H,
        fill=RGBColor(0x06, 0x0C, 0x1C), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, RP_L, RP_T, RP_W, Inches(0.3), fill=ACCENT_GOLD)
    txt(slide, "C4 — ASG-Aware Parallel Dispatch", RP_L, RP_T + Inches(0.04), RP_W, Inches(0.24),
        size=9.5, bold=True, color=BG_DARK)

    txt(slide, "When the ASG contains multiple independent sub-graphs — "
               "hosts or services with no shared edges — the Commander can spawn "
               "multiple agents concurrently.",
        RP_L + Inches(0.1), RP_T + Inches(0.38), RP_W - Inches(0.16), Inches(0.7),
        size=9.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Visual: 3 hosts side by side
    host_labels = ["Host A\n192.168.1.10", "Host B\n10.0.1.5", "Host C\n172.16.0.8"]
    host_colors = [ACCENT_LIME, ACCENT_CYAN, ACCENT_GOLD]
    hw = (RP_W - Inches(0.3)) / 3 - Inches(0.06)
    ht = RP_T + Inches(1.18)
    for i, (lbl, clr) in enumerate(zip(host_labels, host_colors)):
        hl = RP_L + Inches(0.1) + i * (hw + Inches(0.06))
        box(slide, hl, ht, hw, Inches(0.6),
            fill=RGBColor(0x08, 0x14, 0x22), line_color=clr, lw=1.2)
        txt(slide, lbl, hl, ht + Inches(0.06), hw, Inches(0.5),
            size=8.5, bold=True, color=clr)
        # Arrow down to agent
        arr(slide, hl + hw / 2, ht + Inches(0.6), hl + hw / 2, ht + Inches(0.96),
            color=clr, lw=1.2)
        box(slide, hl, ht + Inches(1.0), hw, Inches(0.44),
            fill=RGBColor(0x10, 0x0C, 0x22), line_color=clr, lw=1.0)
        txt(slide, "Analysis\nAgent", hl, ht + Inches(1.02), hw, Inches(0.4),
            size=8, bold=True, color=clr)

    txt(slide, "ASG edges determine dependency — agents on independent sub-graphs run concurrently.\n"
               "The sequence diagrams in this deck show the single-agent case for clarity.",
        RP_L + Inches(0.1), RP_T + Inches(2.0), RP_W - Inches(0.16), Inches(0.7),
        size=9, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Risk tier summary
    box(slide, RP_L + Inches(0.1), RP_T + Inches(2.82), RP_W - Inches(0.2), Inches(0.26),
        fill=RGBColor(0x12, 0x10, 0x06))
    txt(slide, "RISK TIER  ·  TOOL CLASS", RP_L + Inches(0.18), RP_T + Inches(2.85),
        RP_W - Inches(0.3), Inches(0.22), size=8, bold=True, color=ACCENT_GOLD)

    tiers = [
        ("LOW",  "Passive discovery — execute immediately",        ACCENT_LIME),
        ("MED",  "Active enumeration — LLM classifier → exec",    ACCENT_GOLD),
        ("HIGH", "Exploit operations — Commander Mailbox approval", ACCENT_RED),
    ]
    for i, (tier, action, clr) in enumerate(tiers):
        tt = RP_T + Inches(3.14) + i * Inches(0.52)
        box(slide, RP_L + Inches(0.1), tt, RP_W - Inches(0.2), Inches(0.48),
            fill=RGBColor(0x0C, 0x0E, 0x22), line_color=clr, lw=0.7)
        box(slide, RP_L + Inches(0.1), tt, Inches(0.5), Inches(0.48), fill=clr)
        txt(slide, tier, RP_L + Inches(0.12), tt + Inches(0.06), Inches(0.46), Inches(0.36),
            size=9, bold=True, color=BG_DARK)
        txt(slide, action, RP_L + Inches(0.66), tt + Inches(0.06), RP_W - Inches(0.8), Inches(0.36),
            size=9, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "08", ACCENT_PURP)
