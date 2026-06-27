"""
Slide 8 — Real-World Scenario Walkthrough: shopvault.io
=========================================================
Redesign: Replace the 4 vertical panel columns with a cleaner layout:

  TOP ROW:   Timeline header with phase arrows connecting them (flow reading)
  BODY:      4 panels but with a more structured grid — each panel has a clear
             "Input → Agent → Output → ASG/APG delta" micro-structure.
  BOTTOM:    Stats bar with final mission outcome counts.

Key improvement: Each phase panel now shows the DATA FLOW (what went in, what
came out) not just a list — making the autonomy visible.
"""
from palette import *
import pptx.enum.shapes


def arr_h(slide, x1, y1, x2, y2, color=None, lw=2.0):
    """Horizontal connector arrow (x1,y1) → (x2,y2)."""
    return arr(slide, x1, y1, x2, y2, color=color or GREY_MID, lw=lw)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome ─────────────────────────────────────────────────────────────────────
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_GOLD)
    box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)
    box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_GOLD)

    # ── Title ──────────────────────────────────────────────────────────────────────
    txt(slide, "REAL-WORLD SCENARIO", Inches(0.3), Inches(0.07), Inches(5), Inches(0.24),
        size=10, bold=True, color=ACCENT_GOLD)
    txt(slide, "Mission: shopvault.io — Black-Box Assessment — Zero Manual Commands",
        Inches(0.3), Inches(0.31), Inches(12.5), Inches(0.48),
        size=24, bold=True, color=WHITE)

    # ── Context banner ─────────────────────────────────────────────────────────────
    box(slide, Inches(0.18), Inches(0.84), SLIDE_W-Inches(0.36), Inches(0.38),
        fill=RGBColor(0x10,0x10,0x04), line_color=ACCENT_GOLD, lw=0.8)
    txt(slide,
        "Scope: all subdomains · web applications · REST APIs of shopvault.io  "
        "·  Mode: Black-Box (zero prior knowledge)  "
        "·  Operator configures root domain + scope → presses start.",
        Inches(0.35), Inches(0.9), SLIDE_W-Inches(0.65), Inches(0.3),
        size=9, italic=True, color=ACCENT_GOLD)

    # ═══════════════════════════════════════════════════════════════════════════════
    #  PHASE DEFINITIONS
    # ═══════════════════════════════════════════════════════════════════════════════
    phases = [
        {
            "label": "PHASE 1\nReconnaissance", "color": ACCENT_LIME,
            "bg": RGBColor(0x04,0x14,0x08),
            "agent": "🕵️ Recon Agent",
            "tools": ["Amass", "httpx", "Nmap"],
            "findings": [
                ("Amass",  "14 subdomains discovered\napi · admin · staging · pay · ..."),
                ("httpx",  "11 live hosts confirmed\nstaging returns unexpected 200"),
                ("Nmap",   "Ports 80, 443, 8080, 8443 open\npay.shopvault.io: expired TLS"),
            ],
            "delta": "ASG ← 37 new nodes",
            "delta_color": ACCENT_LIME,
        },
        {
            "label": "PHASE 2\nAnalysis + Intel", "color": ACCENT_CYAN,
            "bg": RGBColor(0x04,0x10,0x18),
            "agent": "🔬 Analysis + 🔍 Research",
            "tools": ["WhatWeb", "Gobuster", "ffuf", "Nuclei", "ZAP"],
            "findings": [
                ("WhatWeb",  "WordPress 5.9.3 · Django API\n→ Research: CVE-2022-21661 (CVSS 8.8)"),
                ("Gobuster", "/backup/db_export_2023.sql EXPOSED\n/admin/users · /admin/login"),
                ("ffuf",     "IDOR: user_id param unsanitized\n/api/v1/internal/users found"),
                ("ZAP",      "XSS on /search?q=\nSQL error on staging login"),
            ],
            "delta": "ASG ← 61 nodes\nAPG ← 3 chains seeded",
            "delta_color": ACCENT_CYAN,
        },
        {
            "label": "PHASE 3\nValidation", "color": ACCENT_RED,
            "bg": RGBColor(0x18,0x06,0x06),
            "agent": "🎯 Validation + 📸 Evidence",
            "tools": ["SQLMap", "Metasploit", "EyeWitness"],
            "findings": [
                ("Chain-01\n(risk 8.8)", "SQLMap → SQLi confirmed\nAdmin hash cracked → Metasploit → RCE\nStatus: VALIDATED (escalated 9.1)"),
                ("Chain-02\n(risk 7.5)", "SQLMap on user_id → IDOR confirmed\nAny customer orders accessible"),
                ("Chain-03\n(risk 8.1)", "Blind SQLi on staging → DB creds\nCredential reuse risk flagged"),
            ],
            "delta": "APG ← 3 chains VALIDATED\nASG ← Evidence nodes linked",
            "delta_color": ACCENT_RED,
        },
        {
            "label": "PHASE 4\nReport", "color": ACCENT_PURP,
            "bg": RGBColor(0x10,0x08,0x1C),
            "agent": "📝 Report Agent",
            "tools": ["Reads ASG + APG"],
            "findings": [
                ("Output", "Executive summary · 4 validated chains\nRCE + IDOR + SQLi + exposed DB backup\n11 vuln entries · full attack surface map\nEvidence at every ChainStep"),
            ],
            "delta": "ZERO manual commands",
            "delta_color": ACCENT_PURP,
        },
    ]

    # ── Phase layout calculations ─────────────────────────────────────────────────
    PANEL_TOP = Inches(1.3)
    PANEL_BOT = SLIDE_H - Inches(0.72)
    PANEL_H   = PANEL_BOT - PANEL_TOP

    # widths: Phase1=2.6, Phase2=3.8, Phase3=4.0, Phase4=2.5
    phase_widths = [Inches(2.62), Inches(3.80), Inches(3.98), Inches(2.62)]
    GAP = Inches(0.06)
    start_l = Inches(0.18)

    panel_lefts = []
    xl = start_l
    for i, pw in enumerate(phase_widths):
        panel_lefts.append(xl)
        xl += pw + GAP

    for i, (ph, pl, pw) in enumerate(zip(phases, panel_lefts, phase_widths)):
        clr = ph["color"]
        # Panel background
        box(slide, pl, PANEL_TOP, pw, PANEL_H, fill=ph["bg"], line_color=clr, lw=1.2)
        # Phase header
        box(slide, pl, PANEL_TOP, pw, Inches(0.5), fill=RGBColor(0x08,0x10,0x1C), line_color=clr, lw=1.2)
        txt(slide, ph["label"], pl+Inches(0.08), PANEL_TOP+Inches(0.06),
            pw-Inches(0.12), Inches(0.42), size=10, bold=True, color=clr, align=PP_ALIGN.CENTER)
        # Agent name
        txt(slide, ph["agent"], pl+Inches(0.08), PANEL_TOP+Inches(0.56),
            pw-Inches(0.12), Inches(0.22), size=8.5, bold=True, color=clr)
        # Tool strip
        tools_str = "  ·  ".join(ph["tools"])
        box(slide, pl+Inches(0.08), PANEL_TOP+Inches(0.8), pw-Inches(0.16), Inches(0.2),
            fill=RGBColor(0x08,0x10,0x22), line_color=clr, lw=0.4)
        txt(slide, tools_str, pl+Inches(0.12), PANEL_TOP+Inches(0.81), pw-Inches(0.2), Inches(0.18),
            size=7, bold=True, color=clr)
        # Finding rows
        findings = ph["findings"]
        row_h_avail = PANEL_H - Inches(1.08) - Inches(0.52)
        row_h = row_h_avail / len(findings)
        for j, (tool_lbl, detail) in enumerate(findings):
            ft = PANEL_TOP + Inches(1.06) + j * row_h
            fh = row_h - Inches(0.06)
            box(slide, pl+Inches(0.08), ft, pw-Inches(0.16), fh,
                fill=RGBColor(0x08,0x0E,0x1C), line_color=clr, lw=0.4)
            # Tool badge (top-left corner colour strip)
            box(slide, pl+Inches(0.08), ft, Inches(0.05), fh, fill=clr)
            txt(slide, tool_lbl, pl+Inches(0.16), ft+Inches(0.04),
                pw-Inches(0.3), Inches(0.26), size=8, bold=True, color=clr, wrap=True)
            txt(slide, detail, pl+Inches(0.16), ft+Inches(0.3),
                pw-Inches(0.3), fh-Inches(0.34), size=8, color=GREY_MID, wrap=True)
        # Delta strip at bottom
        delta_t = PANEL_BOT - Inches(0.44)
        box(slide, pl, delta_t, pw, Inches(0.44), fill=RGBColor(0x06,0x0C,0x18), line_color=clr, lw=0.8)
        txt(slide, ph["delta"], pl+Inches(0.1), delta_t+Inches(0.06),
            pw-Inches(0.15), Inches(0.34), size=8.5, bold=True, color=ph["delta_color"])

        # Arrow connecting phases
        if i < 3:
            next_l = panel_lefts[i+1]
            arr_h(slide, pl+pw, PANEL_TOP+Inches(0.25), next_l, PANEL_TOP+Inches(0.25),
                  color=phases[i+1]["color"], lw=1.8)

    # ── Bottom summary strip ───────────────────────────────────────────────────────
    box(slide, Inches(0.18), SLIDE_H-Inches(0.68), SLIDE_W-Inches(0.36), Inches(0.54),
        fill=RGBColor(0x08,0x10,0x22), line_color=ACCENT_CYAN, lw=1.2)
    stats = [
        ("14", "subdomains", ACCENT_LIME),
        ("11", "live hosts", ACCENT_LIME),
        ("28", "open ports", ACCENT_CYAN),
        ("19", "endpoints", ACCENT_CYAN),
        ("11", "vulnerabilities", ACCENT_GOLD),
        ("4",  "validated chains", ACCENT_RED),
        ("0",  "manual commands", ACCENT_PURP),
    ]
    sw = (SLIDE_W - Inches(0.6)) / len(stats)
    for i, (num, label, clr) in enumerate(stats):
        sl = Inches(0.22) + i * sw
        st = SLIDE_H - Inches(0.65)
        txt(slide, num, sl, st, sw, Inches(0.28), size=18, bold=True, color=clr)
        txt(slide, label, sl, st+Inches(0.24), sw, Inches(0.2), size=7.5, color=GREY_MID)

    txt(slide, "15", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.52),
        Inches(0.35), Inches(0.42), size=13, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.RIGHT)


