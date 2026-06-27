"""
Slide 10 — Cross-Mission Learning (Experience Store + Attack Strategy Library)
"""
from palette import *
import pptx.enum.shapes


def arrow_h(slide, x1, y, x2, color=None, lw=1.5):
    """Horizontal arrow from (x1,y) to (x2,y) — used by slide_10 flow diagram."""
    return arr(slide, x1, y, x2, y, color=color or GREY_MID, lw=lw)


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome ─────────────────────────────────────────────────────────────────────
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_TEAL)
    box(slide, Inches(0.06), Inches(0), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_TEAL)
    box(slide, Inches(0.06), SLIDE_H-Inches(0.04), SLIDE_W-Inches(0.06), Inches(0.04), fill=ACCENT_TEAL)

    txt(slide, "CROSS-MISSION LEARNING", Inches(0.3), Inches(0.08), Inches(8), Inches(0.3),
        size=11, bold=True, color=ACCENT_TEAL)
    txt(slide, "Experience Store + Attack Strategy Library — CMatrix Gets Smarter Each Mission",
        Inches(0.3), Inches(0.38), Inches(12.5), Inches(0.62), size=26, bold=True, color=WHITE)
    box(slide, Inches(0.3), Inches(1.0), Inches(7), Inches(0.03), fill=ACCENT_TEAL)

    # ══════════════════════════════════════════════
    #  FLOW DIAGRAM: Mission N → Store → Mission N+1
    # ══════════════════════════════════════════════
    flow_t = Inches(1.1)
    flow_h = Inches(1.0)

    # Mission N box
    box(slide, Inches(0.3), flow_t, Inches(2.5), flow_h,
        fill=RGBColor(0x06,0x14,0x20), line_color=ACCENT_CYAN, lw=1.5)
    txt(slide, "Mission N\n(completed)", Inches(0.45), flow_t+Inches(0.15),
        Inches(2.2), Inches(0.7), size=13, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)

    # Arrow → Report Agent writes
    arrow_h(slide, Inches(2.8), flow_t+flow_h/2, Inches(4.0), ACCENT_CYAN)
    txt(slide, "Report Agent writes\nvalidated chains",
        Inches(2.82), flow_t+Inches(0.1), Inches(1.12), Inches(0.75),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # Cross-Mission Experience Store box
    box(slide, Inches(4.1), flow_t-Inches(0.3), Inches(2.6), flow_h+Inches(0.6),
        fill=RGBColor(0x04,0x20,0x1E), line_color=ACCENT_TEAL, lw=2.0)
    txt(slide, "Cross-Mission\nExperience Store",
        Inches(4.22), flow_t-Inches(0.18), Inches(2.35), Inches(0.5),
        size=12, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)
    txt(slide, "RAG-backed · persistent\nacross all missions",
        Inches(4.22), flow_t+Inches(0.35), Inches(2.35), Inches(0.4),
        size=9, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # Arrow → Crystallize (if 2+ same fingerprint)
    arrow_h(slide, Inches(6.7), flow_t+flow_h/2, Inches(7.9), ACCENT_GOLD)
    txt(slide, "Crystallize\n(>=2 missions)",
        Inches(6.72), flow_t+Inches(0.05), Inches(1.12), Inches(0.75),
        size=8.5, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

    # Attack Strategy Library box
    box(slide, Inches(8.0), flow_t-Inches(0.3), Inches(2.6), flow_h+Inches(0.6),
        fill=RGBColor(0x20,0x18,0x04), line_color=ACCENT_GOLD, lw=2.0)
    txt(slide, "Attack Strategy\nLibrary",
        Inches(8.12), flow_t-Inches(0.18), Inches(2.35), Inches(0.5),
        size=12, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    txt(slide, "Named strategies · confidence scores\ntechnology-fingerprint indexed",
        Inches(8.12), flow_t+Inches(0.35), Inches(2.35), Inches(0.4),
        size=9, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # Arrow → Mission N+1
    arrow_h(slide, Inches(10.6), flow_t+flow_h/2, Inches(11.8), ACCENT_LIME)
    txt(slide, "Pre-ranked seeds\nat mission start",
        Inches(10.62), flow_t+Inches(0.1), Inches(1.12), Inches(0.75),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # Mission N+1 box
    box(slide, Inches(11.9), flow_t, Inches(1.25), flow_h,
        fill=RGBColor(0x06,0x20,0x10), line_color=ACCENT_LIME, lw=1.5)
    txt(slide, "Mission\nN+1", Inches(11.95), flow_t+Inches(0.12),
        Inches(1.15), Inches(0.75), size=12, bold=True, color=ACCENT_LIME, align=PP_ALIGN.CENTER)

    # Also: Experience Store feeds Mission N+1
    arrow_h(slide, Inches(5.4), flow_t+flow_h*0.9, Inches(11.9), ACCENT_TEAL)
    txt(slide, "Candidate chain hypotheses injected into Commander context",
        Inches(6.6), flow_t+flow_h*0.9+Inches(0.04), Inches(4.0), Inches(0.28),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════
    #  TWO CARDS BELOW: Experience Store vs Library
    # ══════════════════════════════════════════════
    card_t = Inches(2.42)
    card_h = Inches(4.55)

    # Experience Store card
    es_l = Inches(0.3); es_w = Inches(6.2)
    box(slide, es_l, card_t, es_w, card_h,
        fill=RGBColor(0x04,0x18,0x18), line_color=ACCENT_TEAL, lw=1.5)
    box(slide, es_l, card_t, es_w, Inches(0.4), fill=ACCENT_TEAL)
    txt(slide, "C10  ·  Cross-Mission Experience Store",
        es_l+Inches(0.12), card_t+Inches(0.06), es_w-Inches(0.2), Inches(0.28),
        size=12, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

    es_items = [
        ("What it stores",
         "Per-mission · per-chain raw records of validated exploitation outcomes.\n"
         "Target technology fingerprint · CVE · successful tool invocation · ChainStep sequence · mission outcome."),
        ("When it's written",
         "Report Agent writes one entry for every chain with terminal status VALIDATED at mission close."),
        ("When it's read",
         "Commander queries immediately after Recon Agent writes first Technology nodes to ASG — "
         "before Analysis begins. Retrieved records become candidate chain hypotheses."),
        ("Research origin",
         "Generalizes AutoAttacker's within-mission experience-reuse mechanism to cross-mission scope. "
         "AutoAttacker reuses subtasks within one session; CMatrix accumulates across every mission ever run."),
    ]
    for i, (heading, body) in enumerate(es_items):
        t = card_t + Inches(0.5) + i * Inches(0.98)
        box(slide, es_l+Inches(0.12), t, es_w-Inches(0.24), Inches(0.94),
            fill=RGBColor(0x06,0x22,0x22), line_color=ACCENT_TEAL, lw=0.5)
        txt(slide, heading, es_l+Inches(0.22), t+Inches(0.06), es_w-Inches(0.38), Inches(0.26),
            size=11, bold=True, color=ACCENT_TEAL)
        txt(slide, body, es_l+Inches(0.22), t+Inches(0.34), es_w-Inches(0.38), Inches(0.56),
            size=9.5, color=GREY_MID, wrap=True)

    # Attack Strategy Library card
    sl_l = Inches(6.75); sl_w = Inches(6.38)
    box(slide, sl_l, card_t, sl_w, card_h,
        fill=RGBColor(0x1C,0x14,0x04), line_color=ACCENT_GOLD, lw=1.5)
    box(slide, sl_l, card_t, sl_w, Inches(0.4), fill=ACCENT_GOLD)
    txt(slide, "C11  ·  Attack Strategy Library",
        sl_l+Inches(0.12), card_t+Inches(0.06), sl_w-Inches(0.2), Inches(0.28),
        size=12, bold=True, color=BG_DARK)

    sl_items = [
        ("What it stores",
         "Generalized attack procedures distilled from multiple validated outcomes against the same "
         "target technology class. Named strategies (e.g. STRAT-WP-SQLI-001) with confidence scores."),
        ("Crystallization threshold",
         "Same technology fingerprint (e.g. WordPress 5.x + WooCommerce + Nginx) produces validated "
         "AttackChain in 2 or more independent missions → Commander triggers crystallization."),
        ("When it's used",
         "Retrieved at mission start. Strategies are injected as pre-ranked APG AttackChain seeds — "
         "prioritized above zero-prior chains because they carry validated track record."),
        ("Why it matters",
         "No existing autonomous VAPT system accumulates and generalizes validated exploitation "
         "procedures across sessions. Every prior system resets to zero knowledge on each mission. "
         "CMatrix's library makes it measurably more efficient on repeat target-type engagements."),
    ]
    for i, (heading, body) in enumerate(sl_items):
        t = card_t + Inches(0.5) + i * Inches(0.98)
        box(slide, sl_l+Inches(0.12), t, sl_w-Inches(0.24), Inches(0.94),
            fill=RGBColor(0x22,0x18,0x06), line_color=ACCENT_GOLD, lw=0.5)
        txt(slide, heading, sl_l+Inches(0.22), t+Inches(0.06), sl_w-Inches(0.38), Inches(0.26),
            size=11, bold=True, color=ACCENT_GOLD)
        txt(slide, body, sl_l+Inches(0.22), t+Inches(0.34), sl_w-Inches(0.38), Inches(0.56),
            size=9.5, color=GREY_MID, wrap=True)

    # ── C12 Trajectory Export callout ──────────────────────────────────────────
    box(slide, Inches(0.3), SLIDE_H-Inches(0.62), SLIDE_W-Inches(0.6), Inches(0.46),
        fill=RGBColor(0x08,0x14,0x22), line_color=ACCENT_PURP, lw=1.0)
    txt(slide,
        "C12 · Trajectory Export:  Every mission produces a machine-readable decision log capturing "
        "ASG/APG triggers, Commander rationale, and actions taken — directly usable as SFT/RL training data for security LLMs.",
        Inches(0.5), SLIDE_H-Inches(0.59), SLIDE_W-Inches(0.9), Inches(0.38),
        size=9, italic=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT, wrap=True)

    txt(slide, "14", SLIDE_W-Inches(0.4), SLIDE_H-Inches(0.55),
        Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_TEAL, align=PP_ALIGN.RIGHT)


