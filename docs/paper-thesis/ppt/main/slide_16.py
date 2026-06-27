"""
Slide 16 — Chain-01 Traceability
==================================
Full end-to-end traceability walkthrough for Chain-01 (CVE-2022-21661 → SQLi → RCE → PII):
Shows how the ASG Vulnerability node seeds the APG AttackChain,
how each ChainStep is linked to Evidence, and how the Report Agent reads both graphs.
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_RED)
    slide_header(slide, "CHAIN-01 TRACEABILITY",
                 "CVE-2022-21661 → SQL Injection → RCE → Customer PII — Full Evidence Chain",
                 ACCENT_RED, title_size=23, divider_w=11)

    # ── Top: ASG path that seeded the chain ──────────────────────────────────
    asg_t = Inches(0.98)
    asg_h = Inches(0.78)
    box(slide, Inches(0.22), asg_t, SLIDE_W - Inches(0.44), asg_h,
        fill=RGBColor(0x04, 0x18, 0x0C), line_color=ACCENT_LIME, lw=1.4)
    box(slide, Inches(0.22), asg_t, Inches(1.8), asg_h, fill=ACCENT_LIME)
    txt(slide, "ASG\nDISCOVERY", Inches(0.26), asg_t + Inches(0.06), Inches(1.72), asg_h - Inches(0.1),
        size=8.5, bold=True, color=BG_DARK)

    asg_nodes = [
        ("Domain\nshopvault.io", ACCENT_LIME),
        ("→ Host\n192.168.1.10", ACCENT_LIME),
        ("→ Port\n:443", ACCENT_LIME),
        ("→ Service\nNginx 1.18.0", ACCENT_LIME),
        ("→ Technology\nWordPress 5.9.3", ACCENT_CYAN),
        ("→ Vulnerability\nCVE-2022-21661\nCVSS 8.8", ACCENT_RED),
    ]
    node_w = (SLIDE_W - Inches(2.44)) / len(asg_nodes)
    for i, (label, clr) in enumerate(asg_nodes):
        nl = Inches(2.12) + i * node_w
        box(slide, nl, asg_t + Inches(0.06), node_w - Inches(0.08), asg_h - Inches(0.12),
            fill=RGBColor(0x06, 0x24, 0x10), line_color=clr, lw=0.8)
        txt(slide, label, nl, asg_t + Inches(0.1), node_w - Inches(0.08), asg_h - Inches(0.14),
            size=8, bold=True, color=clr)

    # Arrow ASG → APG seeding
    arr(slide, SLIDE_W / 2, asg_t + asg_h, SLIDE_W / 2, asg_t + asg_h + Inches(0.2),
        color=ACCENT_GOLD, lw=1.6)
    txt(slide, "starts_at: Vulnerability node → Commander seeds APG AttackChain",
        SLIDE_W / 2 + Inches(0.08), asg_t + asg_h + Inches(0.02),
        Inches(5.5), Inches(0.2), size=8, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

    # ── Middle: APG AttackChain + ChainSteps ─────────────────────────────────
    apg_t = asg_t + asg_h + Inches(0.28)
    apg_h = SLIDE_H - apg_t - Inches(1.02)

    box(slide, Inches(0.22), apg_t, SLIDE_W - Inches(0.44), apg_h,
        fill=RGBColor(0x14, 0x0C, 0x02), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, Inches(0.22), apg_t, SLIDE_W - Inches(0.44), Inches(0.28), fill=ACCENT_GOLD)
    txt(slide,
        "APG  ·  AttackChain: Chain-01  ·  status: VALIDATED  ·  risk_score: 9.1",
        Inches(0.34), apg_t + Inches(0.04), SLIDE_W - Inches(0.6), Inches(0.22),
        size=9.5, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

    chain_steps = [
        (
            "ChainStep 1\nRecognise", "VALIDATED", ACCENT_GOLD,
            "SQLMap → WP_Query SQLi confirmed\nParameter: /?s= | Backend: MySQL 8.0",
            "Evidence\nscreenshot_sqli_01.png\nASG Evidence node",
        ),
        (
            "ChainStep 2\nEscalate", "VALIDATED", ACCENT_GOLD,
            "Admin password hash extracted\nCracked offline (MD5) → admin:Welcome1!",
            "Evidence\nhash_dump.txt\nASG Evidence node",
        ),
        (
            "ChainStep 3\nExploit", "VALIDATED", ACCENT_RED,
            "Metasploit WP-Admin upload → shell\nRCE confirmed on production host",
            "Evidence\nshell_access.png\nASG Evidence node",
        ),
        (
            "Impact\nDemonstrated", "DEMONSTRATED", ACCENT_PURP,
            "Full server access\nCustomer PII: 14,000 records accessible",
            "Evidence\npii_sample.json\nASG Evidence node",
        ),
    ]

    step_w = (SLIDE_W - Inches(0.44)) / len(chain_steps)
    step_t = apg_t + Inches(0.32)
    step_h = apg_h - Inches(0.36)

    for i, (label, status, clr, action, evidence) in enumerate(chain_steps):
        sl = Inches(0.22) + i * step_w
        sw = step_w - Inches(0.06)

        # Step card
        box(slide, sl, step_t, sw, step_h,
            fill=RGBColor(0x1E, 0x14, 0x02), line_color=clr, lw=1.4)
        box(slide, sl, step_t, sw, Inches(0.24), fill=clr)
        txt(slide, label, sl + Inches(0.06), step_t + Inches(0.03),
            sw - Inches(0.4), Inches(0.2), size=8, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)

        # Status badge
        badge_w = Inches(1.05)
        badge_clr = ACCENT_LIME if status == "VALIDATED" else ACCENT_PURP
        box(slide, sl + sw - badge_w - Inches(0.04), step_t + Inches(0.02), badge_w, Inches(0.2), fill=badge_clr)
        txt(slide, status, sl + sw - badge_w - Inches(0.04), step_t + Inches(0.03),
            badge_w, Inches(0.18), size=6.5, bold=True, color=BG_DARK)

        # Action detail
        txt(slide, action, sl + Inches(0.1), step_t + Inches(0.3), sw - Inches(0.16), Inches(0.55),
            size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

        # Evidence section
        box(slide, sl + Inches(0.08), step_t + step_h - Inches(0.72), sw - Inches(0.14), Inches(0.66),
            fill=RGBColor(0x10, 0x06, 0x1E), line_color=ACCENT_PURP, lw=0.6)
        txt(slide, evidence, sl + Inches(0.14), step_t + step_h - Inches(0.7),
            sw - Inches(0.22), Inches(0.62),
            size=8, color=ACCENT_PURP, align=PP_ALIGN.LEFT, wrap=True)
        # supported_by annotation
        txt(slide, "↗ supported_by\nASG Evidence",
            sl + sw - Inches(0.06), step_t + step_h - Inches(0.68),
            Inches(0.88), Inches(0.36),
            size=6.5, italic=True, color=ACCENT_PURP, align=PP_ALIGN.RIGHT)

        # Arrow →
        if i < len(chain_steps) - 1:
            arr(slide, sl + sw, step_t + step_h / 2, sl + step_w, step_t + step_h / 2,
                color=clr, lw=1.4)
            txt(slide, "next_step", sl + sw + Inches(0.02), step_t + step_h / 2 - Inches(0.16),
                Inches(0.7), Inches(0.14), size=6.5, italic=True, color=GREY_MID)

    # ── Bottom bar ────────────────────────────────────────────────────────────
    bot_t = SLIDE_H - Inches(0.56)
    box(slide, Inches(0.22), bot_t, SLIDE_W - Inches(0.44), Inches(0.38),
        fill=RGBColor(0x06, 0x0E, 0x20), line_color=ACCENT_CYAN, lw=1.0)
    txt(slide,
        "Every step fully traceable: ASG Vulnerability → APG AttackChain → ChainSteps → ASG Evidence nodes.  "
        "Report Agent reads both graphs and outputs an evidence-backed narrative without human intervention.",
        Inches(0.4), bot_t + Inches(0.04), SLIDE_W - Inches(0.7), Inches(0.3),
        size=9, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "16", ACCENT_RED)
