"""
Slide 10 — Tool Adapter Layer / Risk Gate
==========================================
Shows the full PreToolUse decision tree:
  Scope Check → Risk Classification (LOW/MED/HIGH) →
    LOW: execute immediately
    MED: LLM Permission Classifier (Fast Filter + CoT) → APPROVE/REJECT/MODIFY
    HIGH: Commander Mailbox → APPROVE/REJECT
Correct REJECT annotation: failure annotated to ASG Vulnerability node (not APG).
Shows all 6 lifecycle hooks.
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_RED)
    slide_header(slide, "TOOL ADAPTER LAYER / RISK GATE",
                 "Every Tool Call Passes Through a Risk Classifier Before Execution",
                 ACCENT_RED, title_size=24, divider_w=11)

    # ── LEFT: Decision tree ───────────────────────────────────────────────────
    DT_L, DT_W = Inches(0.22), Inches(8.2)
    DT_T, DT_H = Inches(0.98), SLIDE_H - Inches(1.26)
    box(slide, DT_L, DT_T, DT_W, DT_H,
        fill=RGBColor(0x08, 0x06, 0x1A), line_color=ACCENT_RED, lw=1.2)
    txt(slide, "PRE-TOOL-USE DECISION TREE", DT_L + Inches(0.1), DT_T + Inches(0.06),
        DT_W - Inches(0.2), Inches(0.2), size=9, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)

    # Step 1: Tool call received
    s1_t = DT_T + Inches(0.36)
    box(slide, DT_L + Inches(0.3), s1_t, DT_W - Inches(0.6), Inches(0.4),
        fill=RGBColor(0x10, 0x0C, 0x24), line_color=ACCENT_PURP, lw=1.2)
    txt(slide, "Agent requests tool call  (e.g. gobuster -u target.com -w wordlist.txt)",
        DT_L + Inches(0.44), s1_t + Inches(0.06), DT_W - Inches(0.72), Inches(0.3),
        size=9, color=WHITE, align=PP_ALIGN.LEFT)

    arr(slide, DT_L + DT_W / 2, s1_t + Inches(0.4), DT_L + DT_W / 2, s1_t + Inches(0.62),
        color=ACCENT_RED, lw=1.4)

    # Step 2: Scope check
    s2_t = s1_t + Inches(0.68)
    box(slide, DT_L + Inches(0.3), s2_t, DT_W - Inches(0.6), Inches(0.42),
        fill=RGBColor(0x18, 0x06, 0x06), line_color=ACCENT_RED, lw=1.4)
    box(slide, DT_L + Inches(0.3), s2_t, DT_W - Inches(0.6), Inches(0.16), fill=ACCENT_RED)
    txt(slide, "SCOPE CHECK", DT_L + Inches(0.4), s2_t + Inches(0.01),
        DT_W - Inches(0.75), Inches(0.16), size=8, bold=True, color=BG_DARK)
    txt(slide, "Target in declared scope?  No → BLOCKED immediately, reason logged.",
        DT_L + Inches(0.4), s2_t + Inches(0.18), DT_W - Inches(0.72), Inches(0.2),
        size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT)

    arr(slide, DT_L + DT_W / 2, s2_t + Inches(0.42), DT_L + DT_W / 2, s2_t + Inches(0.64),
        color=ACCENT_GOLD, lw=1.4)

    # Step 3: Risk classification
    s3_t = s2_t + Inches(0.70)
    box(slide, DT_L + Inches(0.3), s3_t, DT_W - Inches(0.6), Inches(0.42),
        fill=RGBColor(0x18, 0x14, 0x04), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, DT_L + Inches(0.3), s3_t, DT_W - Inches(0.6), Inches(0.16), fill=ACCENT_GOLD)
    txt(slide, "RISK CLASSIFICATION", DT_L + Inches(0.4), s3_t + Inches(0.01),
        DT_W - Inches(0.75), Inches(0.16), size=8, bold=True, color=BG_DARK)
    txt(slide, "Three axes: Scope Alignment  ·  Chain Intent  ·  Parameter Safety  →  LOW / MED / HIGH",
        DT_L + Inches(0.4), s3_t + Inches(0.18), DT_W - Inches(0.72), Inches(0.2),
        size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT)

    # Three branches
    branch_y = s3_t + Inches(0.42)
    arr(slide, DT_L + DT_W / 2, branch_y, DT_L + DT_W / 2, branch_y + Inches(0.22),
        color=GREY_MID, lw=1.2)

    branches = [
        (DT_L + Inches(0.3), Inches(2.06), ACCENT_LIME, "LOW\nPassive Discovery",
         "Execute immediately\nNo approval required"),
        (DT_L + Inches(3.08), Inches(2.06), ACCENT_GOLD, "MED\nActive Enumeration",
         "LLM Permission Classifier\n  Fast Filter → CoT → APPROVE/MODIFY/REJECT"),
        (DT_L + Inches(5.86), Inches(2.06), ACCENT_RED, "HIGH\nExploit Operations",
         "Commander Mailbox\n  Human-in-the-loop APPROVE / REJECT"),
    ]
    bw = Inches(2.22)
    bh = Inches(1.5)
    bt = branch_y + Inches(0.22)
    for bl, _, clr, title, detail in branches:
        box(slide, bl, bt, bw, bh, fill=RGBColor(0x08, 0x0C, 0x22), line_color=clr, lw=1.6)
        box(slide, bl, bt, bw, Inches(0.28), fill=clr)
        txt(slide, title, bl, bt + Inches(0.03), bw, Inches(0.24), size=8.5, bold=True, color=BG_DARK)
        txt(slide, detail, bl + Inches(0.1), bt + Inches(0.34), bw - Inches(0.16), bh - Inches(0.38),
            size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Outcomes
    out_t = bt + bh + Inches(0.16)
    outcomes = [
        (DT_L + Inches(0.3), "EXECUTE", ACCENT_LIME),
        (DT_L + Inches(3.08), "EXECUTE / MODIFY / REJECT", ACCENT_GOLD),
        (DT_L + Inches(5.86), "APPROVE / REJECT", ACCENT_RED),
    ]
    for ol, label, clr in outcomes:
        arr(slide, ol + bw / 2, bt + bh, ol + bw / 2, out_t, color=clr, lw=1.2)
        box(slide, ol, out_t, bw, Inches(0.3), fill=clr)
        txt(slide, label, ol, out_t + Inches(0.04), bw, Inches(0.24), size=8, bold=True, color=BG_DARK)

    # REJECT annotation note (correct target: ASG Vulnerability node)
    rej_t = out_t + Inches(0.35)
    box(slide, DT_L + Inches(0.3), rej_t, DT_W - Inches(0.6), Inches(0.42),
        fill=RGBColor(0x18, 0x06, 0x06), line_color=ACCENT_RED, lw=0.8)
    txt(slide,
        "REJECT outcome:  Tool call cancelled.  Failure reason annotated to ASG Vulnerability node "
        "(not the APG). Commander re-reads ASG on next cycle.",
        DT_L + Inches(0.4), rej_t + Inches(0.04), DT_W - Inches(0.72), Inches(0.34),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Hooks summary
    hook_t = rej_t + Inches(0.5)
    box(slide, DT_L + Inches(0.3), hook_t, DT_W - Inches(0.6), Inches(0.28),
        fill=RGBColor(0x10, 0x0A, 0x22), line_color=ACCENT_PURP, lw=0.8)
    txt(slide,
        "Lifecycle Hooks (all 6):  PreToolUse  ·  PostToolUse  ·  PreAgentSpawn  ·  "
        "PostAgentReturn  ·  PreAPGUpdate  ·  PostMissionTerminate",
        DT_L + Inches(0.4), hook_t + Inches(0.04), DT_W - Inches(0.72), Inches(0.2),
        size=8, italic=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

    # ── RIGHT: Tool risk tier table ───────────────────────────────────────────
    TP_L = DT_L + DT_W + Inches(0.14)
    TP_W = SLIDE_W - TP_L - Inches(0.22)
    TP_T, TP_H = Inches(0.98), SLIDE_H - Inches(1.26)

    box(slide, TP_L, TP_T, TP_W, TP_H,
        fill=RGBColor(0x06, 0x08, 0x18), line_color=ACCENT_PURP, lw=1.2)
    box(slide, TP_L, TP_T, TP_W, Inches(0.28), fill=ACCENT_PURP)
    txt(slide, "TOOL RISK TIERS", TP_L, TP_T + Inches(0.04), TP_W, Inches(0.22),
        size=9.5, bold=True, color=BG_DARK)

    tool_tiers = [
        (ACCENT_LIME, "LOW — Passive",    ["Amass", "httpx", "WhatWeb", "EyeWitness"]),
        (ACCENT_GOLD, "MED — Active",     ["Nmap", "Gobuster", "ffuf", "Nuclei", "OWASP ZAP"]),
        (ACCENT_RED,  "HIGH — Exploit",   ["SQLMap", "Metasploit"]),
    ]
    tier_h = (TP_H - Inches(0.32)) / 3
    for i, (clr, label, tools) in enumerate(tool_tiers):
        tt = TP_T + Inches(0.32) + i * tier_h
        box(slide, TP_L + Inches(0.1), tt, TP_W - Inches(0.2), tier_h - Inches(0.06),
            fill=RGBColor(0x0C, 0x0A, 0x1E), line_color=clr, lw=0.9)
        box(slide, TP_L + Inches(0.1), tt, TP_W - Inches(0.2), Inches(0.26), fill=clr)
        txt(slide, label, TP_L + Inches(0.18), tt + Inches(0.03),
            TP_W - Inches(0.3), Inches(0.22), size=9, bold=True, color=BG_DARK, align=PP_ALIGN.LEFT)
        tool_str = "\n".join(f"• {t}" for t in tools)
        txt(slide, tool_str, TP_L + Inches(0.18), tt + Inches(0.32),
            TP_W - Inches(0.3), tier_h - Inches(0.38),
            size=9, color=clr, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "10", ACCENT_RED)
