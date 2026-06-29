"""
Slide 11 — Autonomous Planning Cycle
======================================
Shows the Commander's per-cycle loop:
  Read ASG + APG → Evaluate Termination → Cycle Guard check → Plan → Dispatch → Wait → Re-read
  With Reflector as a separate corrective path on repeated failures.
Correctly separates Cycle Guard (detection) from Reflector (corrective prompt injection).
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_CYAN)
    slide_header(slide, "AUTONOMOUS PLANNING CYCLE",
                 "Commander Loop — Graph-Driven Re-Planning on Every Cycle",
                 ACCENT_CYAN, title_size=26, divider_w=10)

    # ── Main cycle loop ───────────────────────────────────────────────────────
    # Draw as a vertical flow with loop-back on the right side

    FLOW_L = Inches(0.22)
    FLOW_W = Inches(8.0)
    FLOW_T = Inches(0.98)
    FLOW_H = SLIDE_H - Inches(1.26)
    box(slide, FLOW_L, FLOW_T, FLOW_W, FLOW_H,
        fill=RGBColor(0x04, 0x0A, 0x1C), line_color=ACCENT_CYAN, lw=1.2)
    txt(slide, "COMMANDER PLANNING LOOP (executed on every cycle)", FLOW_L + Inches(0.1), FLOW_T + Inches(0.05),
        FLOW_W - Inches(0.2), Inches(0.2), size=9, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

    steps = [
        ("1  READ", "Read full ASG node state + APG chain priorities", ACCENT_LIME,  RGBColor(0x04, 0x18, 0x0C)),
        ("2  TERMINATE?", "Both conditions met?  ASG exhausted + all APG chains resolved → END", ACCENT_CYAN, RGBColor(0x04, 0x12, 0x20)),
        ("3  CYCLE GUARD", "Detect repeated identical tool call sequences → flag as fixation", ACCENT_GOLD, RGBColor(0x18, 0x12, 0x02)),
        ("4  PLAN", "Select highest-priority HYPOTHESIZED chain · decide next ChainStep · choose agent", ACCENT_PURP, RGBColor(0x10, 0x06, 0x1C)),
        ("5  DISPATCH", "Spawn agent with scoped context (ASG slice + APG slice + task)", ACCENT_LIME, RGBColor(0x04, 0x18, 0x0C)),
        ("6  AWAIT DELTA", "Receive ASG delta from returned agent · write new nodes + edges", ACCENT_CYAN, RGBColor(0x04, 0x14, 0x22)),
        ("7  APG UPDATE", "Re-evaluate chains: new Vulnerability node? → seed new AttackChain", ACCENT_GOLD, RGBColor(0x1C, 0x12, 0x02)),
    ]

    step_h = (FLOW_H - Inches(0.34)) / len(steps)
    cx = FLOW_L + FLOW_W * 0.45

    for i, (label, detail, clr, bg) in enumerate(steps):
        st = FLOW_T + Inches(0.28) + i * step_h
        sh = step_h - Inches(0.06)
        sl = FLOW_L + Inches(0.22)
        sw = FLOW_W - Inches(0.44)

        box(slide, sl, st, sw, sh, fill=bg, line_color=clr, lw=1.2)
        # Step number badge
        box(slide, sl, st, Inches(1.1), sh, fill=clr)
        txt(slide, label, sl, st + Inches(0.05), Inches(1.1), sh - Inches(0.08),
            size=8.5, bold=True, color=BG_DARK)
        txt(slide, detail, sl + Inches(1.16), st + Inches(0.05),
            sw - Inches(1.24), sh - Inches(0.08),
            size=9, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

        if i < len(steps) - 1:
            arr(slide, cx, st + sh, cx, st + sh + Inches(0.06), color=clr, lw=1.4)

    # Loop-back arrow (7 → 1)
    loop_x = FLOW_L + FLOW_W - Inches(0.3)
    loop_top = FLOW_T + Inches(0.28)
    loop_bot = FLOW_T + FLOW_H - Inches(0.06)
    box(slide, loop_x, loop_top, Inches(0.06), loop_bot - loop_top, fill=ACCENT_CYAN)
    arr(slide, loop_x, loop_top, FLOW_L + Inches(0.22) + Inches(1.1), loop_top, color=ACCENT_CYAN, lw=1.4)
    txt(slide, "next cycle", loop_x + Inches(0.08), loop_top + Inches(0.04),
        Inches(0.8), Inches(0.2), size=7.5, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

    # ── RIGHT: Cycle Guard + Reflector detail ─────────────────────────────────
    RP_L = FLOW_L + FLOW_W + Inches(0.14)
    RP_W = SLIDE_W - RP_L - Inches(0.22)
    RP_T, RP_H = Inches(0.98), SLIDE_H - Inches(1.26)

    # Cycle Guard box
    CG_H = RP_H * 0.46
    box(slide, RP_L, RP_T, RP_W, CG_H,
        fill=RGBColor(0x18, 0x14, 0x02), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, RP_L, RP_T, RP_W, Inches(0.3), fill=ACCENT_GOLD)
    txt(slide, "CYCLE GUARD  (detection)", RP_L, RP_T + Inches(0.04), RP_W, Inches(0.24),
        size=9.5, bold=True, color=BG_DARK)

    cg_items = [
        "Monitors Commander's tool call history per cycle",
        "If identical sequence repeats ≥ N times → fixation flag",
        "Trigger: inject Reflector corrective prompt",
        "Prevents infinite loops on stuck states",
    ]
    for i, item in enumerate(cg_items):
        ct = RP_T + Inches(0.38) + i * Inches(0.38)
        box(slide, RP_L + Inches(0.1), ct, RP_W - Inches(0.2), Inches(0.32),
            fill=RGBColor(0x1E, 0x18, 0x04), line_color=ACCENT_GOLD, lw=0.4)
        txt(slide, f"→  {item}", RP_L + Inches(0.18), ct + Inches(0.04),
            RP_W - Inches(0.3), Inches(0.24), size=9, color=GREY_MID, align=PP_ALIGN.LEFT)

    # Reflector box
    RF_T = RP_T + CG_H + Inches(0.12)
    RF_H = RP_H - CG_H - Inches(0.12)
    box(slide, RP_L, RF_T, RP_W, RF_H,
        fill=RGBColor(0x08, 0x14, 0x22), line_color=ACCENT_CYAN, lw=1.4)
    box(slide, RP_L, RF_T, RP_W, Inches(0.3), fill=ACCENT_CYAN)
    txt(slide, "REFLECTOR  (corrective — separate from Cycle Guard)",
        RP_L, RF_T + Inches(0.04), RP_W, Inches(0.24), size=9, bold=True, color=BG_DARK)

    rf_items = [
        "Triggered by: Cycle Guard fixation flag OR Validation Agent cap exhaustion",
        "Action: injects a corrective context window into Commander's next prompt",
        "Prompt content: what was tried, why it failed, which alternative approaches remain",
        "Outcome: Commander re-plans with explicit knowledge of stuck state",
        "NOT the same as the planning loop — fires only on anomaly",
    ]
    for i, item in enumerate(rf_items):
        rt = RF_T + Inches(0.38) + i * Inches(0.42)
        box(slide, RP_L + Inches(0.1), rt, RP_W - Inches(0.2), Inches(0.36),
            fill=RGBColor(0x06, 0x14, 0x26), line_color=ACCENT_CYAN, lw=0.4)
        txt(slide, f"→  {item}", RP_L + Inches(0.18), rt + Inches(0.04),
            RP_W - Inches(0.3), Inches(0.28), size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "11", ACCENT_CYAN)
