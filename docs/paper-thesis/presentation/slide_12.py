"""
Slide 12 — Dual Termination + Context Compaction
==================================================
LEFT: Dual termination condition diagram with both conditions (ASG exhaustion + APG resolution)
      shown as a logical AND gate — neither alone is sufficient.
RIGHT: Three-layer Context Compaction scheme with motivation sentence explaining
       why it's necessary for long missions (context fills in 20-30 calls without it).
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_TEAL)
    slide_header(slide, "DUAL TERMINATION + CONTEXT COMPACTION",
                 "Mission Ends When — and Only When — Both Conditions Are Simultaneously True",
                 ACCENT_TEAL, title_size=22, divider_w=11)

    # ── LEFT: Dual Termination ────────────────────────────────────────────────
    DT_L, DT_W = Inches(0.22), Inches(7.2)
    DT_T, DT_H = Inches(0.98), SLIDE_H - Inches(1.26)
    box(slide, DT_L, DT_T, DT_W, DT_H,
        fill=RGBColor(0x04, 0x10, 0x18), line_color=ACCENT_TEAL, lw=1.4)
    box(slide, DT_L, DT_T, DT_W, Inches(0.3), fill=ACCENT_TEAL)
    txt(slide, "C8 — DUAL TERMINATION CONDITION", DT_L, DT_T + Inches(0.04), DT_W, Inches(0.24),
        size=9.5, bold=True, color=BG_DARK)

    # Condition 1: ASG exhaustion
    c1_t = DT_T + Inches(0.38)
    c1_l = DT_L + Inches(0.2)
    c1_w = DT_W - Inches(0.4)
    c1_h = Inches(1.6)
    box(slide, c1_l, c1_t, c1_w, c1_h, fill=RGBColor(0x04, 0x1E, 0x0E), line_color=ACCENT_LIME, lw=1.6)
    box(slide, c1_l, c1_t, c1_w, Inches(0.28), fill=ACCENT_LIME)
    txt(slide, "CONDITION 1  —  ASG Exhaustion", c1_l, c1_t + Inches(0.04), c1_w, Inches(0.22),
        size=9.5, bold=True, color=BG_DARK)
    txt(slide,
        "Every node in the ASG has status INVESTIGATED.\n"
        "No node remains in DISCOVERED or IN_PROGRESS state.\n"
        "Formally: ∀ n ∈ V(ASG) : n.status = INVESTIGATED",
        c1_l + Inches(0.12), c1_t + Inches(0.34), c1_w - Inches(0.2), Inches(1.18),
        size=10, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # AND gate
    and_t = c1_t + c1_h + Inches(0.08)
    box(slide, DT_L + DT_W / 2 - Inches(0.5), and_t, Inches(1.0), Inches(0.38),
        fill=ACCENT_TEAL, line_color=WHITE, lw=1.0)
    txt(slide, "AND", DT_L + DT_W / 2 - Inches(0.5), and_t + Inches(0.04), Inches(1.0), Inches(0.3),
        size=13, bold=True, color=BG_DARK)

    # Condition 2: APG resolution
    c2_t = and_t + Inches(0.44)
    c2_l = c1_l
    c2_w = c1_w
    c2_h = Inches(1.6)
    box(slide, c2_l, c2_t, c2_w, c2_h, fill=RGBColor(0x1C, 0x14, 0x02), line_color=ACCENT_GOLD, lw=1.6)
    box(slide, c2_l, c2_t, c2_w, Inches(0.28), fill=ACCENT_GOLD)
    txt(slide, "CONDITION 2  —  APG Resolution", c2_l, c2_t + Inches(0.04), c2_w, Inches(0.22),
        size=9.5, bold=True, color=BG_DARK)
    txt(slide,
        "Every AttackChain in the APG has terminal status.\n"
        "No chain remains HYPOTHESIZED or PARTIALLY_VALIDATED.\n"
        "Formally: ∀ c ∈ V(APG) : c.status ∈ {VALIDATED, RULED_OUT}",
        c2_l + Inches(0.12), c2_t + Inches(0.34), c2_w - Inches(0.2), Inches(1.18),
        size=10, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Result
    res_t = c2_t + c2_h + Inches(0.14)
    box(slide, c1_l, res_t, c1_w, Inches(0.44), fill=ACCENT_TEAL)
    txt(slide, "→  MISSION TERMINATES — Report Agent invoked",
        c1_l + Inches(0.14), res_t + Inches(0.06), c1_w - Inches(0.2), Inches(0.32),
        size=10.5, bold=True, color=BG_DARK)

    # Novelty note
    nov_t = res_t + Inches(0.52)
    box(slide, c1_l, nov_t, c1_w, Inches(0.98),
        fill=RGBColor(0x04, 0x18, 0x20), line_color=ACCENT_TEAL, lw=0.8)
    txt(slide,
        "Why this is novel:\n"
        "Timer-based systems (PentestGPT) terminate by wall-clock — cannot express APG resolution.\n"
        "Queue-based systems (VulnBot) terminate when queue is empty — cannot express 'all chains resolved'.\n"
        "CMatrix is the only system where both must be true simultaneously.",
        c1_l + Inches(0.1), nov_t + Inches(0.06), c1_w - Inches(0.18), Inches(0.88),
        size=8.5, italic=True, color=ACCENT_TEAL, align=PP_ALIGN.LEFT, wrap=True)

    # ── RIGHT: Context Compaction ─────────────────────────────────────────────
    CP_L = DT_L + DT_W + Inches(0.14)
    CP_W = SLIDE_W - CP_L - Inches(0.22)
    CP_T, CP_H = Inches(0.98), SLIDE_H - Inches(1.26)

    box(slide, CP_L, CP_T, CP_W, CP_H,
        fill=RGBColor(0x04, 0x0A, 0x18), line_color=ACCENT_PURP, lw=1.4)
    box(slide, CP_L, CP_T, CP_W, Inches(0.3), fill=ACCENT_PURP)
    txt(slide, "C6 — LOSSLESS CONTEXT COMPACTION", CP_L, CP_T + Inches(0.04), CP_W, Inches(0.24),
        size=9.5, bold=True, color=BG_DARK)

    # Motivation
    box(slide, CP_L + Inches(0.1), CP_T + Inches(0.36), CP_W - Inches(0.2), Inches(0.52),
        fill=RGBColor(0x14, 0x08, 0x22), line_color=ACCENT_PURP, lw=0.7)
    txt(slide,
        "Motivation:  Without compaction, context window fills in 20–30 tool calls "
        "on real-world missions (multiple phases, many subdomains). "
        "Raw history cannot be truncated without losing discovered facts.",
        CP_L + Inches(0.18), CP_T + Inches(0.4), CP_W - Inches(0.32), Inches(0.44),
        size=9, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Three layers
    layers = [
        (
            "LAYER 1  —  Permanent State\n(ASG itself)",
            "The graph IS the memory. All discovered facts are in typed nodes + edges.\n"
            "Never discarded. Available to all agents via graph query.",
            ACCENT_LIME, RGBColor(0x04, 0x18, 0x0C),
        ),
        (
            "LAYER 2  —  Phase Summary\n(structured digest)",
            "After each phase, a structured JSON summary of that phase's findings is stored.\n"
            "Replaces raw tool output. ASG-queryable for follow-on agents.",
            ACCENT_GOLD, RGBColor(0x18, 0x10, 0x02),
        ),
        (
            "LAYER 3  —  Active Window\n(last N exchanges only)",
            "Commander's live context contains only the last N tool calls + current task.\n"
            "All older history has been compacted into Layers 1 + 2. Zero loss of findings.",
            ACCENT_TEAL, RGBColor(0x04, 0x18, 0x18),
        ),
    ]
    layer_h = (CP_H - Inches(1.02)) / 3
    for i, (title, body, clr, bg) in enumerate(layers):
        lt = CP_T + Inches(0.98) + i * layer_h
        lh = layer_h - Inches(0.06)
        box(slide, CP_L + Inches(0.1), lt, CP_W - Inches(0.2), lh,
            fill=bg, line_color=clr, lw=1.2)
        box(slide, CP_L + Inches(0.1), lt, CP_W - Inches(0.2), Inches(0.26), fill=clr)
        txt(slide, title, CP_L + Inches(0.18), lt + Inches(0.03),
            CP_W - Inches(0.3), Inches(0.22), size=8.5, bold=True, color=BG_DARK)
        txt(slide, body, CP_L + Inches(0.18), lt + Inches(0.32),
            CP_W - Inches(0.3), lh - Inches(0.36),
            size=9, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "12", ACCENT_TEAL)
