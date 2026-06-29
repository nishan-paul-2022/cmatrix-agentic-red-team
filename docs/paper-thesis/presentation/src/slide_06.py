"""
Slide 5 — Dual-Graph World Model (Deep Dive — Visual Graph Diagram)
====================================================================
LEFT  (ASG):  Knowledge graph: Domain→Host→Port→Service→Vulnerability→Evidence
              + branches: Port→Technology, Service→Endpoint→Parameter
RIGHT (APG):  Attack chain: AttackChain container with ChainStep nodes,
              validation badges, edge labels, supported_by annotations
CENTRE:       Strict Separation barrier showing write ownership
BOTTOM:       Architectural principle banner
"""
from palette import *
import pptx.enum.shapes


def node(slide, cx, cy, w, h, label, sublabel="", fill_col=None, border_col=None, label_size=9):
    """Draw a graph node (rectangle) centred at cx, cy."""
    fc = fill_col or RGBColor(0x10, 0x18, 0x2C)
    bc = border_col or ACCENT_CYAN
    l = cx - w/2; t = cy - h/2
    box(slide, l, t, w, h, fill=fc, line_color=bc, lw=1.4)
    ty = t + Inches(0.04) if sublabel else t
    txt(slide, label, l, ty, w, h/2 if sublabel else h-Inches(0.06),
        size=label_size, bold=True, color=bc)
    if sublabel:
        txt(slide, sublabel, l, cy, w, h/2 - Inches(0.04), size=7, color=GREY_MID, italic=True)


def connector(slide, x1, y1, x2, y2, color=None, lw=1.2, label="", bidirectional=False):
    """Draw a connector with optional mid-point label."""
    c = arr(slide, x1, y1, x2, y2, color=color or GREY_MID, lw=lw, bidirectional=bidirectional)
    if label:
        mx = (x1 + x2) / 2; my = (y1 + y2) / 2
        txt(slide, label, mx - Inches(0.55), my - Inches(0.22),
            Inches(1.1), Inches(0.22), size=7, italic=True, color=GREY_MID)
    return c


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome (lime) ─────────────────────────────────────────────────────────
    chrome(slide, ACCENT_LIME)

    # ── Title ─────────────────────────────────────────────────────────────────
    txt(slide, "DUAL-GRAPH WORLD MODEL", Inches(0.3), Inches(0.06), Inches(6), Inches(0.26),
        size=10, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "Two Strictly Separated Knowledge Layers — Visualised",
        Inches(0.3), Inches(0.32), Inches(11), Inches(0.48),
        size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # ── Zone backgrounds ──────────────────────────────────────────────────────
    box(slide, Inches(0.18), Inches(0.88), Inches(5.85), Inches(5.62),
        fill=RGBColor(0x04,0x18,0x0A), line_color=ACCENT_LIME, lw=1.5)
    txt(slide, "ASG — Attack Surface Graph", Inches(0.22), Inches(0.9),
        Inches(3.5), Inches(0.26), size=11, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "What does the target look like?  (Discovered Reality)",
        Inches(0.22), Inches(1.14), Inches(5.5), Inches(0.24),
        size=9, italic=True, color=RGBColor(0x60,0xC0,0x70), align=PP_ALIGN.LEFT)

    box(slide, Inches(7.3), Inches(0.88), Inches(5.85), Inches(5.62),
        fill=RGBColor(0x1C,0x10,0x04), line_color=ACCENT_GOLD, lw=1.5)
    txt(slide, "APG — Attack Path Graph", Inches(7.34), Inches(0.9),
        Inches(3.5), Inches(0.26), size=11, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "What can be done to it?  (Inferred Opportunity)",
        Inches(7.34), Inches(1.14), Inches(5.5), Inches(0.24),
        size=9, italic=True, color=RGBColor(0xC0,0x90,0x20), align=PP_ALIGN.LEFT)

    # ── Separation barrier ────────────────────────────────────────────────────
    sep_cx = Inches(6.667)
    box(slide, sep_cx - Inches(0.55), Inches(0.88), Inches(1.1), Inches(5.62),
        fill=RGBColor(0x08,0x10,0x18), line_color=GREY_DARK, lw=0.8)
    box(slide, sep_cx - Inches(0.01), Inches(0.88), Inches(0.02), Inches(5.62), fill=GREY_DARK)
    txt(slide, "STRICT SEPARATION", sep_cx - Inches(0.52), Inches(2.5),
        Inches(1.04), Inches(0.7), size=8, bold=True, color=GREY_MID)
    txt(slide, "No agent crosses this boundary", sep_cx - Inches(0.5), Inches(3.3),
        Inches(1.0), Inches(0.75), size=7.5, italic=True, color=GREY_MID)
    connector(slide, sep_cx - Inches(0.04), Inches(4.25),
              Inches(4.5), Inches(4.25), color=ACCENT_LIME, lw=1.2, label="Discovery Agents write")
    connector(slide, sep_cx + Inches(0.04), Inches(4.6),
              Inches(8.5), Inches(4.6), color=ACCENT_GOLD, lw=1.2, label="Commander writes")

    # ── ASG node graph ────────────────────────────────────────────────────────
    NW = Inches(1.22); NH = Inches(0.44)

    ASG_NODES = {
        "Domain":        (Inches(1.5),  Inches(1.72)),
        "Host":          (Inches(1.5),  Inches(2.42)),
        "Port":          (Inches(1.5),  Inches(3.12)),
        "Service":       (Inches(1.5),  Inches(3.82)),
        "Technology":    (Inches(3.1),  Inches(3.12)),
        "Endpoint":      (Inches(3.1),  Inches(3.82)),
        "Parameter":     (Inches(3.1),  Inches(4.52)),
        "Vulnerability": (Inches(1.5),  Inches(4.52)),
        "Evidence":      (Inches(1.5),  Inches(5.22)),
    }
    ASG_COLORS = {
        "Domain":        (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
        "Host":          (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
        "Port":          (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
        "Service":       (RGBColor(0x06,0x22,0x12), ACCENT_LIME),
        "Technology":    (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
        "Endpoint":      (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
        "Parameter":     (RGBColor(0x08,0x20,0x18), ACCENT_CYAN),
        "Vulnerability": (RGBColor(0x22,0x06,0x06), ACCENT_RED),
        "Evidence":      (RGBColor(0x12,0x08,0x20), ACCENT_PURP),
    }
    ASG_SUBLABELS = {
        "Domain":        "shopvault.io",
        "Host":          "192.168.1.10",
        "Port":          ":443 / :8080",
        "Service":       "Nginx 1.18.0",
        "Technology":    "WordPress 5.9.3",
        "Endpoint":      "/api/v1/orders",
        "Parameter":     "user_id=?",
        "Vulnerability": "CVE-2022-21661",
        "Evidence":      "screenshot.png",
    }
    for name, (cx, cy) in ASG_NODES.items():
        fc, bc = ASG_COLORS[name]
        node(slide, cx, cy, NW, NH, name, ASG_SUBLABELS[name], fc, bc, label_size=8.5)

    # Vertical spine edges
    chain = ["Domain", "Host", "Port", "Service", "Vulnerability", "Evidence"]
    edge_labels = {
        ("Domain","Host"): "has_host", ("Host","Port"): "has_port",
        ("Port","Service"): "runs", ("Service","Vulnerability"): "affected_by",
        ("Vulnerability","Evidence"): "validated_by",
    }
    for i in range(len(chain)-1):
        src = ASG_NODES[chain[i]]; dst = ASG_NODES[chain[i+1]]
        lbl = edge_labels.get((chain[i], chain[i+1]), "")
        connector(slide, src[0], src[1]+NH/2, dst[0], dst[1]-NH/2, color=ACCENT_LIME, lw=0.9, label=lbl)

    # Branch edges
    connector(slide, ASG_NODES["Port"][0]+NW/2, ASG_NODES["Port"][1],
              ASG_NODES["Technology"][0]-NW/2, ASG_NODES["Technology"][1],
              color=ACCENT_CYAN, lw=0.8, label="uses")
    connector(slide, ASG_NODES["Service"][0]+NW/2, ASG_NODES["Service"][1],
              ASG_NODES["Endpoint"][0]-NW/2, ASG_NODES["Endpoint"][1],
              color=ACCENT_CYAN, lw=0.8, label="has_endpoint")
    connector(slide, ASG_NODES["Endpoint"][0], ASG_NODES["Endpoint"][1]+NH/2,
              ASG_NODES["Parameter"][0], ASG_NODES["Parameter"][1]-NH/2,
              color=ACCENT_CYAN, lw=0.8, label="has_parameter")

    # ASG edge legend
    box(slide, Inches(0.22), Inches(5.72), Inches(5.78), Inches(0.62),
        fill=RGBColor(0x06,0x1C,0x0C), line_color=ACCENT_LIME, lw=0.6)
    txt(slide, "EDGES:  has_host  ·  has_port  ·  runs  ·  uses  ·  has_endpoint  ·  "
        "has_parameter  ·  affected_by  ·  validated_by",
        Inches(0.3), Inches(5.78), Inches(5.6), Inches(0.46),
        size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # ── APG attack chain visual ───────────────────────────────────────────────
    ac_l, ac_t = Inches(7.42), Inches(1.44)
    ac_w, ac_h = Inches(5.5), Inches(4.7)
    box(slide, ac_l, ac_t, ac_w, ac_h, fill=RGBColor(0x14,0x0C,0x02), line_color=ACCENT_GOLD, lw=1.4)
    box(slide, ac_l, ac_t, ac_w, Inches(0.28), fill=ACCENT_GOLD)
    txt(slide, "AttackChain  ·  risk_score: 9.1  ·  status: VALIDATED",
        ac_l, ac_t+Inches(0.04), ac_w, Inches(0.22), size=8, bold=True, color=BG_DARK)
    txt(slide, "Chain-01: CVE-2022-21661 → SQL Injection → RCE → Customer PII",
        ac_l+Inches(0.12), ac_t+Inches(0.32), ac_w-Inches(0.2), Inches(0.22),
        size=8.5, italic=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)

    steps = [
        ("STEP 1", "SQLMap → WP_Query SQLi confirmed",       ACCENT_GOLD, "VALIDATED"),
        ("STEP 2", "Admin hash extracted + cracked offline", ACCENT_GOLD, "VALIDATED"),
        ("STEP 3", "Metasploit → Web shell RCE achieved",    ACCENT_RED,  "VALIDATED"),
        ("IMPACT", "Full server access Customer PII exposed", ACCENT_PURP, "DEMONSTRATED"),
    ]
    snw = Inches(2.1); snh = Inches(0.65)
    s_start_t = ac_t + Inches(0.65)
    s_cx = ac_l + ac_w/2

    for i, (step_lbl, detail, clr, status) in enumerate(steps):
        sy = s_start_t + i * (snh + Inches(0.22))
        sl = s_cx - snw/2
        is_impact = step_lbl == "IMPACT"
        bg = RGBColor(0x16,0x08,0x20) if is_impact else RGBColor(0x1E,0x14,0x02)
        box(slide, sl, sy, snw, snh, fill=bg, line_color=clr, lw=1.6)
        txt(slide, step_lbl, sl+Inches(0.1), sy+Inches(0.04), Inches(0.8), Inches(0.28),
            size=8.5, bold=True, color=clr, align=PP_ALIGN.LEFT)
        badge_w = Inches(1.05)
        badge_color = ACCENT_LIME if status == "VALIDATED" else ACCENT_PURP
        box(slide, sl+snw-badge_w-Inches(0.06), sy+Inches(0.06), badge_w, Inches(0.2), fill=badge_color)
        txt(slide, status, sl+snw-badge_w-Inches(0.06), sy+Inches(0.07),
            badge_w, Inches(0.18), size=6.5, bold=True, color=BG_DARK)
        txt(slide, detail, sl+Inches(0.1), sy+Inches(0.3), snw-Inches(0.16), Inches(0.34),
            size=8, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)
        if i < len(steps)-1:
            connector(slide, s_cx, sy+snh, s_cx, sy+snh+Inches(0.22), color=clr, lw=1.2)
            txt(slide, "next_step", s_cx+Inches(0.06), sy+snh+Inches(0.02),
                Inches(0.8), Inches(0.2), size=6.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)
        txt(slide, "↗ supported_by ASG Evidence",
            sl+snw+Inches(0.08), sy+Inches(0.15), Inches(0.85), Inches(0.36),
            size=6.5, italic=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)

    txt(slide, "starts_at: ASG Vulnerability node",
        ac_l+Inches(0.12), ac_t+Inches(0.54), ac_w-Inches(0.2), Inches(0.2),
        size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # APG edge legend
    box(slide, Inches(7.34), Inches(5.72), Inches(5.78), Inches(0.62),
        fill=RGBColor(0x1C,0x10,0x02), line_color=ACCENT_GOLD, lw=0.6)
    txt(slide, "EDGES:  starts_at  ·  next_step  ·  achieves  ·  supported_by\n"
        "VALIDATION STATUS:  HYPOTHESIZED → PARTIALLY_VALIDATED → VALIDATED / RULED_OUT",
        Inches(7.42), Inches(5.78), Inches(5.6), Inches(0.5),
        size=7.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # ── Bottom principle bar ──────────────────────────────────────────────────
    box(slide, Inches(0.18), Inches(6.42), SLIDE_W-Inches(0.36), Inches(0.72),
        fill=RGBColor(0x06,0x14,0x22), line_color=ACCENT_CYAN, lw=1.5)
    txt(slide,
        "Separation Principle:  Discovery agents write only to the ASG — they never reason about chains.  "
        "The Commander writes only to the APG — it never runs tools.  "
        "No agent conflates discovered facts with hypothesised attack reasoning.  "
        "Graph (not RAG): typed edges express that CVE-X \u2018affected_by\u2019 Service-Y — a relationship that vector similarity cannot represent.",
        Inches(0.5), Inches(6.5), SLIDE_W-Inches(0.9), Inches(0.58),
        size=10.5, italic=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "06", ACCENT_LIME)
