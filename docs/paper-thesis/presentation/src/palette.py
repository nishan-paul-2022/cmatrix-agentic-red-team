"""
palette.py — Shared constants and drawing helpers for all CMatrix slides.

Import from any slide module:
    from palette import *
"""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import pptx.enum.shapes

# ── Slide dimensions ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Colour palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0A, 0x0D, 0x1A)   # near-black navy
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)   # electric cyan
ACCENT_LIME  = RGBColor(0x39, 0xFF, 0x14)   # neon green
ACCENT_GOLD  = RGBColor(0xFF, 0xD7, 0x00)   # amber gold
ACCENT_RED   = RGBColor(0xFF, 0x45, 0x45)   # alert red
ACCENT_PURP  = RGBColor(0xBD, 0x93, 0xF9)   # soft purple
ACCENT_TEAL  = RGBColor(0x00, 0xBF, 0xD8)   # teal
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_MID     = RGBColor(0xA0, 0xAA, 0xB8)
GREY_DARK    = RGBColor(0x30, 0x38, 0x48)
CARD_BG      = RGBColor(0x10, 0x16, 0x2B)
CARD_DARK    = RGBColor(0x0A, 0x0D, 0x1A)


# ── Drawing helpers ───────────────────────────────────────────────────────────

def set_bg(slide, color):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line_color=None, lw=1.0):
    """Add a rectangle. fill/line_color are RGBColor or None."""
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def txt(slide, text, l, t, w, h, size=11, bold=False, italic=False,
        color=None, align=PP_ALIGN.CENTER, wrap=True, font="Calibri"):
    """Add a single-paragraph text box."""
    if color is None:
        color = WHITE
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def arr(slide, x1, y1, x2, y2, color=None, lw=1.2, bidirectional=False):
    """Draw a straight connector with arrowhead(s)."""
    if color is None:
        color = GREY_MID
    c = slide.shapes.add_connector(
        pptx.enum.shapes.MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(lw)
    ln = c.line._ln
    he = etree.SubElement(ln, qn('a:headEnd'))
    he.set('type', 'arrow'); he.set('w', 'med'); he.set('len', 'med')
    te = etree.SubElement(ln, qn('a:tailEnd'))
    if bidirectional:
        te.set('type', 'arrow'); te.set('w', 'med'); te.set('len', 'med')
    else:
        te.set('type', 'none')
    return c


def chrome(slide, accent_color):
    """Draw the standard left-bar + top-line + bottom-line chrome."""
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=accent_color)
    box(slide, Inches(0.06), Inches(0), SLIDE_W - Inches(0.06), Inches(0.04), fill=accent_color)
    box(slide, Inches(0.06), SLIDE_H - Inches(0.04), SLIDE_W - Inches(0.06), Inches(0.04), fill=accent_color)


def slide_header(slide, tag, title, accent_color, title_size=26, divider_w=None):
    """Draw the standard section-tag + title + optional underline."""
    txt(slide, tag, Inches(0.3), Inches(0.07), Inches(8), Inches(0.26),
        size=10, bold=True, color=accent_color, align=PP_ALIGN.LEFT)
    txt(slide, title, Inches(0.3), Inches(0.32), Inches(12.5), Inches(0.52),
        size=title_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if divider_w:
        box(slide, Inches(0.3), Inches(0.9), Inches(divider_w), Inches(0.03), fill=accent_color)


def slide_number(slide, num_str, accent_color):
    """Draw slide number badge in bottom-right corner."""
    txt(slide, num_str, SLIDE_W - Inches(0.45), SLIDE_H - Inches(0.55),
        Inches(0.42), Inches(0.45), size=13, bold=True, color=accent_color,
        align=PP_ALIGN.RIGHT)
