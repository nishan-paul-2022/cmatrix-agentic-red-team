"""
Slide 16 — Chain-01 Full Traceability
"""
from palette import *


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # Gold border frame (no standard chrome — this slide uses a frame)
    box(slide, Inches(0.08), Inches(0.08), Inches(13.173), Inches(7.34),
        fill=None, line_color=ACCENT_GOLD, lw=1.5)

    # ── Header ────────────────────────────────────────────────────────────────
    box(slide, Inches(0.5), Inches(0.28), Inches(8.0), Inches(0.28),
        fill=BG_DARK, line_color=None)
    txt(slide, "CHAIN-01 TRACEABILITY",
        Inches(0.5), Inches(0.28), Inches(8.0), Inches(0.28),
        size=10, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "16",
        Inches(12.233), Inches(0.28), Inches(0.60), Inches(0.28),
        size=14, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "Chain-01 Full Traceability — From Recon to Evidence",
        Inches(0.5), Inches(0.58), Inches(12.333), Inches(0.55),
        size=28, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "Every claim in the final report follows this exact path through the dual graph — nothing is ever asserted without evidence.",
        Inches(0.5), Inches(1.15), Inches(12.333), Inches(0.30),
        size=12.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # ── ASG Discovery row label ───────────────────────────────────────────────
    txt(slide, "ASG · DISCOVERY  —  RECON TO VULNERABILITY",
        Inches(0.5), Inches(1.58), Inches(12.333), Inches(0.22),
        size=10.5, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)

    # ── ASG node chips ────────────────────────────────────────────────────────
    NODE_TOP = Inches(1.84)
    NODE_H = Inches(0.62)
    CHIP_W = Inches(1.872)
    ARROW_W = Inches(0.22)

    nodes = [
        ("DOMAIN",       "shopvault.io",                 False),
        ("HOST",         "192.168.1.10",                 False),
        ("PORT",         ":443",                         False),
        ("SERVICE",      "Nginx 1.18.0",                 False),
        ("TECHNOLOGY",   "WordPress 5.9.3",              False),
        ("VULNERABILITY","CVE-2022-21661  ·  CVSS 8.8",  True),
    ]
    node_gap = Inches(0.22)  # arrow width
    node_l = Inches(0.5)

    for i, (label, value, is_vuln) in enumerate(nodes):
        fill_c = RGBColor(0x22, 0x08, 0x08) if is_vuln else RGBColor(0x10, 0x14, 0x1F)
        line_c = ACCENT_RED if is_vuln else RGBColor(0x5A, 0x63, 0x78)
        lw = 1.2 if is_vuln else 0.8
        label_c = ACCENT_RED if is_vuln else RGBColor(0x7E, 0x8A, 0xA0)
        box(slide, node_l, NODE_TOP, CHIP_W, NODE_H,
            fill=fill_c, line_color=line_c, lw=lw)
        txt(slide, label,
            node_l + Inches(0.08), NODE_TOP + Inches(0.05), Inches(1.712), Inches(0.20),
            size=8.5, bold=True, color=label_c, align=PP_ALIGN.LEFT)
        txt(slide, value,
            node_l + Inches(0.08), NODE_TOP + Inches(0.27), Inches(1.712), Inches(0.30),
            size=10.5 if is_vuln else 11.0, bold=True,
            color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
        node_l += CHIP_W
        if i < len(nodes) - 1:
            txt(slide, "→",
                node_l, NODE_TOP, ARROW_W, NODE_H,
                size=15, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.CENTER)
            node_l += ARROW_W

    # starts_at label
    txt(slide, "starts_at:  Vulnerability node  →  Commander seeds APG AttackChain",
        Inches(0.5), Inches(2.52), Inches(12.333), Inches(0.20),
        size=9.5, italic=True, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.LEFT)

    # ── APG Chain + ChainSteps ────────────────────────────────────────────────
    CHAIN_TOP = Inches(2.74)
    CHAIN_H = Inches(1.55)
    CBOX_W = Inches(2.291)
    CARR_W = Inches(0.22)

    chain_items = [
        # (fill, line_color, header_color, header_text, title_lines, status, arrow_label)
        (RGBColor(0x1E,0x18,0x04), ACCENT_GOLD,  ACCENT_GOLD,  "🟡 APG · ATTACKCHAIN",
         ["Chain-01", "risk_score: 9.1"], "status: VALIDATED", "starts_at"),
        (RGBColor(0x0E,0x1A,0x06), ACCENT_LIME,  ACCENT_LIME,  "CHAINSTEP 1",
         ["SQLMap →", "WP_Query SQLi"], "status: VALIDATED", "next_step"),
        (RGBColor(0x0E,0x1A,0x06), ACCENT_LIME,  ACCENT_LIME,  "CHAINSTEP 2",
         ["SQLMap dump →", "hash cracked"], "status: VALIDATED", "next_step"),
        (RGBColor(0x0E,0x1A,0x06), ACCENT_LIME,  ACCENT_LIME,  "CHAINSTEP 3",
         ["Metasploit →", "Web shell"], "status: VALIDATED", "achieves"),
        (RGBColor(0x15,0x0C,0x28), ACCENT_PURP,  ACCENT_PURP,  "🟣 IMPACT",
         ["RCE on shopvault.io"], "Customer PII accessible", None),
    ]

    offsets = [Inches(0.5), Inches(3.011), Inches(5.521), Inches(8.032), Inches(10.542)]
    arrow_label_offsets = [Inches(2.391), Inches(4.901), Inches(7.412), Inches(9.922)]

    for i, (fill_c, line_c, hdr_c, hdr_txt, title_lines, status_txt, arr_lbl) in enumerate(chain_items):
        cl = offsets[i]
        box(slide, cl, CHAIN_TOP, CBOX_W, CHAIN_H,
            fill=fill_c, line_color=line_c, lw=1.0)
        txt(slide, hdr_txt,
            cl + Inches(0.14), CHAIN_TOP + Inches(0.12), Inches(2.01), Inches(0.22),
            size=9.5, bold=True, color=hdr_c, align=PP_ALIGN.LEFT)
        for ti, tline in enumerate(title_lines):
            txt(slide, tline,
                cl + Inches(0.14), CHAIN_TOP + Inches(0.40) + ti * Inches(0.26),
                Inches(2.01), Inches(0.26),
                size=13.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
        txt(slide, status_txt,
            cl + Inches(0.14), CHAIN_TOP + Inches(1.13), Inches(2.01), Inches(0.34),
            size=9.5, color=GREY_MID, align=PP_ALIGN.LEFT)
        # Arrow between boxes
        if i < len(chain_items) - 1:
            arr_l = cl + CBOX_W
            txt(slide, "→",
                arr_l, CHAIN_TOP, CARR_W, CHAIN_H,
                size=15, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.CENTER)
            if arr_lbl and i < len(arrow_label_offsets):
                txt(slide, arr_lbl,
                    arrow_label_offsets[i], CHAIN_TOP + CHAIN_H + Inches(0.02),
                    Inches(1.02), Inches(0.18),
                    size=7.5, italic=True, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.LEFT)

    # ── Evidence nodes (below each ChainStep) ─────────────────────────────────
    EV_ITEMS = [
        (Inches(3.011), "sqli-extraction.txt"),
        (Inches(5.521), "user-table-dump.png"),
        (Inches(8.032), "webshell-rce.png"),
        (Inches(10.542), "pii-sample.json"),
    ]
    EV_TOP_LABEL = Inches(4.35)
    EV_TOP_BOX = Inches(4.71)
    EV_H = Inches(0.85)

    for ev_l, ev_file in EV_ITEMS:
        txt(slide, "supported_by",
            ev_l, EV_TOP_LABEL, CBOX_W, Inches(0.18),
            size=7.5, italic=True, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.LEFT)
        # vertical connector stub
        from pptx.oxml.ns import qn
        from lxml import etree
        import pptx.enum.shapes as _sh
        c = slide.shapes.add_connector(_sh.MSO_CONNECTOR.STRAIGHT,
                                       ev_l + CBOX_W / 2, CHAIN_TOP + CHAIN_H,
                                       ev_l + CBOX_W / 2, EV_TOP_BOX)
        c.line.color.rgb = RGBColor(0x5A, 0x63, 0x78)
        c.line.width = Pt(1.0)

        box(slide, ev_l, EV_TOP_BOX, CBOX_W, EV_H,
            fill=RGBColor(0x15, 0x0C, 0x28), line_color=RGBColor(0x3A, 0x2A, 0x5C), lw=1.0)
        txt(slide, "📎 ASG · EVIDENCE",
            ev_l + Inches(0.12), EV_TOP_BOX + Inches(0.10), Inches(2.051), Inches(0.20),
            size=8.5, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)
        txt(slide, ev_file,
            ev_l + Inches(0.12), EV_TOP_BOX + Inches(0.36), Inches(2.051), Inches(0.40),
            size=10.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)

    # ── Reading guide ─────────────────────────────────────────────────────────
    box(slide, Inches(0.5), Inches(5.76), Inches(12.333), Inches(0.62),
        fill=RGBColor(0x10, 0x14, 0x1F), line_color=RGBColor(0x5A, 0x63, 0x78), lw=0.8)
    txt(slide, "Reading this diagram:  ",
        Inches(0.68), Inches(5.83), Inches(2.0), Inches(0.48),
        size=10.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "start at the recon chips (ASG facts) → follow starts_at to the Chain (APG reasoning) → step through the ChainSteps in order → arrive at the Impact → follow supported_by back to the Evidence files (ASG proof).",
        Inches(2.6), Inches(5.83), Inches(9.9), Inches(0.48),
        size=10.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # ── Bottom stats ──────────────────────────────────────────────────────────
    stats = [
        ("8.8 → 9.1", "risk_score escalated",      ACCENT_GOLD, Inches(0.5)),
        ("6",          "recon-to-CVE hops traced",  ACCENT_RED,  Inches(2.967)),
        ("3",          "VALIDATED ChainSteps",       ACCENT_LIME, Inches(5.433)),
        ("4",          "Evidence artifacts linked",  ACCENT_PURP, Inches(7.9)),
        ("0",          "manual commands",            ACCENT_CYAN, Inches(10.366)),
    ]
    for val, label, color, sl in stats:
        txt(slide, val,
            sl, Inches(6.56), Inches(2.467), Inches(0.42),
            size=24, bold=True, color=color, align=PP_ALIGN.LEFT)
        txt(slide, label,
            sl, Inches(6.98), Inches(2.467), Inches(0.22),
            size=9.0, color=RGBColor(0x7E, 0x8A, 0xA0), align=PP_ALIGN.LEFT)
