"""
Slide 12 — Dual Termination & Context Compaction
"""
from palette import *


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    chrome(slide, ACCENT_CYAN)

    # ── Header ────────────────────────────────────────────────────────────────
    box(slide, Inches(0.3), Inches(0.07), Inches(8), Inches(0.27),
        fill=BG_DARK, line_color=None)
    txt(slide, "AUTONOMOUS PLANNING CYCLE",
        Inches(0.3), Inches(0.07), Inches(8), Inches(0.27),
        size=10, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    box(slide, Inches(0.3), Inches(0.31), Inches(11.5), Inches(0.50),
        fill=BG_DARK, line_color=None)
    txt(slide, "Dual Termination & Context Compaction",
        Inches(0.3), Inches(0.31), Inches(11.5), Inches(0.50),
        size=24, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)

    # ── Vertical divider line ─────────────────────────────────────────────────
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree
    import pptx.enum.shapes as _sh
    c = slide.shapes.add_connector(_sh.MSO_CONNECTOR.STRAIGHT,
                                   Inches(6.6), Inches(1.0), Inches(6.6), Inches(7.15))
    c.line.color.rgb = RGBColor(0x23, 0x28, 0x3A)
    c.line.width = Pt(1.0)

    # ═══════════════════════════════════════════════════════
    # LEFT PANEL — DUAL TERMINATION CONDITION
    # ═══════════════════════════════════════════════════════
    # Section header
    box(slide, Inches(0.3), Inches(1.02), Inches(6.1), Inches(0.30),
        fill=BG_DARK, line_color=None)
    txt(slide, "DUAL TERMINATION CONDITION",
        Inches(0.3), Inches(1.02), Inches(3.5), Inches(0.30),
        size=13, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "  —  why both must be true",
        Inches(3.5), Inches(1.02), Inches(2.9), Inches(0.30),
        size=11, bold=False, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # Top condition cards
    # ASG exhausted?
    box(slide, Inches(0.3), Inches(1.42), Inches(2.85), Inches(1.32),
        fill=RGBColor(0x0E, 0x1A, 0x06), line_color=ACCENT_LIME, lw=1.0)
    txt(slide, "ASG exhausted?",
        Inches(0.42), Inches(1.50), Inches(2.61), Inches(0.30),
        size=11.5, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "Every Domain, Host, Port, Service, Technology, Endpoint, and Parameter node has been investigated by the appropriate agent.",
        Inches(0.42), Inches(1.82), Inches(2.61), Inches(0.88),
        size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # APG resolved?
    box(slide, Inches(3.55), Inches(1.42), Inches(2.85), Inches(1.32),
        fill=RGBColor(0x1E, 0x18, 0x04), line_color=ACCENT_GOLD, lw=1.0)
    txt(slide, "APG resolved?",
        Inches(3.67), Inches(1.50), Inches(2.61), Inches(0.30),
        size=11.5, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "Every AttackChain sits in a terminal state — VALIDATED or RULED_OUT. None left HYPOTHESIZED or PARTIALLY_VALIDATED.",
        Inches(3.67), Inches(1.82), Inches(2.61), Inches(0.88),
        size=8.5, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Failure cases row 1
    box(slide, Inches(0.3), Inches(2.92), Inches(2.85), Inches(1.18),
        fill=RGBColor(0x22, 0x08, 0x08), line_color=ACCENT_RED, lw=1.0)
    txt(slide, "❌ ASG only — NOT DONE",
        Inches(0.42), Inches(2.99), Inches(2.61), Inches(0.26),
        size=9.5, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)
    txt(slide, "Surface fully explored but chains still open. Attack reasoning is unfinished.",
        Inches(0.42), Inches(3.26), Inches(2.61), Inches(0.78),
        size=8.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    box(slide, Inches(3.55), Inches(2.92), Inches(2.85), Inches(1.18),
        fill=RGBColor(0x22, 0x08, 0x08), line_color=ACCENT_RED, lw=1.0)
    txt(slide, "❌ APG only — NOT DONE",
        Inches(3.67), Inches(2.99), Inches(2.61), Inches(0.26),
        size=9.5, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)
    txt(slide, "All chains resolved, but new ASG nodes were just written — they might seed new chains.",
        Inches(3.67), Inches(3.26), Inches(2.61), Inches(0.78),
        size=8.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Failure + Success row 2
    box(slide, Inches(0.3), Inches(4.18), Inches(2.85), Inches(1.18),
        fill=RGBColor(0x22, 0x08, 0x08), line_color=ACCENT_RED, lw=1.0)
    txt(slide, "❌ Neither — NOT DONE",
        Inches(0.42), Inches(4.25), Inches(2.61), Inches(0.26),
        size=9.5, bold=True, color=ACCENT_RED, align=PP_ALIGN.LEFT)
    txt(slide, "Both conditions incomplete. Mission continues.",
        Inches(0.42), Inches(4.52), Inches(2.61), Inches(0.78),
        size=8.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    box(slide, Inches(3.55), Inches(4.18), Inches(2.85), Inches(1.18),
        fill=RGBColor(0x0C, 0x20, 0x08), line_color=ACCENT_LIME, lw=1.0)
    txt(slide, "✅ Both true — MISSION COMPLETE",
        Inches(3.67), Inches(4.25), Inches(2.61), Inches(0.26),
        size=9.5, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "ASG fully mapped, every attack opportunity proven or disproven. Report Agent spawned.",
        Inches(3.67), Inches(4.52), Inches(2.61), Inches(0.78),
        size=8.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Why existing systems fail
    box(slide, Inches(0.3), Inches(5.50), Inches(6.1), Inches(1.10),
        fill=RGBColor(0x15, 0x0C, 0x28), line_color=ACCENT_PURP, lw=1.0)
    txt(slide, "Why existing systems fail:  ",
        Inches(0.45), Inches(5.60), Inches(1.5), Inches(0.90),
        size=9.0, bold=True, color=ACCENT_PURP, align=PP_ALIGN.LEFT)
    txt(slide, "timer-based execution stops mid-chain; task-queue-based systems can't express APG resolution at all. CMatrix is the only architecture that evaluates both conditions simultaneously.",
        Inches(1.85), Inches(5.60), Inches(4.45), Inches(0.90),
        size=9.0, color=RGBColor(0xD8, 0xC9, 0xF5), align=PP_ALIGN.LEFT, wrap=True)

    # ═══════════════════════════════════════════════════════
    # RIGHT PANEL — CONTEXT COMPACTION
    # ═══════════════════════════════════════════════════════
    txt(slide, "CONTEXT COMPACTION",
        Inches(6.85), Inches(1.02), Inches(3.5), Inches(0.30),
        size=13, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "  —  how long missions stay sharp",
        Inches(9.8), Inches(1.02), Inches(3.23), Inches(0.30),
        size=11, bold=False, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # NORMAL OPERATION block
    box(slide, Inches(6.85), Inches(1.42), Inches(6.18), Inches(0.95),
        fill=RGBColor(0x0E, 0x1A, 0x06), line_color=ACCENT_LIME, lw=1.0)
    txt(slide, "NORMAL OPERATION",
        Inches(7.0), Inches(1.50), Inches(5.88), Inches(0.24),
        size=9.0, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "MicroCompact — every tool call",
        Inches(7.0), Inches(1.74), Inches(5.88), Inches(0.26),
        size=11.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "Raw tool output is discarded immediately; the agent sees only a 3-line summary.",
        Inches(7.0), Inches(2.02), Inches(5.88), Inches(0.27),
        size=9.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Arrow label
    txt(slide, "context grows  ↓",
        Inches(6.85), Inches(2.40), Inches(6.18), Inches(0.20),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # AUTOCOMPACT block
    box(slide, Inches(6.85), Inches(2.62), Inches(6.18), Inches(0.95),
        fill=RGBColor(0x1E, 0x18, 0x04), line_color=ACCENT_GOLD, lw=1.0)
    txt(slide, "AUTOCOMPACT  ·  @ 60% context",
        Inches(7.0), Inches(2.70), Inches(5.88), Inches(0.24),
        size=9.0, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    txt(slide, "Older turns summarized by a scoped LLM call",
        Inches(7.0), Inches(2.94), Inches(5.88), Inches(0.26),
        size=11.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "The summary replaces raw conversation turns. The agent continues uninterrupted.",
        Inches(7.0), Inches(3.22), Inches(5.88), Inches(0.27),
        size=9.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Arrow label
    txt(slide, "context grows  ↓",
        Inches(6.85), Inches(3.60), Inches(6.18), Inches(0.20),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # FULLCOMPACT block
    box(slide, Inches(6.85), Inches(3.82), Inches(6.18), Inches(1.15),
        fill=RGBColor(0x06, 0x1A, 0x28), line_color=ACCENT_CYAN, lw=1.0)
    txt(slide, "FULLCOMPACT  ·  @ 85% context",
        Inches(7.0), Inches(3.90), Inches(5.88), Inches(0.24),
        size=9.0, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)
    txt(slide, "Entire history rebuilt from scratch",
        Inches(7.0), Inches(4.14), Inches(5.88), Inches(0.26),
        size=11.5, bold=True, color=RGBColor(0xE8, 0xEC, 0xF4), align=PP_ALIGN.LEFT)
    txt(slide, "Reconstructed from: current ASG snapshot · current APG priorities · last N tool results. Zero intelligence lost — everything important already lives in the graph.",
        Inches(7.0), Inches(4.42), Inches(5.88), Inches(0.47),
        size=9.0, color=GREY_MID, align=PP_ALIGN.LEFT, wrap=True)

    # Loop label
    txt(slide, "fresh context  →  loops back to Normal Operation  ↺",
        Inches(6.85), Inches(5.02), Inches(6.18), Inches(0.22),
        size=8.5, italic=True, color=GREY_MID, align=PP_ALIGN.LEFT)

    # ASG key insight block
    box(slide, Inches(6.85), Inches(5.50), Inches(6.18), Inches(1.10),
        fill=RGBColor(0x0C, 0x20, 0x08), line_color=ACCENT_LIME, lw=1.0)
    txt(slide, "The ASG is the key:  ",
        Inches(7.0), Inches(5.60), Inches(1.5), Inches(0.90),
        size=9.0, bold=True, color=ACCENT_LIME, align=PP_ALIGN.LEFT)
    txt(slide, "conversation history is expendable because every discovery lives permanently in the graph. That's what makes FullCompact safe — nothing important is ever lost.",
        Inches(8.4), Inches(5.60), Inches(4.53), Inches(0.90),
        size=9.0, color=RGBColor(0xC8, 0xF5, 0xBD), align=PP_ALIGN.LEFT, wrap=True)

    slide_number(slide, "12", ACCENT_CYAN)
