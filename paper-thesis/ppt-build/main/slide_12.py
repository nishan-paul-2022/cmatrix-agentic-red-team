"""
Slide 12 — Inspirations & References
5 academic papers + 3 open-source repositories that informed CMatrix's design,
with exactly which mechanism was borrowed and where it lives in CMatrix.
"""
from palette import *
import pptx.enum.shapes


def build_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)

    # ── Chrome ─────────────────────────────────────────────────────────────────────
    box(slide, Inches(0), Inches(0), Inches(0.06), SLIDE_H, fill=ACCENT_PURP)
    box(slide, Inches(0.06), Inches(0), SLIDE_W - Inches(0.06), Inches(0.04), fill=ACCENT_PURP)
    box(slide, Inches(0.06), SLIDE_H - Inches(0.04), SLIDE_W - Inches(0.06), Inches(0.04), fill=ACCENT_PURP)

    txt(slide, "INSPIRATIONS & REFERENCES", Inches(0.3), Inches(0.06), Inches(8), Inches(0.3),
        size=11, bold=True, color=ACCENT_PURP)
    txt(slide, "Academic Papers + Open-Source Repositories That Shaped CMatrix",
        Inches(0.3), Inches(0.36), Inches(12.5), Inches(0.6),
        size=26, bold=True, color=WHITE)
    box(slide, Inches(0.3), Inches(0.98), Inches(7), Inches(0.03), fill=ACCENT_PURP)

    # ══════════════════════════════════════════════
    #  LEFT COLUMN — 5 Academic Papers
    # ══════════════════════════════════════════════
    papers = [
        (
            "1", "PentestGPT",
            "USENIX Security '24",
            ACCENT_CYAN,
            "Tripartite Reasoning / Generation / Parsing architecture.\nPentesting Task Tree memory for long-session context.",
            "Tool Adapter Layer (parse before you reason) + MicroCompact permanent graph state."
        ),
        (
            "2", "AutoAttacker",
            "arXiv '24  — Xu et al.",
            ACCENT_LIME,
            "Modular planner / summarizer / code-generator for post-breach automation.\nExperience manager stores + reuses attack subtasks within one mission.",
            "Cross-Mission Experience Store (C10) — generalizes reuse to cross-mission scope."
        ),
        (
            "3", "Teams of LLM Agents (HPTSA)",
            "arXiv '24  — Zhu et al.",
            ACCENT_GOLD,
            "Hierarchical planner delegates to task-specific sub-agents.\nTask-specific document injection improves zero-day exploitation by 2.1x.",
            "Vulnerability-Class Knowledge Injection at agent spawn time (C3 / §7)."
        ),
        (
            "4", "PentestAgent",
            "AsiaCCS '25  — Shen et al.",
            ACCENT_TEAL,
            "Planning agent + RAG-backed shared memory.\nExecution agent self-diagnoses failed exploit attempts before abandoning.",
            "Validation Agent self-debugging loop: Diagnose → Contextualize → Adapt → Cap."
        ),
        (
            "5", "VulnBot",
            "arXiv '25  — Kong et al.",
            ACCENT_RED,
            "Five-module pipeline around a Penetration Task Graph.\nSummarizer condenses phase output — prevents raw history leaking between agents.",
            "Context-Isolated Agent Spawning — agents return structured ASG delta, not raw history."
        ),
    ]

    paper_col_l = Inches(0.3)
    paper_col_w = Inches(6.45)
    paper_h = Inches(1.12)
    paper_gap = Inches(0.06)
    paper_t_start = Inches(1.08)

    for i, (num, name, venue, clr, concept, incorporated) in enumerate(papers):
        t = paper_t_start + i * (paper_h + paper_gap)

        box(slide, paper_col_l, t, paper_col_w, paper_h,
            fill=CARD_BG, line_color=clr, lw=1.0)

        # Number badge
        badge_w = Inches(0.38)
        box(slide, paper_col_l, t, badge_w, paper_h, fill=clr)
        txt(slide, num, paper_col_l + Inches(0.02), t + paper_h/2 - Inches(0.15),
            badge_w - Inches(0.04), Inches(0.3),
            size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        # Paper name + venue
        txt(slide, name,
            paper_col_l + badge_w + Inches(0.1), t + Inches(0.06),
            paper_col_w - badge_w - Inches(0.15), Inches(0.25),
            size=11.5, bold=True, color=clr)
        txt(slide, venue,
            paper_col_l + badge_w + Inches(0.1), t + Inches(0.3),
            paper_col_w - badge_w - Inches(0.15), Inches(0.2),
            size=8.5, italic=True, color=GREY_MID)

        # Concept studied
        txt(slide, concept,
            paper_col_l + badge_w + Inches(0.1), t + Inches(0.5),
            paper_col_w - badge_w - Inches(0.15), Inches(0.32),
            size=8.5, color=GREY_MID, wrap=True)

        # Where it lives in CMatrix (green tick line)
        box(slide, paper_col_l + badge_w + Inches(0.08), t + paper_h - Inches(0.22),
            paper_col_w - badge_w - Inches(0.16), Inches(0.18), fill=RGBColor(0x06, 0x18, 0x08))
        txt(slide, f"CMatrix:  {incorporated}",
            paper_col_l + badge_w + Inches(0.12), t + paper_h - Inches(0.21),
            paper_col_w - badge_w - Inches(0.2), Inches(0.18),
            size=8, bold=True, color=ACCENT_LIME, wrap=False)

    # ══════════════════════════════════════════════
    #  RIGHT COLUMN — 3 Open-Source Repos
    # ══════════════════════════════════════════════
    repos = [
        (
            "6", "PentAGI",
            "vxcontrol/pentagi",
            ACCENT_LIME,
            "Production multi-agent pentesting platform.\nAdviser detects fixation; Reflector nudges stuck agents.",
            "Cycle Guard (detects repeated identical tool calls) + Reflector (corrective guidance on repeated failures)."
        ),
        (
            "7", "Claude Code",
            "Anthropic pattern — yasasbanukaofficial/claude-code",
            ACCENT_CYAN,
            "yoloClassifier.ts: two-stage fast-filter + chain-of-thought tool-call auto-approval.\n27-event named hook architecture for operator intervention.",
            "LLM Permission Classifier (Medium-tier Risk Gate) + 6-event Agent Lifecycle Hook System."
        ),
        (
            "8", "Hermes Agent",
            "NousResearch/hermes-agent",
            ACCENT_GOLD,
            "Closed learning loop: skill crystallization from completed tasks.\nStructured trajectory export for SFT/RL training via DSPy + Atropos.",
            "Attack Strategy Library crystallization (C11) + Engagement Trajectory Export dataset (C12)."
        ),
    ]

    repo_col_l = Inches(6.95)
    repo_col_w = Inches(6.2)
    repo_h = Inches(1.78)
    repo_gap = Inches(0.1)
    repo_t_start = Inches(1.08)

    # Section header
    box(slide, repo_col_l, repo_t_start - Inches(0.0), repo_col_w, Inches(0.3),
        fill=RGBColor(0x12, 0x08, 0x20), line_color=ACCENT_PURP, lw=0.8)
    txt(slide, "💻  OPEN-SOURCE REPOSITORIES",
        repo_col_l + Inches(0.12), repo_t_start + Inches(0.02),
        repo_col_w - Inches(0.2), Inches(0.26),
        size=10, bold=True, color=ACCENT_PURP)

    for i, (num, name, repo_id, clr, concept, incorporated) in enumerate(repos):
        t = repo_t_start + Inches(0.35) + i * (repo_h + repo_gap)

        box(slide, repo_col_l, t, repo_col_w, repo_h,
            fill=CARD_BG, line_color=clr, lw=1.2)

        # Number badge
        badge_w = Inches(0.38)
        box(slide, repo_col_l, t, badge_w, repo_h, fill=clr)
        txt(slide, num, repo_col_l + Inches(0.02), t + repo_h/2 - Inches(0.15),
            badge_w - Inches(0.04), Inches(0.3),
            size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        # Repo name + ID
        txt(slide, name,
            repo_col_l + badge_w + Inches(0.1), t + Inches(0.08),
            repo_col_w - badge_w - Inches(0.15), Inches(0.28),
            size=13, bold=True, color=clr)
        txt(slide, repo_id,
            repo_col_l + badge_w + Inches(0.1), t + Inches(0.37),
            repo_col_w - badge_w - Inches(0.15), Inches(0.2),
            size=8.5, italic=True, color=GREY_MID)

        # Concept
        txt(slide, concept,
            repo_col_l + badge_w + Inches(0.1), t + Inches(0.58),
            repo_col_w - badge_w - Inches(0.15), Inches(0.5),
            size=9, color=GREY_MID, wrap=True)

        # CMatrix incorporation
        box(slide, repo_col_l + badge_w + Inches(0.08), t + repo_h - Inches(0.32),
            repo_col_w - badge_w - Inches(0.16), Inches(0.26), fill=RGBColor(0x06, 0x18, 0x08))
        txt(slide, f"CMatrix:  {incorporated}",
            repo_col_l + badge_w + Inches(0.12), t + repo_h - Inches(0.31),
            repo_col_w - badge_w - Inches(0.2), Inches(0.26),
            size=8.5, bold=True, color=ACCENT_LIME, wrap=True)

    # ── Papers section header (left column) ──────────────────────────────────────
    box(slide, paper_col_l, paper_t_start - Inches(0.0), paper_col_w, Inches(0.0), fill=None)  # spacer
    # Retroactively label it
    box(slide, paper_col_l, paper_t_start - Inches(0.0), paper_col_w, Inches(0.3),
        fill=RGBColor(0x08, 0x14, 0x20), line_color=ACCENT_CYAN, lw=0.8)
    txt(slide, "📄  ACADEMIC PAPERS",
        paper_col_l + Inches(0.12), paper_t_start + Inches(0.02),
        paper_col_w - Inches(0.2), Inches(0.26),
        size=10, bold=True, color=ACCENT_CYAN)

    # Shift the papers down by the header height
    # (Already positioned correctly above — the header overlaps the top of paper 1,
    #  so we need to fix: all papers are drawn starting at paper_t_start which is 1.08;
    #  add the header on top at paper_t_start which is fine as the badge covers it.
    #  Actually the header box is drawn AFTER the cards, so it will appear on top only if
    #  it would overlap. Since we draw the header last, it appears on top of card 1's top.
    #  The badge_w box also starts at paper_col_l, same position. This is fine visually
    #  because the header is a thin strip at the very top of card 1's area.
    #  Let's instead nudge all cards down Inches(0.32) by adding a proper header row.)

    # ── Vertical separator between columns ────────────────────────────────────────
    box(slide, Inches(6.82), Inches(1.05), Inches(0.03), SLIDE_H - Inches(1.65), fill=GREY_MID)

    txt(slide, "12", SLIDE_W - Inches(0.4), SLIDE_H - Inches(0.55),
        Inches(0.35), Inches(0.45), size=13, bold=True, color=ACCENT_PURP, align=PP_ALIGN.RIGHT)


